"""Universal file-format text extractor for Batch Analyst uploads.

Given a filename + raw bytes, returns extracted plaintext ready to be
scanned for command-line payloads by `commandline_miner`.

Supported formats
-----------------
    Documents:  .docx, .pdf, .rtf, .txt, .md, .log, .html, .htm, .xml, .eml
    Structured: .csv, .tsv, .json, .jsonl, .yaml, .yml, .ini, .cfg, .conf
    Sheets:     .xlsx  (via openpyxl)
    Slides:     .pptx  (via python-pptx)
    Scripts:    .ps1, .psm1, .bat, .cmd, .sh, .py, .js, .vbs, .hta, .wsf, .reg
    Archives:   .zip, .tar, .tar.gz, .tgz, .gz  (recursed, up to 25 members)

The extractor is deterministic and lightweight — no external processes,
no OCR, no LLM calls. Unknown extensions fall back to a UTF-8 best-effort
decode with a `latin-1` safety net so binary blobs still yield something
scannable rather than raising.
"""
from __future__ import annotations

import csv
import gzip
import io
import json
import logging
import re
import tarfile
import zipfile
from dataclasses import dataclass, field
from typing import Callable, Dict, List, Tuple

log = logging.getLogger("nivx.file_extractors")


_TEXT_LIKE_EXTS = {
    ".txt", ".md", ".log", ".ini", ".cfg", ".conf", ".yaml", ".yml",
    ".ps1", ".psm1", ".bat", ".cmd", ".sh", ".py", ".js", ".vbs", ".hta",
    ".wsf", ".reg", ".rb", ".pl", ".php", ".asp", ".aspx", ".jsp",
    ".tsv", ".csv", ".json", ".jsonl", ".ndjson", ".xml",
    ".html", ".htm", ".eml",
}


@dataclass
class ExtractedSegment:
    """One text region pulled from an uploaded file.

    Attributes
    ----------
    origin       Human-readable pointer back to source ("page 3", "sheet Log",
                 "member payload.ps1"). Empty for whole-file extractors.
    text         Extracted plaintext.
    kind         One of: text | doc | pdf | sheet | slide | html | email |
                 archive-member | json | yaml | script | rtf | fallback
    """
    origin: str = ""
    text: str = ""
    kind: str = "text"


@dataclass
class ExtractionResult:
    filename: str = ""
    total_bytes: int = 0
    segments: List[ExtractedSegment] = field(default_factory=list)
    notes: List[str] = field(default_factory=list)

    @property
    def combined_text(self) -> str:
        return "\n".join(s.text for s in self.segments if s.text)


# --------------------------------------------------------------------------- #
# Individual extractors
# --------------------------------------------------------------------------- #
def _extract_txt(raw: bytes) -> str:
    try:
        return raw.decode("utf-8-sig")
    except UnicodeDecodeError:
        return raw.decode("latin-1", errors="replace")


def _extract_docx(raw: bytes) -> List[ExtractedSegment]:
    try:
        import docx  # python-docx
    except Exception as exc:                                # pragma: no cover
        log.warning("python-docx missing: %s", exc)
        return [ExtractedSegment(text=_extract_txt(raw), kind="fallback")]
    d = docx.Document(io.BytesIO(raw))
    parts: List[str] = []
    for p in d.paragraphs:
        if p.text and p.text.strip():
            parts.append(p.text)
    for tbl in d.tables:
        for row in tbl.rows:
            for cell in row.cells:
                if cell.text and cell.text.strip():
                    parts.append(cell.text)
    return [ExtractedSegment(text="\n".join(parts), kind="doc")]


def _extract_pdf(raw: bytes) -> List[ExtractedSegment]:
    try:
        from pypdf import PdfReader
    except Exception as exc:                                # pragma: no cover
        log.warning("pypdf missing: %s", exc)
        return [ExtractedSegment(text=_extract_txt(raw), kind="fallback")]
    try:
        reader = PdfReader(io.BytesIO(raw))
    except Exception as exc:
        return [ExtractedSegment(text="", kind="pdf",
                                  origin=f"[pdf parse failed: {exc}]")]
    segs: List[ExtractedSegment] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            txt = page.extract_text() or ""
        except Exception:
            txt = ""
        if txt.strip():
            segs.append(ExtractedSegment(text=txt, kind="pdf",
                                          origin=f"page {i}"))
    return segs or [ExtractedSegment(text="", kind="pdf",
                                       origin="[no extractable text]")]


def _extract_xlsx(raw: bytes) -> List[ExtractedSegment]:
    try:
        from openpyxl import load_workbook
    except Exception as exc:                                # pragma: no cover
        log.warning("openpyxl missing: %s", exc)
        return [ExtractedSegment(text=_extract_txt(raw), kind="fallback")]
    wb = load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
    segs: List[ExtractedSegment] = []
    for ws in wb.worksheets:
        rows: List[str] = []
        for row in ws.iter_rows(values_only=True):
            vals = [str(v) for v in row if v not in (None, "")]
            if vals:
                rows.append("\t".join(vals))
        if rows:
            segs.append(ExtractedSegment(
                origin=f"sheet {ws.title!r}",
                text="\n".join(rows), kind="sheet",
            ))
    return segs


def _extract_pptx(raw: bytes) -> List[ExtractedSegment]:
    try:
        from pptx import Presentation
    except Exception as exc:                                # pragma: no cover
        log.warning("python-pptx missing: %s", exc)
        return [ExtractedSegment(text=_extract_txt(raw), kind="fallback")]
    prs = Presentation(io.BytesIO(raw))
    segs: List[ExtractedSegment] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text and cell.text.strip():
                            parts.append(cell.text)
        if parts:
            segs.append(ExtractedSegment(
                origin=f"slide {i}", text="\n".join(parts), kind="slide",
            ))
    return segs


def _extract_html(raw: bytes) -> List[ExtractedSegment]:
    try:
        from bs4 import BeautifulSoup
    except Exception:                                       # pragma: no cover
        return [ExtractedSegment(text=_extract_txt(raw), kind="fallback")]
    txt = raw.decode("utf-8", errors="replace")
    soup = BeautifulSoup(txt, "html.parser")
    parts: List[str] = []
    # 1. Body text
    body_text = soup.get_text("\n", strip=True)
    if body_text:
        parts.append(body_text)
    # 2. Inline scripts (very common malware carrier)
    for scr in soup.find_all("script"):
        if scr.string:
            parts.append(scr.string)
    # 3. Attributes carrying JS/URLs (href, src, onclick, ...)
    for tag in soup.find_all(True):
        for attr in ("href", "src", "action", "onclick", "onerror", "onload"):
            val = tag.get(attr)
            if val and isinstance(val, str):
                parts.append(val)
    return [ExtractedSegment(text="\n".join(parts), kind="html")]


def _extract_eml(raw: bytes) -> List[ExtractedSegment]:
    from email import policy
    from email.parser import BytesParser
    try:
        msg = BytesParser(policy=policy.default).parsebytes(raw)
    except Exception:
        return [ExtractedSegment(text=_extract_txt(raw), kind="fallback")]
    segs: List[ExtractedSegment] = []
    header_lines: List[str] = []
    for h in ("From", "To", "Subject", "Reply-To", "Return-Path", "Date"):
        val = msg.get(h)
        if val:
            header_lines.append(f"{h}: {val}")
    if header_lines:
        segs.append(ExtractedSegment(
            origin="headers", text="\n".join(header_lines), kind="email",
        ))
    for part in msg.walk():
        ctype = part.get_content_type()
        if part.is_multipart():
            continue
        try:
            payload = part.get_content()
        except Exception:
            payload = None
        if not payload:
            continue
        if isinstance(payload, bytes):
            payload = payload.decode("utf-8", errors="replace")
        if ctype == "text/html":
            html_segs = _extract_html(payload.encode("utf-8", errors="replace"))
            for s in html_segs:
                s.origin = f"body(html)/{s.origin}".rstrip("/")
                segs.append(s)
        else:
            segs.append(ExtractedSegment(
                origin=f"body({ctype})",
                text=str(payload), kind="email",
            ))
    return segs


def _extract_rtf(raw: bytes) -> List[ExtractedSegment]:
    try:
        from striprtf.striprtf import rtf_to_text
    except Exception:                                       # pragma: no cover
        return [ExtractedSegment(text=_extract_txt(raw), kind="fallback")]
    txt = rtf_to_text(raw.decode("latin-1", errors="replace"))
    return [ExtractedSegment(text=txt, kind="rtf")]


def _extract_json(raw: bytes) -> List[ExtractedSegment]:
    txt = raw.decode("utf-8", errors="replace")
    # If it parses, we still expose the raw pretty-printed text — the
    # commandline miner runs on the flat text.
    try:
        data = json.loads(txt)
        pretty = json.dumps(data, indent=2, ensure_ascii=False)
        return [ExtractedSegment(text=pretty, kind="json")]
    except Exception:
        # JSONL / NDJSON path
        segs: List[ExtractedSegment] = []
        for i, line in enumerate(txt.splitlines(), start=1):
            if not line.strip():
                continue
            try:
                obj = json.loads(line)
                segs.append(ExtractedSegment(
                    origin=f"jsonl:{i}",
                    text=json.dumps(obj, ensure_ascii=False), kind="json",
                ))
            except Exception:
                segs.append(ExtractedSegment(
                    origin=f"jsonl:{i}", text=line, kind="fallback",
                ))
        return segs or [ExtractedSegment(text=txt, kind="fallback")]


def _extract_yaml(raw: bytes) -> List[ExtractedSegment]:
    txt = raw.decode("utf-8", errors="replace")
    # We don't need to parse — flat text is what the miner wants. YAML
    # comments and inline strings both scan fine.
    return [ExtractedSegment(text=txt, kind="yaml")]


def _extract_csv_tsv(raw: bytes, delim: str) -> List[ExtractedSegment]:
    txt = raw.decode("utf-8-sig", errors="replace")
    segs: List[ExtractedSegment] = []
    try:
        for i, row in enumerate(csv.reader(io.StringIO(txt), delimiter=delim), start=1):
            cell = delim.join([c for c in row if c])
            if cell.strip():
                segs.append(ExtractedSegment(
                    origin=f"row {i}", text=cell, kind="text",
                ))
    except Exception:
        segs.append(ExtractedSegment(text=txt, kind="fallback"))
    return segs


def _extract_gz(raw: bytes) -> List[ExtractedSegment]:
    try:
        decompressed = gzip.decompress(raw)
    except Exception as exc:
        return [ExtractedSegment(text="", kind="archive-member",
                                  origin=f"[gz decompress failed: {exc}]")]
    return _dispatch("inner.txt", decompressed, depth=1)


def _extract_zip(raw: bytes) -> List[ExtractedSegment]:
    segs: List[ExtractedSegment] = []
    try:
        z = zipfile.ZipFile(io.BytesIO(raw))
    except Exception as exc:
        return [ExtractedSegment(text="", kind="archive-member",
                                  origin=f"[zip open failed: {exc}]")]
    for i, info in enumerate(z.infolist()):
        if info.is_dir():
            continue
        if i >= 25:                                          # safety cap
            segs.append(ExtractedSegment(
                text="", origin="[archive truncated at 25 members]",
                kind="archive-member",
            ))
            break
        try:
            with z.open(info) as f:
                member_bytes = f.read(8 * 1024 * 1024)      # 8 MB per member
        except Exception as exc:
            segs.append(ExtractedSegment(
                text="", kind="archive-member",
                origin=f"[read failed: {info.filename}: {exc}]",
            ))
            continue
        inner_segs = _dispatch(info.filename, member_bytes, depth=1)
        for s in inner_segs:
            s.origin = f"zip:{info.filename}" + (f"/{s.origin}" if s.origin else "")
            segs.append(s)
    return segs


def _extract_tar(raw: bytes) -> List[ExtractedSegment]:
    segs: List[ExtractedSegment] = []
    try:
        tf = tarfile.open(fileobj=io.BytesIO(raw))
    except Exception as exc:
        return [ExtractedSegment(text="", kind="archive-member",
                                  origin=f"[tar open failed: {exc}]")]
    count = 0
    for member in tf:
        if not member.isfile():
            continue
        if count >= 25:
            segs.append(ExtractedSegment(
                text="", origin="[archive truncated at 25 members]",
                kind="archive-member",
            ))
            break
        try:
            fp = tf.extractfile(member)
            if not fp:
                continue
            member_bytes = fp.read(8 * 1024 * 1024)
        except Exception as exc:
            segs.append(ExtractedSegment(
                text="", kind="archive-member",
                origin=f"[read failed: {member.name}: {exc}]",
            ))
            continue
        for s in _dispatch(member.name, member_bytes, depth=1):
            s.origin = f"tar:{member.name}" + (f"/{s.origin}" if s.origin else "")
            segs.append(s)
        count += 1
    return segs


# --------------------------------------------------------------------------- #
# Dispatch
# --------------------------------------------------------------------------- #
_EXT_TABLE: Dict[str, Callable[[bytes], List[ExtractedSegment]]] = {
    ".docx": _extract_docx,
    ".pdf":  _extract_pdf,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
    ".html": _extract_html,
    ".htm":  _extract_html,
    ".eml":  _extract_eml,
    ".rtf":  _extract_rtf,
    ".json": _extract_json,
    ".jsonl": _extract_json,
    ".ndjson": _extract_json,
    ".yaml": _extract_yaml,
    ".yml":  _extract_yaml,
    ".csv":  lambda b: _extract_csv_tsv(b, ","),
    ".tsv":  lambda b: _extract_csv_tsv(b, "\t"),
    ".gz":   _extract_gz,
    ".zip":  _extract_zip,
    ".tar":  _extract_tar,
    ".tgz":  _extract_tar,
}


def _ext(filename: str) -> str:
    filename = (filename or "").lower()
    for compound in (".tar.gz", ".tar.bz2", ".tar.xz"):
        if filename.endswith(compound):
            return ".tgz" if compound == ".tar.gz" else compound
    idx = filename.rfind(".")
    return filename[idx:] if idx >= 0 else ""


def _dispatch(filename: str, raw: bytes, depth: int = 0) -> List[ExtractedSegment]:
    if depth > 2:                                             # archive-in-archive
        return [ExtractedSegment(
            text="", kind="archive-member",
            origin=f"[max nesting depth reached for {filename!r}]",
        )]
    ext = _ext(filename)
    handler = _EXT_TABLE.get(ext)
    if handler:
        try:
            return handler(raw)
        except Exception as exc:
            log.warning("Extractor %s raised on %s: %s", ext, filename, exc)
            return [ExtractedSegment(
                text=_extract_txt(raw), kind="fallback",
                origin=f"[{ext} extractor error, fell back to raw text: {exc}]",
            )]
    # Text-like or unknown → best-effort raw decode
    return [ExtractedSegment(text=_extract_txt(raw), kind="text")]


# --------------------------------------------------------------------------- #
# Public API
# --------------------------------------------------------------------------- #
def extract(filename: str, raw: bytes) -> ExtractionResult:
    """Universal entry point.

    Returns an ExtractionResult with one or more `ExtractedSegment`s. Never
    raises — pathological inputs still surface a `fallback` segment with the
    best-effort raw text so the miner has something to work with.
    """
    res = ExtractionResult(filename=filename or "upload", total_bytes=len(raw))
    if not raw:
        res.notes.append("empty upload")
        return res
    segs = _dispatch(filename or "", raw)
    res.segments = [s for s in segs if s is not None]
    return res


def is_supported(filename: str) -> bool:
    """Returns True if the extension has a dedicated extractor OR is a
    known text-like extension. Callers can display a friendlier error
    for other extensions if they want, but `extract()` will still try."""
    ext = _ext(filename)
    return ext in _EXT_TABLE or ext in _TEXT_LIKE_EXTS
