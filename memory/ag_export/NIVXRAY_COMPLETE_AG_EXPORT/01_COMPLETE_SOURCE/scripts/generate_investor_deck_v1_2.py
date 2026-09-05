"""
NivXRay Investor Pitch Deck — Generator (v1.3.1 · sourced from locked Master Positioning)

Produces: /app/deck_assets/NivXRay_Investor_Deck_v1_2.pptx

Design rules (owner-locked):
- Visual distinction between 🟢 TODAY (verified) and 🔵 ROADMAP (future)
- Never claim full SIEM/EDR/XDR/SOAR in present tense
- Battle-cry preserved on every slide footer
- No 5.6/10 score anywhere (belongs in DD, not deck)
- Every slide footer cites the master positioning doc section
- No emoji characters in body text — color-coded dots/squares used instead

Source of truth: /app/memory/NivXRay_Strategic_Master_Positioning.md v1.3.1 (LOCKED)
Evidence base: /app/memory/NivXRay_360_Product_Market_Posture.md
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION


# ─── Palette (dark deterministic look — evidence · trust · precision) ───────────
BG          = RGBColor(0x0B, 0x0F, 0x14)   # near-black background
INK         = RGBColor(0xF5, 0xF7, 0xFA)   # off-white primary text
MUTED       = RGBColor(0x8B, 0x94, 0xA6)   # secondary text
GOLD        = RGBColor(0xE8, 0xB6, 0x4C)   # brand accent (verdict / battle-cry)
TODAY_GREEN = RGBColor(0x22, 0xC5, 0x5E)   # verified · TODAY
ROADMAP_BLUE = RGBColor(0x38, 0xBD, 0xF8)  # roadmap / vision
DIVIDER     = RGBColor(0x1F, 0x2A, 0x37)   # subtle dividers
RED_LINE    = RGBColor(0xF8, 0x71, 0x71)   # never-say rail (used sparingly)

# ─── Layout constants (16:9 · 13.333 x 7.5 in) ─────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.7)
MARGIN_R = Inches(0.7)
FONT_HEAD = "Helvetica Neue"    # falls back to system Sans on non-Mac viewers
FONT_BODY = "Helvetica Neue"
FONT_MONO = "Menlo"

OUT_PATH = Path("/app/deck_assets/NivXRay_Investor_Deck_v1_2.pptx")
OUT_PATH.parent.mkdir(parents=True, exist_ok=True)


# ══════════════════════════════════════════════════════════════════════════════
# Helpers
# ══════════════════════════════════════════════════════════════════════════════
def new_slide(prs, layout_idx=6):
    """Blank layout · we own the composition."""
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def paint_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.line.fill.background()
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        if i > 0:
            p = tf.add_paragraph()
            p.alignment = align
        run = p.add_run()
        run.text = line
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def add_dot(slide, x, y, color, size=Inches(0.14)):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    dot.line.fill.background()
    dot.fill.solid()
    dot.fill.fore_color.rgb = color


def slide_footer(slide, page_no, total, cite):
    """Bottom rail on every slide: page · brand · citation to master doc."""
    y = Inches(7.02)
    add_rect(slide, Inches(0.7), y, Inches(11.933), Emu(9525), fill=DIVIDER)
    # left: page number
    add_text(slide, Inches(0.7), Inches(7.12), Inches(1.6), Inches(0.3),
             f"{page_no:02d} / {total:02d}",
             size=9, color=MUTED, font=FONT_MONO)
    # centre: battle-cry
    add_text(slide, Inches(4.4), Inches(7.12), Inches(4.5), Inches(0.3),
             "Verdict, cited. Every time.",
             size=9, color=GOLD, italic=True, align=PP_ALIGN.CENTER)
    # right: master-doc citation
    add_text(slide, Inches(9.0), Inches(7.12), Inches(3.633), Inches(0.3),
             f"Source: Strategic Master Positioning v1.3.1 · {cite}",
             size=8, color=MUTED, font=FONT_MONO, align=PP_ALIGN.RIGHT)


def slide_header(slide, eyebrow, headline, *, tone=INK):
    """Top block: eyebrow + big headline."""
    add_text(slide, MARGIN_L, Inches(0.55), Inches(11.5), Inches(0.3),
             eyebrow, size=10, color=MUTED, font=FONT_MONO, bold=True)
    add_text(slide, MARGIN_L, Inches(0.85), Inches(11.5), Inches(1.1),
             headline, size=32, color=tone, bold=True, font=FONT_HEAD)


def status_badge(slide, x, y, kind):
    """TODAY / ROADMAP visual chip."""
    if kind == "TODAY":
        color, label = TODAY_GREEN, "TODAY · VERIFIED"
    elif kind == "ROADMAP":
        color, label = ROADMAP_BLUE, "ROADMAP · VISION"
    else:
        color, label = MUTED, kind
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.1), Inches(0.32))
    chip.line.fill.background()
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    tf = chip.text_frame
    tf.margin_top = tf.margin_bottom = 0
    tf.margin_left = tf.margin_right = Pt(6)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = FONT_MONO
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = BG


# ══════════════════════════════════════════════════════════════════════════════
# Slide constructors
# ══════════════════════════════════════════════════════════════════════════════
def slide_00_title(prs, total):
    s = new_slide(prs)
    paint_background(s)
    # gold rule at top
    add_rect(s, Inches(0.7), Inches(0.55), Inches(0.6), Inches(0.04), fill=GOLD)
    add_text(s, Inches(0.7), Inches(0.68), Inches(6), Inches(0.35),
             "NIVXRAY", size=11, color=GOLD, font=FONT_MONO, bold=True)

    add_text(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(1.2),
             "Evidence-Driven",
             size=68, color=INK, bold=True, font=FONT_HEAD)
    add_text(s, Inches(0.7), Inches(2.8), Inches(11.9), Inches(1.2),
             "Security Investigation.",
             size=68, color=INK, bold=True, font=FONT_HEAD)

    add_text(s, Inches(0.7), Inches(4.15), Inches(11.9), Inches(0.5),
             "Deterministic-first.   AI-optional.",
             size=22, color=TODAY_GREEN, bold=True, font=FONT_HEAD)

    add_text(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(0.5),
             '"Verdict, cited. Every time."',
             size=20, color=GOLD, italic=True, font=FONT_HEAD)

    add_text(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(0.4),
             "Investigate any evidence from your existing security stack.",
             size=13, color=MUTED, font=FONT_BODY)
    add_text(s, Inches(0.7), Inches(5.9), Inches(11.9), Inches(0.4),
             "Investor Overview · Seed Round · Q1 2026",
             size=11, color=MUTED, font=FONT_MONO)

    slide_footer(s, 1, total, "§ 1 · § 2.0 hierarchy")


def slide_01_problem(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "01 · THE PROBLEM",
                 "Security teams drown in evidence.\nInvestigation remains hard.")

    add_text(s, MARGIN_L, Inches(2.35), Inches(11.9), Inches(0.4),
             "Every enterprise now generates alerts, logs, telemetry and artefacts from a dozen tools.",
             size=16, color=INK)
    add_text(s, MARGIN_L, Inches(2.85), Inches(11.9), Inches(0.4),
             "But when an incident hits, the analyst still has to reconstruct what actually happened — by hand.",
             size=16, color=INK)

    # three pain points
    top = Inches(3.9); w = Inches(3.85); gap = Inches(0.2); h = Inches(2.3)
    labels = [
        ("Fragmented evidence",
         "SIEM · EDR · XDR · cloud · identity · network — each tool tells one slice of the story."),
        ("Analyst-intensive",
         "L1/L2 analysts spend hours per incident manually pivoting, decoding and correlating."),
        ("Fragile AI shortcuts",
         "LLM-based summaries produce plausible answers that cite nothing. Not defensible."),
    ]
    x = MARGIN_L
    for title, body in labels:
        add_rect(s, x, top, w, h, fill=DIVIDER)
        add_text(s, x + Inches(0.35), top + Inches(0.3), w - Inches(0.7), Inches(0.5),
                 title, size=16, bold=True, color=INK)
        add_text(s, x + Inches(0.35), top + Inches(0.9), w - Inches(0.7), h - Inches(1.1),
                 body, size=13, color=MUTED)
        x += w + gap

    slide_footer(s, 2, total, "§ 1.1 · POSTURE § 5, 18")


def slide_02_existing_stack(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "02 · THE EXISTING STACK",
                 "The evidence already exists.")

    add_text(s, MARGIN_L, Inches(2.35), Inches(11.9), Inches(0.4),
             "Every organisation already owns EDR, XDR, SIEM, cloud, identity, network and threat-intel tools.",
             size=16, color=INK)
    add_text(s, MARGIN_L, Inches(2.85), Inches(11.9), Inches(0.4),
             "They are generating exactly the evidence a great investigation needs.",
             size=16, color=INK)

    # stack grid
    cats = [
        ("SIEM",        "Splunk · Sentinel · QRadar · Elastic"),
        ("XDR",         "XSIAM · Falcon · Sentinel XDR · Cybereason"),
        ("EDR",         "CrowdStrike · SentinelOne · Defender"),
        ("Cloud audit", "CloudTrail · Azure Activity · GCP Audit"),
        ("Identity",    "Okta · Entra ID · AD-DS"),
        ("Network",     "Zeek · NetFlow · IDS/IPS · NDR"),
        ("Threat intel","VT · AbuseIPDB · URLscan · TI feeds"),
        ("Endpoint",    "Sysmon · EDR native · WMI"),
    ]
    cols = 4
    cw = Inches(3.0); ch = Inches(1.15); gx = Inches(0.03); gy = Inches(0.15)
    ox = MARGIN_L; oy = Inches(3.8)
    for i, (name, tools) in enumerate(cats):
        r, c = divmod(i, cols)
        x = ox + c * (cw + gx)
        y = oy + r * (ch + gy)
        add_rect(s, x, y, cw, ch, fill=DIVIDER)
        add_text(s, x + Inches(0.25), y + Inches(0.14), cw - Inches(0.4), Inches(0.35),
                 name, size=13, bold=True, color=INK)
        add_text(s, x + Inches(0.25), y + Inches(0.5), cw - Inches(0.4), Inches(0.6),
                 tools, size=10, color=MUTED)

    slide_footer(s, 3, total, "§ 2.2 relationships · § 4.3 wedge diagram")


def slide_03_gap(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "03 · THE GAP",
                 "Detection is solved.\nReconstruction is not.")

    add_text(s, MARGIN_L, Inches(2.35), Inches(11.9), Inches(0.4),
             "Every existing tool answers one question: DID something happen?",
             size=16, color=INK)
    add_text(s, MARGIN_L, Inches(2.85), Inches(11.9), Inches(0.4),
             "None answer the one that matters most: WHAT actually happened, in evidence?",
             size=16, color=GOLD)

    # side-by-side comparison
    x1 = MARGIN_L
    x2 = Inches(7.0)
    top = Inches(3.9); w = Inches(5.6); h = Inches(2.4)

    add_rect(s, x1, top, w, h, fill=DIVIDER)
    add_text(s, x1 + Inches(0.3), top + Inches(0.25), w - Inches(0.6), Inches(0.4),
             "What the stack does today", size=13, bold=True, color=INK)
    for i, line in enumerate([
        "Emits alerts and telemetry",
        "Runs detections against streaming logs",
        "Presents a queue of tickets",
        "Leaves the reconstruction to the analyst",
    ]):
        add_text(s, x1 + Inches(0.3), top + Inches(0.75 + i*0.36), w - Inches(0.6), Inches(0.34),
                 f"— {line}", size=12, color=MUTED)

    add_rect(s, x2, top, w, h, fill=DIVIDER, line=GOLD)
    add_text(s, x2 + Inches(0.3), top + Inches(0.25), w - Inches(0.6), Inches(0.4),
             "What analysts still do by hand", size=13, bold=True, color=GOLD)
    for i, line in enumerate([
        "Decode obfuscated payloads",
        "Pivot across SIEM / EDR / XDR / cloud",
        "Map behaviour to ATT&CK techniques",
        "Write a defensible, cited investigation report",
    ]):
        add_text(s, x2 + Inches(0.3), top + Inches(0.75 + i*0.36), w - Inches(0.6), Inches(0.34),
                 f"— {line}", size=12, color=INK)

    slide_footer(s, 4, total, "§ 3.1 pillars · § 5")


def slide_04_nivxray_today(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "04 · NIVXRAY TODAY",
                 "Evidence-Driven Security Investigation.")

    status_badge(s, Inches(11.15), Inches(0.85), "TODAY")

    add_text(s, MARGIN_L, Inches(2.0), Inches(11.9), Inches(0.45),
             "Deterministic-first.   AI-optional.",
             size=18, color=TODAY_GREEN, bold=True)
    add_text(s, MARGIN_L, Inches(2.5), Inches(11.9), Inches(0.4),
             "A deterministic investigation layer that sits on top of your existing security stack.",
             size=14, color=INK)
    add_text(s, MARGIN_L, Inches(2.9), Inches(11.9), Inches(0.4),
             "Every finding cites its evidence. Every verdict is defensible. AI is optional augmentation, never required.",
             size=14, color=GOLD, italic=True)

    # 10 verified pillars — 5 x 2 grid
    pillars = [
        ("Deterministic first", "Rules R21 · R22 · zero LLM in critical path"),
        ("Multi-language AST",  "PowerShell · CMD · Bash · Python · JS · VBS"),
        ("Recursive decode",    "12 layers · 12 codecs"),
        ("Single-pass correlation", "ICE R21 · 1385 loc · one deterministic pass"),
        ("MITRE ATT&CK",        "154 mappings · 79 display names · code-frozen"),
        ("11-field narrative",  "Executive · Analyst · Verdict · fully deterministic"),
        ("9-card + 8-tab UI",   "Analyst Brief + L4 Investigation Session"),
        ("Evidence Explorer",   "Every row cites its source input"),
        ("NIST IR export",      "MD + PDF straight from the investigation"),
        ("Wire discipline",     "Slim + allow-list + SHA-256-only IOC policy"),
    ]
    top = Inches(3.65); cols = 5
    cw = Inches(2.42); ch = Inches(1.55); gx = Inches(0.08); gy = Inches(0.12)
    ox = MARGIN_L
    for i, (t, sub) in enumerate(pillars):
        r, c = divmod(i, cols)
        x = ox + c * (cw + gx)
        y = top + r * (ch + gy)
        add_rect(s, x, y, cw, ch, fill=DIVIDER)
        add_dot(s, x + Inches(0.22), y + Inches(0.24), TODAY_GREEN, size=Inches(0.13))
        add_text(s, x + Inches(0.55), y + Inches(0.18), cw - Inches(0.8), Inches(0.45),
                 t, size=11, bold=True, color=INK)
        add_text(s, x + Inches(0.22), y + Inches(0.68), cw - Inches(0.4), Inches(0.85),
                 sub, size=9, color=MUTED)

    slide_footer(s, 5, total, "§ 3.1 · 10 pillars")


def slide_05_how_it_works(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "05 · HOW IT WORKS",
                 "The evidence spine.")

    status_badge(s, Inches(11.15), Inches(0.85), "TODAY")

    add_text(s, MARGIN_L, Inches(2.0), Inches(11.9), Inches(0.5),
             "The same deterministic spine that runs today's wedge scales into tomorrow's platform.",
             size=15, color=MUTED, italic=True)

    # 7 stages in a horizontal pipeline
    stages = [
        "Input",
        "Canonical\nEvidence",
        "Semantic\nAnalysis",
        "Correlation",
        "Investigation",
        "Attack\nReconstruction",
        "Verdict",
    ]
    top = Inches(3.3); h = Inches(1.4)
    total_w = Inches(12.0)
    n = len(stages)
    gap = Inches(0.15)
    box_w = Emu((total_w.emu - gap.emu * (n - 1)) // n)
    ox = MARGIN_L
    for i, name in enumerate(stages):
        x = ox + i * (box_w.emu + gap.emu)
        add_rect(s, Emu(int(x)), top, box_w, h,
                 fill=DIVIDER, line=GOLD if i == n - 1 else None)
        add_text(s, Emu(int(x) + Inches(0.15).emu), top + Inches(0.28),
                 Emu(box_w.emu - Inches(0.3).emu), Inches(0.85),
                 name, size=13, bold=True, color=INK, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        # arrow between boxes
        if i < n - 1:
            arrow_x = Emu(int(x) + box_w.emu + Emu(int(gap.emu * 0.1)).emu)
            ay = top + Inches(0.62)
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                     arrow_x, ay,
                                     Emu(int(gap.emu * 0.8)), Inches(0.14))
            arr.line.fill.background()
            arr.fill.solid()
            arr.fill.fore_color.rgb = MUTED

    add_text(s, MARGIN_L, Inches(5.15), Inches(11.9), Inches(0.4),
             "Every stage is deterministic. Every output is cited to its evidence.",
             size=14, color=INK)
    add_text(s, MARGIN_L, Inches(5.6), Inches(11.9), Inches(0.4),
             "LLM overlay is optional and rate-capped. It never enters the critical decision path.",
             size=13, color=MUTED)

    # footer citation
    add_text(s, MARGIN_L, Inches(6.25), Inches(11.9), Inches(0.4),
             "Verified: services/ice/correlate.py:701 (Rule R21) · services/session/summary_narrative.py (Rule R22)",
             size=9, color=MUTED, font=FONT_MONO)

    slide_footer(s, 6, total, "§ 2.6 spine · POSTURE § 3, 10")


def slide_06_why_different(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "06 · WHY DIFFERENT",
                 "Deterministic-first. AI-optional.")

    status_badge(s, Inches(11.15), Inches(0.85), "TODAY")

    add_text(s, MARGIN_L, Inches(2.05), Inches(11.9), Inches(0.5),
             "Core investigation, correlation and verdicts stay identical if the LLM overlay is removed.",
             size=16, color=INK)
    add_text(s, MARGIN_L, Inches(2.5), Inches(11.9), Inches(0.4),
             "AI/LLMs are augmentation — never foundation, never in the critical decision path.",
             size=13, color=MUTED, italic=True)

    # side-by-side: Deterministic core vs Optional AI augmentation
    top = Inches(3.15); w = Inches(5.85); h = Inches(3.2); gap = Inches(0.15)
    x1 = MARGIN_L
    x2 = x1 + w + gap

    # left column · deterministic core
    add_rect(s, x1, top, w, h, fill=DIVIDER, line=TODAY_GREEN)
    add_rect(s, x1, top, w, Inches(0.4), fill=TODAY_GREEN)
    add_text(s, x1, top + Inches(0.06), w, Inches(0.35),
             "DETERMINISTIC CORE  ·  the identity", size=12, bold=True, color=BG,
             font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for i, item in enumerate([
        "Detection reasoning",
        "Correlation (ICE Rule R21)",
        "Investigation graph",
        "ATT&CK technique attribution",
        "Verdict + confidence",
        "Evidence provenance chain",
    ]):
        add_text(s, x1 + Inches(0.35), top + Inches(0.65 + i * 0.4),
                 w - Inches(0.7), Inches(0.35),
                 f"—  {item}", size=13, color=INK)

    # right column · optional AI
    add_rect(s, x2, top, w, h, fill=DIVIDER, line=MUTED)
    add_rect(s, x2, top, w, Inches(0.4), fill=MUTED)
    add_text(s, x2, top + Inches(0.06), w, Inches(0.35),
             "OPTIONAL AI  ·  augmentation only", size=12, bold=True, color=BG,
             font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for i, item in enumerate([
        "Summarisation",
        "Analyst assistance",
        "Hunting assistance",
        "Report drafting",
        "Natural-language queries",
        "Investigation suggestions",
    ]):
        add_text(s, x2 + Inches(0.35), top + Inches(0.65 + i * 0.4),
                 w - Inches(0.7), Inches(0.35),
                 f"—  {item}", size=13, color=MUTED)

    add_text(s, MARGIN_L, Inches(6.5), Inches(11.9), Inches(0.4),
             "If the LLM hallucinates, changes behaviour, becomes unavailable or expensive — NivXRay still works, identically.",
             size=11, color=GOLD, italic=True, align=PP_ALIGN.CENTER)

    slide_footer(s, 7, total, "§ Permanent Positioning Rule · § 30")


def slide_07_wedge(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "07 · THE WEDGE",
                 "Collect directly. Investigate independently.\nIntegrate everywhere.")

    # dual badge — today core + roadmap ingestion / downstream integration
    # (green + blue side-by-side)
    add_rect(s, Inches(9.1), Inches(0.85), Inches(1.9), Inches(0.32), fill=TODAY_GREEN)
    add_text(s, Inches(9.1), Inches(0.88), Inches(1.9), Inches(0.28),
             "CORE · TODAY", size=9, bold=True, color=BG, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(11.05), Inches(0.85), Inches(1.9), Inches(0.32), fill=ROADMAP_BLUE)
    add_text(s, Inches(11.05), Inches(0.88), Inches(1.9), Inches(0.28),
             "INGEST + I/O · ROADMAP", size=9, bold=True, color=BG, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, MARGIN_L, Inches(2.05), Inches(11.9), Inches(0.4),
             "NivXRay ingests telemetry, logs, artefacts and events directly from the environment —",
             size=13, color=INK)
    add_text(s, MARGIN_L, Inches(2.4), Inches(11.9), Inches(0.4),
             "and can also consume evidence from existing SIEM · XDR · EDR platforms.",
             size=13, color=INK)
    add_text(s, MARGIN_L, Inches(2.75), Inches(11.9), Inches(0.4),
             "Detects · correlates · investigates independently, then integrates back with SIEM · ServiceNow · SOAR.",
             size=13, color=GOLD, italic=True)

    # ── Sources row (7 domains) — roadmap-blue chips ──
    sources = ["Endpoint", "Network", "Identity", "Cloud · IAM", "Web / API", "Email · Artefacts", "Applications"]
    top = Inches(3.35)
    n = len(sources); cw = Inches(1.65); gap = Inches(0.09)
    total_w = n * cw.emu + (n - 1) * gap.emu
    ox = Emu((SLIDE_W.emu - total_w) // 2)
    for i, name in enumerate(sources):
        x = ox + i * (cw.emu + gap.emu)
        add_rect(s, Emu(int(x)), top, cw, Inches(0.5), fill=DIVIDER, line=ROADMAP_BLUE)
        add_text(s, Emu(int(x)), top + Inches(0.09), cw, Inches(0.35),
                 name, size=10, bold=True, color=ROADMAP_BLUE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # arrow down
    arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                             Inches(6.5), Inches(3.98), Inches(0.3), Inches(0.32))
    arr.line.fill.background()
    arr.fill.solid()
    arr.fill.fore_color.rgb = GOLD

    # NivXRay independent processing box — TODAY green
    nvx_w = Inches(9.5); nvx_h = Inches(1.15)
    nvx_x = Emu((SLIDE_W.emu - nvx_w.emu) // 2)
    nvx_y = Inches(4.42)
    add_rect(s, Emu(int(nvx_x)), nvx_y, nvx_w, nvx_h, fill=DIVIDER, line=TODAY_GREEN)
    add_rect(s, Emu(int(nvx_x)), nvx_y, nvx_w, Inches(0.32), fill=TODAY_GREEN)
    add_text(s, Emu(int(nvx_x)), nvx_y + Inches(0.03), nvx_w, Inches(0.28),
             "NivXRay · Independent Detection · Correlation · Investigation · Verdict",
             size=11, bold=True, color=BG, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # sub-pipeline
    sub = ["Parse", "Canonical Evidence", "Deterministic Analysis", "Correlation", "Investigation Graph", "Verdict"]
    sy = nvx_y + Inches(0.5)
    sw = Emu((nvx_w.emu - Inches(0.4).emu) // len(sub))
    sx = Emu(int(nvx_x) + Inches(0.2).emu)
    for i, stage in enumerate(sub):
        x = sx + i * sw.emu
        add_text(s, Emu(int(x)), sy, sw, Inches(0.55),
                 stage, size=9, color=INK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_MONO)

    # arrow down 2
    arr2 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                              Inches(6.5), Inches(5.65), Inches(0.3), Inches(0.32))
    arr2.line.fill.background()
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = GOLD

    # Downstream integration row — 3 chips (SIEM · ServiceNow · SOAR) — roadmap blue
    outs = ["SIEM  (Splunk · Sentinel · QRadar · Elastic)",
            "ServiceNow / ITSM",
            "SOAR / XSOAR / Automation"]
    top2 = Inches(6.1)
    n = len(outs); cw = Inches(3.85); gap = Inches(0.12)
    total_w = n * cw.emu + (n - 1) * gap.emu
    ox = Emu((SLIDE_W.emu - total_w) // 2)
    for i, name in enumerate(outs):
        x = ox + i * (cw.emu + gap.emu)
        add_rect(s, Emu(int(x)), top2, cw, Inches(0.5), fill=DIVIDER, line=ROADMAP_BLUE)
        add_text(s, Emu(int(x)), top2 + Inches(0.09), cw, Inches(0.35),
                 name, size=10, bold=True, color=ROADMAP_BLUE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # caveat strip
    add_text(s, MARGIN_L, Inches(6.75), Inches(11.9), Inches(0.28),
             "Today (verified): 8 adapters — paste · URL · docx · pdf · eml · image · zip + prose recognition.  Native telemetry ingestion + downstream SIEM/SNOW/SOAR connectors are Phase 2-4 roadmap.",
             size=9, color=MUTED, font=FONT_MONO, italic=True, align=PP_ALIGN.CENTER)

    slide_footer(s, 8, total, "§ 2.3 hierarchy · § 4 wedge · Slide 07 arch correction")


def slide_07a_gaps(prs, total):
    """Explicit market gaps · what NivXRay fills · how."""
    s = new_slide(prs); paint_background(s)
    slide_header(s, "08 · MARKET GAPS  ·  HOW NIVXRAY FILLS THEM",
                 "Six real gaps in today's SOC. Each has a concrete NivXRay answer.")

    add_text(s, MARGIN_L, Inches(2.0), Inches(11.9), Inches(0.4),
             "Enterprise SOCs spend on tools, still cannot reconstruct incidents. NivXRay closes each gap deterministically — no LLM required.",
             size=11, color=MUTED, italic=True)

    # 6-gap table · 3 columns × 2 rows
    gaps = [
        ("Analyst-skill gap",
         "L1/L2 turnover · inconsistent writeups · 6-12 mo ramp per analyst",
         "9-card brief + 11-field narrative + NIST IR export land junior-analyst work at senior-analyst quality automatically.",
         "40-60 % analyst time saved per incident"),
        ("Investigation lacking",
         "Detection is solved. Reconstruction, correlation and attribution are still manual across 5+ tools.",
         "ICE Rule R21 single-pass correlation + Investigation Knowledge Graph reconstruct attacks deterministically.",
         "Full investigation in seconds, not hours"),
        ("Evidence-not-cited",
         "AI copilot summaries are plausible but cite nothing. Not defensible in DD, audit or court.",
         "Every field in the brief traces to its source (Evidence Explorer + SHA-256 IOC + wire allow-list).",
         "100 % of findings evidence-cited"),
        ("Verdict inconsistency",
         "Same incident · two analysts · three verdicts. Same LLM copilot · two runs · two answers.",
         "Deterministic verdict (Rule R22) · zero LLM in critical path · fixture-locked equivalence harness.",
         "Reproducible verdicts · zero-drift"),
        ("Time consumption",
         "Analyst spends 30-90 min per incident on paste-decode-pivot-write-up.",
         "12-layer recursive decode + 6 AST engines + auto-map to 154 MITRE techniques run in under 5 seconds.",
         "5-30× faster paste-to-brief"),
        ("Licence & data cost",
         "SIEM/EDR spend keeps rising · LLM budgets are unpredictable · retention costs explode.",
         "Deterministic core needs zero LLM budget. Wire-slim reduces payload by ~80 %. No retention lock-in.",
         "No LLM budget required to run"),
    ]
    top = Inches(2.55); cols = 3
    cw = Inches(3.9); ch = Inches(2.15); gx = Inches(0.1); gy = Inches(0.15)
    ox = MARGIN_L
    for i, (title, gap_text, fill_text, kpi) in enumerate(gaps):
        r, c = divmod(i, cols)
        x = ox + c * (cw + gx)
        y = top + r * (ch + gy)
        add_rect(s, x, y, cw, ch, fill=DIVIDER)
        # gold strip left
        add_rect(s, x, y, Inches(0.12), ch, fill=GOLD)
        # title
        add_text(s, x + Inches(0.28), y + Inches(0.15),
                 cw - Inches(0.4), Inches(0.32),
                 title, size=12, bold=True, color=GOLD)
        # gap description
        add_text(s, x + Inches(0.28), y + Inches(0.5),
                 cw - Inches(0.4), Inches(0.6),
                 gap_text, size=9, color=MUTED, italic=True)
        # NivXRay fills
        add_text(s, x + Inches(0.28), y + Inches(1.12),
                 cw - Inches(0.4), Inches(0.75),
                 fill_text, size=9, color=INK)
        # KPI bottom
        add_text(s, x + Inches(0.28), y + Inches(1.85),
                 cw - Inches(0.4), Inches(0.28),
                 kpi, size=9, bold=True, color=TODAY_GREEN, font=FONT_MONO)

    slide_footer(s, 9, total, "§ 25-27 market · § 30 AI strategy · § 15 outputs")


def slide_07b_market(prs, total):
    """Market size + scope + gap + % coverage + use cases (with charts)."""
    s = new_slide(prs); paint_background(s)
    slide_header(s, "09 · MARKET · SCOPE · COVERAGE",
                 "The gap is real. Our coverage is honest.")

    # ── LEFT: Market allocation donut chart ──────────────────────────────
    add_text(s, MARGIN_L, Inches(2.05), Inches(6), Inches(0.35),
             "Global cybersecurity market allocation (2027 est.)",
             size=12, bold=True, color=INK)
    add_text(s, MARGIN_L, Inches(2.4), Inches(6), Inches(0.3),
             "Adjacent categories that share NivXRay's investigation buyer",
             size=10, color=MUTED, italic=True)

    chart_data = CategoryChartData()
    chart_data.categories = [
        "MSSP services (~$40B)",
        "XDR / EDR (~$12B)",
        "IR consulting (~$8B)",
        "SIEM (~$8B)",
        "SOAR (~$2B)",
        "AI SOC copilots (~$1B)",
    ]
    chart_data.add_series("Market $B", (40, 12, 8, 8, 2, 1))

    chart_shp = s.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        MARGIN_L, Inches(2.75), Inches(6.0), Inches(3.6),
        chart_data,
    )
    chart = chart_shp.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    chart.legend.font.color.rgb = INK
    # per-slice colors
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(9)
    plot.data_labels.font.bold = True
    plot.data_labels.font.color.rgb = INK
    slice_colors = [GOLD, ROADMAP_BLUE, RGBColor(0xF4, 0x8F, 0xB2),
                    RGBColor(0x9C, 0x92, 0xEE), RGBColor(0x5E, 0xEA, 0xD4), TODAY_GREEN]
    for i, pt in enumerate(plot.series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = slice_colors[i]
        pt.format.line.color.rgb = BG

    # SAM / SOM callouts to right of chart
    cx = MARGIN_L + Inches(6.15); cy = Inches(2.75)
    add_rect(s, cx, cy, Inches(6.4), Inches(0.85), fill=DIVIDER, line=GOLD)
    add_text(s, cx + Inches(0.25), cy + Inches(0.1),
             Inches(6.0), Inches(0.3),
             "SAM  ·  $2 – 5 B by 2028",
             size=13, bold=True, color=GOLD)
    add_text(s, cx + Inches(0.25), cy + Inches(0.42),
             Inches(6.0), Inches(0.42),
             "AI-assisted SOC investigation + MSSP analyst leverage + IR write-up automation",
             size=9, color=INK)

    add_rect(s, cx, cy + Inches(0.98), Inches(6.4), Inches(0.85), fill=DIVIDER, line=TODAY_GREEN)
    add_text(s, cx + Inches(0.25), cy + Inches(1.08),
             Inches(6.0), Inches(0.3),
             "SOM  ·  $80 – 150 M by 2029",
             size=13, bold=True, color=TODAY_GREEN)
    add_text(s, cx + Inches(0.25), cy + Inches(1.4),
             Inches(6.0), Inches(0.42),
             "3-yr serviceable share · 30 MSSPs + 50 enterprises · deterministic wedge",
             size=9, color=INK)

    # ── RIGHT: Domain coverage bar chart ─────────────────────────────────
    add_text(s, cx, Inches(4.75),
             Inches(6.4), Inches(0.3),
             "Domain Coverage  ·  Today (green) vs Phase 3 target (blue)",
             size=11, bold=True, color=INK)

    bar_data = CategoryChartData()
    bar_data.categories = ["Endpoint\n/Malware", "Email\n/Docs", "Identity\n/Cloud", "Network\n/DNS/WAF", "Web/API\n/DB"]
    bar_data.add_series("Today",   (40, 50, 10, 10,  5))
    bar_data.add_series("Phase 3", (90, 80, 70, 60, 60))

    bar_shp = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        cx, Inches(5.05), Inches(6.4), Inches(1.9),
        bar_data,
    )
    bar = bar_shp.chart
    bar.has_title = False
    bar.has_legend = True
    bar.legend.position = XL_LEGEND_POSITION.BOTTOM
    bar.legend.include_in_layout = False
    bar.legend.font.size = Pt(8)
    bar.legend.font.color.rgb = INK
    ser = bar.plots[0].series
    ser[0].format.fill.solid(); ser[0].format.fill.fore_color.rgb = TODAY_GREEN
    ser[1].format.fill.solid(); ser[1].format.fill.fore_color.rgb = ROADMAP_BLUE
    for cat_axis in [bar.category_axis, bar.value_axis]:
        cat_axis.tick_labels.font.size = Pt(8)
        cat_axis.tick_labels.font.color.rgb = MUTED

    slide_footer(s, 10, total, "§ 25-27 market · § 15 outputs")


def slide_07c_competition(prs, total):
    """Competitor matrix + adoption/flexibility (seller & buyer lens)."""
    s = new_slide(prs); paint_background(s)
    slide_header(s, "10 · COMPETITIVE POSITION",
                 "Deterministic, evidence-cited, integration-flexible.")

    # header sub-line
    add_text(s, MARGIN_L, Inches(2.0), Inches(11.9), Inches(0.4),
             "Where NivXRay wins today · where we intentionally do not compete · why buyers adopt without risk.",
             size=12, color=MUTED, italic=True)

    # Competitor comparison table
    top = Inches(2.5); h = Inches(2.3); w = Inches(11.9)
    add_rect(s, MARGIN_L, top, w, h, fill=DIVIDER)
    add_text(s, MARGIN_L + Inches(0.25), top + Inches(0.1),
             w - Inches(0.5), Inches(0.3),
             "Comparison  ·  what makes NivXRay different from adjacent tools",
             size=12, bold=True, color=INK)
    # column headers
    cols = ["Category", "Example", "What they do", "NivXRay difference"]
    col_x = [Inches(1.7), Inches(4.4), Inches(6.4), Inches(9.7)]
    hy = top + Inches(0.45)
    for i, c in enumerate(cols):
        add_text(s, MARGIN_L + col_x[i], hy, Inches(3.0), Inches(0.28),
                 c, size=10, bold=True, color=GOLD, font=FONT_MONO)
    # rows
    rows = [
        ("LLM copilot",       "Dropzone · Prophet · Radiant", "LLM summaries · plausible",    "Deterministic · cited · reproducible"),
        ("SIEM / XDR",        "Splunk · Sentinel · Falcon",   "Alerts · telemetry",           "Investigation reasoning on top"),
        ("Sandbox",           "Any.Run · Joe · VMRay",        "Detonate · verdict",           "Paste-time reasoning · no detonation"),
        ("Case management",   "IBM Resilient · D3",           "Ticket workflow",              "Evidence-cited reconstruction"),
        ("SOAR",              "XSOAR · Torq · Tines",         "Playbook automation",          "Verdict feeds SOAR · not a replacement"),
    ]
    ry = hy + Inches(0.3)
    for i, r in enumerate(rows):
        y = ry + i * Inches(0.28)
        for j, val in enumerate(r):
            color = INK if j < 3 else TODAY_GREEN
            weight = False if j < 3 else True
            add_text(s, MARGIN_L + col_x[j], y, Inches(3.4), Inches(0.26),
                     val, size=9, color=color, bold=weight)

    # Bottom row: Adoption + Flexibility (seller lens + buyer lens)
    top2 = Inches(4.95); h2 = Inches(1.9); cw = Inches(5.9); gap = Inches(0.15)

    # Seller lens
    add_rect(s, MARGIN_L, top2, cw, h2, fill=DIVIDER, line=GOLD)
    add_text(s, MARGIN_L + Inches(0.25), top2 + Inches(0.12),
             cw - Inches(0.5), Inches(0.32),
             "Seller lens  ·  why NivXRay lands",
             size=11, bold=True, color=GOLD, font=FONT_MONO)
    seller = [
        "Zero-integration paste-first onboarding",
        "No LLM budget required to start (AI-optional)",
        "Deterministic guarantees defensible in DD",
        "MSSP + IR wholesale motion (Phase 3)",
    ]
    for i, item in enumerate(seller):
        add_text(s, MARGIN_L + Inches(0.35),
                 top2 + Inches(0.5) + i * Inches(0.32),
                 cw - Inches(0.7), Inches(0.3),
                 f"—  {item}", size=10, color=INK)

    # Buyer lens
    bx = MARGIN_L + cw + gap
    add_rect(s, bx, top2, cw, h2, fill=DIVIDER, line=TODAY_GREEN)
    add_text(s, bx + Inches(0.25), top2 + Inches(0.12),
             cw - Inches(0.5), Inches(0.32),
             "Buyer lens  ·  why NivXRay is low-risk",
             size=11, bold=True, color=TODAY_GREEN, font=FONT_MONO)
    buyer = [
        "No rip-and-replace · keeps existing SIEM/XDR/EDR",
        "Exports = NIST standard MD/PDF · no vendor lock-in",
        "Single-tenant self-host today (privacy-friendly)",
        "REST API surface · integrates with any downstream (SIEM · SNOW · SOAR)",
    ]
    for i, item in enumerate(buyer):
        add_text(s, bx + Inches(0.35),
                 top2 + Inches(0.5) + i * Inches(0.32),
                 cw - Inches(0.7), Inches(0.3),
                 f"—  {item}", size=10, color=INK)

    slide_footer(s, 11, total, "§ 25 competitive · § 4.4 deployment modes")


def slide_08_expansion(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "11 · EXPANSION",
                 "From wedge to platform.")

    status_badge(s, Inches(11.15), Inches(0.85), "ROADMAP")

    add_text(s, MARGIN_L, Inches(2.0), Inches(11.9), Inches(0.5),
             "The investigation spine expands upward. The deterministic-first invariant never changes.",
             size=15, color=MUTED, italic=True)

    # 5-phase roadmap horizontal
    phases = [
        ("Phase 1", "0–3 mo", "Wedge", "L4 projections · multi-tenant · RBAC · XDR JSON class."),
        ("Phase 2", "3–6 mo", "Native connectors", "First XDR connectors · Sysmon · YARA/Sigma exec · SSO"),
        ("Phase 3", "6–12 mo", "Scale + moat", "Distributed workers · cross-session graph · SOC-2 T1"),
        ("Phase 4", "12–18 mo", "Cross-domain", "SOAR-lite · detection authoring · cross-tenant intel"),
        ("Phase 5", "18–36 mo", "Platform", "SIEM + EDR + XDR + SOAR + Investigation · unified"),
    ]
    top = Inches(2.95); h = Inches(3.5)
    n = len(phases); gap = Inches(0.12)
    box_w = Emu((Inches(11.9).emu - gap.emu * (n - 1)) // n)
    ox = MARGIN_L
    for i, (ph, timing, label, body) in enumerate(phases):
        x = ox + i * (box_w.emu + gap.emu)
        # phase 1-3 are near-term (deeper green tint); 4-5 are vision blue
        accent = TODAY_GREEN if i == 0 else (ROADMAP_BLUE if i >= 3 else GOLD)
        add_rect(s, Emu(int(x)), top, box_w, h, fill=DIVIDER, line=accent)
        # phase strip
        add_rect(s, Emu(int(x)), top, box_w, Inches(0.35), fill=accent)
        add_text(s, Emu(int(x)), top + Inches(0.03), box_w, Inches(0.3),
                 ph, size=11, bold=True, color=BG, font=FONT_MONO,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Emu(int(x) + Inches(0.15).emu),
                 top + Inches(0.55),
                 Emu(box_w.emu - Inches(0.3).emu), Inches(0.35),
                 timing, size=9, color=MUTED, font=FONT_MONO)
        add_text(s, Emu(int(x) + Inches(0.15).emu),
                 top + Inches(0.9),
                 Emu(box_w.emu - Inches(0.3).emu), Inches(0.5),
                 label, size=13, bold=True, color=INK)
        add_text(s, Emu(int(x) + Inches(0.15).emu),
                 top + Inches(1.5),
                 Emu(box_w.emu - Inches(0.3).emu), h - Inches(1.7),
                 body, size=10, color=MUTED)

    add_text(s, MARGIN_L, Inches(6.55), Inches(11.9), Inches(0.4),
             "Every future stage inherits the deterministic-first + evidence-provenance invariants.",
             size=11, color=ROADMAP_BLUE, italic=True, align=PP_ALIGN.CENTER)

    slide_footer(s, 12, total, "§ 6 roadmap · § 2.4 Product Loop")


def slide_09_moat(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "12 · MOAT",
                 "The combination is the moat.")

    add_text(s, MARGIN_L, Inches(2.05), Inches(11.9), Inches(0.5),
             "Not any single pillar. The five together — reproduced retroactively — is prohibitively expensive.",
             size=15, color=MUTED, italic=True)

    # 5 pillars stacked
    pillars = [
        ("Deterministic-first architecture",
         "Foundational design choice. Reproducibility · traceability · explainability."),
        ("Evidence-provenance discipline",
         "Codified from wire boundary to UI. Every claim cites its evidence."),
        ("Investigation Knowledge Graph",
         "Compounding across investigations. Network effect within MSSPs (Phase 3+)."),
        ("Curated AST / decode / ATT&CK corpus",
         "Years of fixture-locked reasoning · 154 mappings · 6 AST engines."),
        ("Accumulated investigation knowledge",
         "Governance harness enforces zero-drift. 56-file equivalence suite."),
    ]
    top = Inches(2.85); w = Inches(11.9); h = Inches(0.75); gy = Inches(0.13)
    for i, (t, body) in enumerate(pillars):
        y = top + i * (h.emu + gy.emu)
        add_rect(s, MARGIN_L, Emu(int(y)), w, h, fill=DIVIDER)
        add_rect(s, MARGIN_L, Emu(int(y)), Inches(0.15), h, fill=GOLD)
        add_text(s, MARGIN_L + Inches(0.4), Emu(int(y) + Inches(0.13).emu),
                 Inches(4.0), Inches(0.5),
                 t, size=13, bold=True, color=INK)
        add_text(s, MARGIN_L + Inches(4.5), Emu(int(y) + Inches(0.16).emu),
                 Inches(7.2), Inches(0.5),
                 body, size=12, color=MUTED)

    add_text(s, MARGIN_L, Inches(6.35), Inches(11.9), Inches(0.4),
             "Defensibility hypothesis: retrofitting all five into an LLM-first codebase requires re-architecture, not a feature bolt-on.",
             size=11, color=GOLD, italic=True, align=PP_ALIGN.CENTER)

    slide_footer(s, 13, total, "§ 7 moat refinement")


def slide_10_vision(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "13 · VISION + INVESTMENT",
                 "Build the Evidence-Driven\nSecurity Operations Platform.")

    status_badge(s, Inches(11.15), Inches(0.85), "ROADMAP")

    # dual-line vision
    add_text(s, MARGIN_L, Inches(2.45), Inches(11.9), Inches(0.4),
             "Today  ·  pre-revenue technical product · verified deterministic investigation core",
             size=13, color=TODAY_GREEN, bold=True)
    add_text(s, MARGIN_L, Inches(2.8), Inches(11.9), Inches(0.4),
             "Year 1–2  ·  own the Evidence-Driven Security Investigation wedge",
             size=13, color=GOLD)
    add_text(s, MARGIN_L, Inches(3.15), Inches(11.9), Inches(0.4),
             "Year 3+   ·  expand the same evidence spine into a unified Security Operations Platform",
             size=13, color=ROADMAP_BLUE)

    # small "You Are Here" pointer
    add_text(s, MARGIN_L, Inches(3.5), Inches(11.9), Inches(0.3),
             "▼  You are here",
             size=9, color=TODAY_GREEN, font=FONT_MONO, italic=True)

    # target scoreboard 4 cols — Phase 0 TODAY + Phase 1-3 funding-accelerated
    cols = [
        ("Phase 0 · TODAY",      ["Verified · Feb 2026",
                                   "608 tests · 154 MITRE",
                                   "6 AST · 12-layer decode",
                                   "8 adapters · 7 OSINT",
                                   "9-card + 8-tab L4",
                                   "1,448 commits · 88 ADRs"]),
        ("Phase 1 · 0–6 mo",     ["Wedge · P0 execution",
                                   "3–5 design partners",
                                   "First 2 native connectors",
                                   "$200k–$500k ARR",
                                   "Multi-tenant + RBAC live",
                                   "SOC-2 T1 kickoff"]),
        ("Phase 2 · 6–12 mo",    ["Expansion · scale",
                                   "15–25 customers",
                                   "6+ connectors · YARA/Sigma exec",
                                   "$1M–$2.5M ARR",
                                   "20–30 team · MSSP wholesale",
                                   "SOC-2 T1 complete"]),
        ("Phase 3 · 12–18 mo",   ["Platform trajectory",
                                   "30–50 customers",
                                   "Distributed · x-session graph",
                                   "$3M–$7M ARR",
                                   "SOAR-lite · negative explain",
                                   "SOC-2 T2 in progress"]),
    ]
    top = Inches(3.85); n = 4; gap = Inches(0.15)
    cw = Emu((Inches(11.9).emu - gap.emu * (n - 1)) // n)
    ox = MARGIN_L
    # 4-color strip: Phase 0 today-green, Phase 1 today-green, Phase 2 gold, Phase 3 roadmap-blue
    strip_colors = [TODAY_GREEN, TODAY_GREEN, GOLD, ROADMAP_BLUE]
    for i, (title, items) in enumerate(cols):
        x = ox + i * (cw.emu + gap.emu)
        add_rect(s, Emu(int(x)), top, cw, Inches(2.7), fill=DIVIDER)
        add_rect(s, Emu(int(x)), top, cw, Inches(0.4), fill=strip_colors[i])
        add_text(s, Emu(int(x)), top + Inches(0.06), cw, Inches(0.35),
                 title, size=11, bold=True, color=BG, font=FONT_MONO,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        for j, item in enumerate(items):
            add_text(s, Emu(int(x) + Inches(0.25).emu),
                     top + Inches(0.6 + j * 0.4),
                     Emu(cw.emu - Inches(0.5).emu), Inches(0.35),
                     f"—  {item}", size=10, color=INK)

    # ask + closing
    add_text(s, MARGIN_L, Inches(6.4), Inches(11.9), Inches(0.35),
             "Strong seed → 6-month execution. Certain floors (SOC-2 audit windows, enterprise sales cycles) are irreducible with capital.",
             size=10, color=MUTED, font=FONT_MONO, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN_L, Inches(6.7), Inches(11.9), Inches(0.35),
             "Raising seed round to fund Phase 1 + Phase 2 execution.",
             size=13, color=GOLD, align=PP_ALIGN.CENTER)

    slide_footer(s, 14, total, "§ 9 phase-funding scoreboard")


def slide_11_close(prs, total):
    s = new_slide(prs); paint_background(s)

    add_rect(s, Inches(0.7), Inches(0.55), Inches(0.6), Inches(0.04), fill=GOLD)
    add_text(s, Inches(0.7), Inches(0.68), Inches(6), Inches(0.35),
             "NIVXRAY", size=11, color=GOLD, font=FONT_MONO, bold=True)

    add_text(s, Inches(0.7), Inches(2.5), Inches(11.9), Inches(1.3),
             '"Verdict, cited.',
             size=76, bold=True, color=INK, font=FONT_HEAD)
    add_text(s, Inches(0.7), Inches(3.5), Inches(11.9), Inches(1.3),
             'Every time."',
             size=76, bold=True, color=GOLD, font=FONT_HEAD)

    add_text(s, Inches(0.7), Inches(5.1), Inches(11.9), Inches(0.5),
             "Thank you.",
             size=22, color=INK, italic=True)
    add_text(s, Inches(0.7), Inches(5.65), Inches(11.9), Inches(0.5),
             "Deterministic-first · Evidence-cited · Investigation-ready.",
             size=14, color=MUTED)

    slide_footer(s, 15, total, "§ 1 battle-cry")


# ══════════════════════════════════════════════════════════════════════════════
# Build
# ══════════════════════════════════════════════════════════════════════════════
def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = 15
    slide_00_title(prs, total)
    slide_01_problem(prs, total)
    slide_02_existing_stack(prs, total)
    slide_03_gap(prs, total)
    slide_04_nivxray_today(prs, total)
    slide_05_how_it_works(prs, total)
    slide_06_why_different(prs, total)
    slide_07_wedge(prs, total)
    slide_07a_gaps(prs, total)
    slide_07b_market(prs, total)
    slide_07c_competition(prs, total)
    slide_08_expansion(prs, total)
    slide_09_moat(prs, total)
    slide_10_vision(prs, total)
    slide_11_close(prs, total)
    prs.save(str(OUT_PATH))
    return OUT_PATH


if __name__ == "__main__":
    p = build()
    print(f"OK · wrote {p} · {p.stat().st_size:,} bytes")
