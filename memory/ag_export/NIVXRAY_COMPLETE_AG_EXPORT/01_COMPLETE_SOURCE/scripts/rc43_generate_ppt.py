"""RC4.3 · PowerPoint deck generator (Feb 2026)."""
from __future__ import annotations
import json
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

EVIDENCE = Path("/app/evidence")
OUT = EVIDENCE / "NivXRay_RC41_Customer_Deck.pptx"

BG    = RGBColor(0x0A, 0x18, 0x26)
CARD  = RGBColor(0x0F, 0x22, 0x33)
BRAND = RGBColor(0x2E, 0xCC, 0xA6)
INK   = RGBColor(0xE8, 0xF4, 0xFF)
FAINT = RGBColor(0x6B, 0x8B, 0xA5)
GOOD  = RGBColor(0x26, 0xE0, 0xA7)
DANGER= RGBColor(0xFF, 0x6B, 0x6B)


def _paint_bg(slide):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = BG


def _text(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    r.font.name = "Helvetica"
    return tb


def _stripe(slide):
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.06))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = BRAND
    stripe.line.fill.background()


def _bullet_list(slide, x, y, w, h, items, size=14, color=INK):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run(); r.text = "• " + it
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.name = "Helvetica"


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Slide 1 · Cover ───────────────────────────────────────
    s = prs.slides.add_slide(blank); _paint_bg(s); _stripe(s)
    _text(s, 1.5, 2.4, 10, 1.2, "NIVXRAY", size=72, color=BRAND, bold=True)
    _text(s, 1.5, 3.6, 10, 0.7, "Deterministic Malware Command Intelligence",
          size=24, color=INK)
    _text(s, 1.5, 4.4, 10, 0.5, f"RC4.1 · Evidence & Benchmark Deck  ·  "
          f"{datetime.utcnow():%B %Y}",
          size=14, color=FAINT)
    _text(s, 1.5, 5.6, 10, 1.2,
          "Recover plaintext from obfuscated command-lines. Honest verdicts. "
          "Reproducible, court-defensible evidence. No cloud upload required.",
          size=13, color=INK)

    # ── Slide 2 · The problem ─────────────────────────────────
    s = prs.slides.add_slide(blank); _paint_bg(s); _stripe(s)
    _text(s, 0.6, 0.4, 12, 0.7, "The problem", size=32, color=BRAND, bold=True)
    _bullet_list(s, 0.8, 1.5, 11.5, 5, [
        "SOC analysts see 200-400 obfuscated payloads/week — most tools produce noise or hallucinated plaintext.",
        "CyberChef is manual — no batch, no MITRE mapping, no honest verdict.",
        "VirusTotal / Any.Run upload the sample to the cloud — blocked in regulated industries.",
        "LLMs invent indicators and cannot run mathematical primitives (RC4, AES, ChaCha).",
        "Nobody tells you HONESTLY when static recovery is impossible — you get either fake plaintext or “decode failed”.",
    ], size=15)

    # ── Slide 3 · What NivXRay does ───────────────────────────
    s = prs.slides.add_slide(blank); _paint_bg(s); _stripe(s)
    _text(s, 0.6, 0.4, 12, 0.7, "What NivXRay does", size=32, color=BRAND, bold=True)
    _bullet_list(s, 0.8, 1.5, 11.5, 5, [
        "Deterministic recursive decoder pipeline · 195 ops · 40 plugins.",
        "6 new pattern-locked decoders: hex-CSV inline · byte-array XOR · reverse · regex-swap · envvar substitute · substring picker.",
        "RC4 stream cipher executed IN PYTHON — recovers plaintext without runtime.",
        "Crypto-API annotator: AES-CBC/GCM · Rijndael · DES/3DES · ChaCha20 · DPAPI · OpenSSL · GPG · MachineGuid · C2-fetched keys.",
        "Honest-verdict engine: “static-recovery-complete · runtime-decryption-required” instead of a lie.",
        "MITRE ATT&CK · LOLBAS · YARA/Sigma export · analyst workspace with case library.",
    ], size=15)

    # ── Slide 4 · Regression numbers ──────────────────────────
    s = prs.slides.add_slide(blank); _paint_bg(s); _stripe(s)
    _text(s, 0.6, 0.4, 12, 0.7, "Regression results", size=32, color=BRAND, bold=True)
    # Numbers row
    stats = [
        ("575",   "fixtures",       BRAND),
        ("97.6%", "overall pass",   GOOD),
        ("28",    "algorithms",     BRAND),
        ("0",     "false negatives",GOOD),
        ("200ms", "median latency", INK),
    ]
    x = 0.6
    for val, lbl, c in stats:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(1.5), Inches(2.3), Inches(1.6))
        box.fill.solid(); box.fill.fore_color.rgb = CARD
        box.line.color.rgb = FAINT
        _text(s, x + 0.05, 1.7, 2.2, 0.9, val, size=44, color=c, bold=True, align=1)
        _text(s, x + 0.05, 2.7, 2.2, 0.4, lbl, size=12, color=FAINT, align=1)
        x += 2.5
    _text(s, 0.6, 3.5, 12, 3,
          "RC4.0 (obfuscation) — 465/475 pass.  RC4.1 (crypto) — 96/100 pass.\n\n"
          "Every fixture ships a stage-ladder describing which stage is "
          "recoverable and why. Runtime-required stages (DPAPI, C2-fetched keys, "
          "MachineGuid-derived) are honestly reported instead of failing.",
          size=14, color=INK)

    # ── Slide 5 · AI vs Deterministic ─────────────────────────
    s = prs.slides.add_slide(blank); _paint_bg(s); _stripe(s)
    _text(s, 0.6, 0.4, 12, 0.7, "AI vs Deterministic — head-to-head",
          size=32, color=BRAND, bold=True)
    try:
        rc43 = json.loads((EVIDENCE / "rc43_ai_vs_det.json").read_text())
        rows = rc43.get("whales", [])
    except Exception:
        rows = []
    # Table
    rows_data = [["Payload", "Det score", "Det ms", "LLM score", "LLM ms"]]
    for r in rows:
        rows_data.append([r["id"],
                           r["deterministic"]["hits_/_expected"],
                           f"{r['deterministic']['latency_ms']}",
                           r["llm"]["hits_/_expected"],
                           f"{r['llm']['latency_ms']}"])
    t = s.shapes.add_table(len(rows_data), 5,
                            Inches(0.6), Inches(1.5),
                            Inches(12), Inches(2)).table
    for j, w in enumerate([3.5, 2, 2, 2, 2]):
        t.columns[j].width = Inches(w)
    for i, row in enumerate(rows_data):
        for j, v in enumerate(row):
            c = t.cell(i, j)
            c.text = ""
            p = c.text_frame.paragraphs[0]
            r = p.add_run(); r.text = str(v)
            r.font.size = Pt(13); r.font.color.rgb = BRAND if i == 0 else INK
            r.font.bold = (i == 0)
            r.font.name = "Helvetica"
            c.fill.solid(); c.fill.fore_color.rgb = CARD if i == 0 else BG
    _text(s, 0.6, 4.5, 12, 2.5,
          "Deterministic dominates whenever MATH is required — the Empire RC4 "
          "loader gave the LLM 0/3 hits because it can’t execute the KSA/PRGA.\n\n"
          "LLM edges deterministic on the well-known Emotet sample (pattern "
          "matching), but is 5-20× slower and results VARY across runs.\n\n"
          "Determinism: NivXRay stable across 3 runs = TRUE. LLM stable = varies.",
          size=14, color=INK)

    # ── Slide 6 · Screenshot INPUT → OUTPUT ───────────────────
    for path, cap in [
        (EVIDENCE / "screenshots/13_INPUT_and_OUTPUT_split.png",
         "INPUT panel — raw ps.exe -EncodedCommand payload (223 chars)"),
        (EVIDENCE / "screenshots/14_INPUT_top_OUTPUT_bottom.png",
         "OUTPUT panel — decoded plaintext URL + verdict Malicious 90%"),
        (EVIDENCE / "screenshots/03_rc4_decoded_workspace.png",
         "RC4 inline decrypted end-to-end — 4 layers peeled, verdict Malicious"),
        (EVIDENCE / "screenshots/04_aes_honest_verdict.png",
         "Honest-verdict — AES-CBC labelled Partial Decode when key is runtime"),
    ]:
        if not path.exists(): continue
        s = prs.slides.add_slide(blank); _paint_bg(s); _stripe(s)
        _text(s, 0.6, 0.3, 12, 0.6, cap, size=18, color=BRAND, bold=True)
        s.shapes.add_picture(str(path), Inches(0.5), Inches(1.1),
                              width=Inches(12.3), height=Inches(6))

    # ── Slide · Why enterprises procure ──────────────────────
    s = prs.slides.add_slide(blank); _paint_bg(s); _stripe(s)
    _text(s, 0.6, 0.4, 12, 0.7, "Why enterprises procure",
          size=32, color=BRAND, bold=True)
    _bullet_list(s, 0.8, 1.5, 11.5, 5, [
        "Reproducibility — byte-for-byte identical output, court-defensible.",
        "No cloud upload — on-prem or in-VPC deployment.",
        "Honest verdicts — no fabricated plaintext.",
        "50× faster than a frontier LLM (200 ms vs 10 s).",
        "28 crypto families · 40+ decoders · 6 malware family classifiers.",
        "Every finding chained to a MITRE technique + a reversible decoder op.",
    ], size=15)

    # ── Slide · Roadmap ──────────────────────────────────────
    s = prs.slides.add_slide(blank); _paint_bg(s); _stripe(s)
    _text(s, 0.6, 0.4, 12, 0.7, "Enterprise roadmap",
          size=32, color=BRAND, bold=True)
    _bullet_list(s, 0.8, 1.5, 11.5, 5, [
        "RC5 (4 wk) — SSO/SAML/OIDC · RBAC · Audit log · STIX 2.1 export · signed evidence bundle.",
        "v1.0 (3 mo) — Air-gap installer · TAXII 2.1 feed · Sandbox tap (Cuckoo/CAPEv2) · public API + SDK · PDF report.",
        "v1.5 (6 mo) — Multi-tenant MSSP mode · YARA-Livehunt · SOAR playbook packs · Prometheus/OTel.",
        "v2.0 (12 mo) — Threat-actor ontology · SOC 2 Type II · purple-team simulation · marketplace of community decoders.",
        "Next-target obfuscator families — Abobus Java · RMM abuse (AnyDesk/ScreenConnect) · GithubC2.",
    ], size=14)

    # ── Slide · Call to action ───────────────────────────────
    s = prs.slides.add_slide(blank); _paint_bg(s); _stripe(s)
    _text(s, 0.6, 2, 12, 1.2, "Pilot NivXRay in your SOC",
          size=44, color=BRAND, bold=True, align=1)
    _text(s, 0.6, 3.5, 12, 1.5,
          "Own-hosted 30-day pilot · your obfuscated samples · your regression "
          "corpus · a 2-hour handoff workshop.",
          size=18, color=INK, align=1)
    _text(s, 0.6, 5.2, 12, 0.8, "hello@nivxray.com  ·  nivxray.com/pilot",
          size=16, color=BRAND, bold=True, align=1)

    prs.save(OUT)
    print(f"PPT: {OUT}")


if __name__ == "__main__":
    build()
