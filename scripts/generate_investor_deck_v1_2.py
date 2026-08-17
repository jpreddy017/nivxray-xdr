"""
NivXRay Investor Pitch Deck — Generator (v1.2 · sourced from locked Master Positioning)

Produces: /app/deck_assets/NivXRay_Investor_Deck_v1_2.pptx

Design rules (owner-locked):
- Visual distinction between 🟢 TODAY (verified) and 🔵 ROADMAP (future)
- Never claim full SIEM/EDR/XDR/SOAR in present tense
- Battle-cry preserved on every slide footer
- No 5.6/10 score anywhere (belongs in DD, not deck)
- Every slide footer cites the master positioning doc section
- No emoji characters in body text — color-coded dots/squares used instead

Source of truth: /app/memory/NivXRay_Strategic_Master_Positioning.md v1.2 (LOCKED)
Evidence base: /app/memory/NivXRay_360_Product_Market_Posture.md
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ─── Palette (dark deterministic look — evidence · trust · precision) ───────────
BG          = RGBColor(0x0B, 0x0F, 0x14)   # near-black background
INK         = RGBColor(0xF5, 0xF7, 0xFA)   # off-white primary text
MUTED       = RGBColor(0x8B, 0x94, 0xA6)   # secondary text
GOLD        = RGBColor(0xE8, 0xB6, 0x4C)   # brand accent (verdict / battle-cry)
TODAY_GREEN = RGBColor(0x22, 0xC5, 0x5E)   # verified · TODAY
ROADMAP_BLUE = RGBColor(0x38, 0xBD, 0xF8)  # roadmap / vision
DIVIDER     = RGBColor(0x1F, 0x2A, 0x37)   # subtle dividers
RED_LINE    = RGBColor(0xF8, 0x71, 0x71)   # never-say rail (used sparingly)

# ─── Layout constants (16:9 · 13.333 x 7.5 in) ─────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.7)
MARGIN_R = Inches(0.7)
FONT_HEAD = "Helvetica Neue"    # falls back to system Sans on non-Mac viewers
FONT_BODY = "Helvetica Neue"
FONT_MONO = "Menlo"

OUT_PATH = Path("/app/deck_assets/NivXRay_Investor_Deck_v1_2.pptx")
OUT_PATH.parent.mkdir(parents=True, exist_ok=True)


# ══════════════════════════════════════════════════════════════════════════════
# Helpers
# ══════════════════════════════════════════════════════════════════════════════
def new_slide(prs, layout_idx=6):
    """Blank layout · we own the composition."""
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def paint_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.line.fill.background()
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        if i > 0:
            p = tf.add_paragraph()
            p.alignment = align
        run = p.add_run()
        run.text = line
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def add_dot(slide, x, y, color, size=Inches(0.14)):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    dot.line.fill.background()
    dot.fill.solid()
    dot.fill.fore_color.rgb = color


def slide_footer(slide, page_no, total, cite):
    """Bottom rail on every slide: page · brand · citation to master doc."""
    y = Inches(7.02)
    add_rect(slide, Inches(0.7), y, Inches(11.933), Emu(9525), fill=DIVIDER)
    # left: page number
    add_text(slide, Inches(0.7), Inches(7.12), Inches(1.6), Inches(0.3),
             f"{page_no:02d} / {total:02d}",
             size=9, color=MUTED, font=FONT_MONO)
    # centre: battle-cry
    add_text(slide, Inches(4.4), Inches(7.12), Inches(4.5), Inches(0.3),
             "Verdict, cited. Every time.",
             size=9, color=GOLD, italic=True, align=PP_ALIGN.CENTER)
    # right: master-doc citation
    add_text(slide, Inches(9.0), Inches(7.12), Inches(3.633), Inches(0.3),
             f"Source: Strategic Master Positioning v1.2 · {cite}",
             size=8, color=MUTED, font=FONT_MONO, align=PP_ALIGN.RIGHT)


def slide_header(slide, eyebrow, headline, *, tone=INK):
    """Top block: eyebrow + big headline."""
    add_text(slide, MARGIN_L, Inches(0.55), Inches(11.5), Inches(0.3),
             eyebrow, size=10, color=MUTED, font=FONT_MONO, bold=True)
    add_text(slide, MARGIN_L, Inches(0.85), Inches(11.5), Inches(1.1),
             headline, size=32, color=tone, bold=True, font=FONT_HEAD)


def status_badge(slide, x, y, kind):
    """TODAY / ROADMAP visual chip."""
    if kind == "TODAY":
        color, label = TODAY_GREEN, "TODAY · VERIFIED"
    elif kind == "ROADMAP":
        color, label = ROADMAP_BLUE, "ROADMAP · VISION"
    else:
        color, label = MUTED, kind
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.1), Inches(0.32))
    chip.line.fill.background()
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    tf = chip.text_frame
    tf.margin_top = tf.margin_bottom = 0
    tf.margin_left = tf.margin_right = Pt(6)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = FONT_MONO
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = BG


# ══════════════════════════════════════════════════════════════════════════════
# Slide constructors
# ══════════════════════════════════════════════════════════════════════════════
def slide_00_title(prs, total):
    s = new_slide(prs)
    paint_background(s)
    # gold rule at top
    add_rect(s, Inches(0.7), Inches(0.55), Inches(0.6), Inches(0.04), fill=GOLD)
    add_text(s, Inches(0.7), Inches(0.68), Inches(6), Inches(0.35),
             "NIVXRAY", size=11, color=GOLD, font=FONT_MONO, bold=True)

    add_text(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(1.2),
             "Evidence-Driven",
             size=68, color=INK, bold=True, font=FONT_HEAD)
    add_text(s, Inches(0.7), Inches(2.8), Inches(11.9), Inches(1.2),
             "Security Investigation.",
             size=68, color=INK, bold=True, font=FONT_HEAD)

    add_text(s, Inches(0.7), Inches(4.15), Inches(11.9), Inches(0.5),
             "Deterministic-first.   AI-optional.",
             size=22, color=TODAY_GREEN, bold=True, font=FONT_HEAD)

    add_text(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(0.5),
             '"Verdict, cited. Every time."',
             size=20, color=GOLD, italic=True, font=FONT_HEAD)

    add_text(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(0.4),
             "Investigate any evidence from your existing security stack.",
             size=13, color=MUTED, font=FONT_BODY)
    add_text(s, Inches(0.7), Inches(5.9), Inches(11.9), Inches(0.4),
             "Investor Overview · Seed Round · Q1 2026",
             size=11, color=MUTED, font=FONT_MONO)

    slide_footer(s, 1, total, "§ 1 · § 2.0 hierarchy")


def slide_01_problem(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "01 · THE PROBLEM",
                 "Security teams drown in evidence.\nInvestigation remains hard.")

    add_text(s, MARGIN_L, Inches(2.35), Inches(11.9), Inches(0.4),
             "Every enterprise now generates alerts, logs, telemetry and artefacts from a dozen tools.",
             size=16, color=INK)
    add_text(s, MARGIN_L, Inches(2.85), Inches(11.9), Inches(0.4),
             "But when an incident hits, the analyst still has to reconstruct what actually happened — by hand.",
             size=16, color=INK)

    # three pain points
    top = Inches(3.9); w = Inches(3.85); gap = Inches(0.2); h = Inches(2.3)
    labels = [
        ("Fragmented evidence",
         "SIEM · EDR · XDR · cloud · identity · network — each tool tells one slice of the story."),
        ("Analyst-intensive",
         "L1/L2 analysts spend hours per incident manually pivoting, decoding and correlating."),
        ("Fragile AI shortcuts",
         "LLM-based summaries produce plausible answers that cite nothing. Not defensible."),
    ]
    x = MARGIN_L
    for title, body in labels:
        add_rect(s, x, top, w, h, fill=DIVIDER)
        add_text(s, x + Inches(0.35), top + Inches(0.3), w - Inches(0.7), Inches(0.5),
                 title, size=16, bold=True, color=INK)
        add_text(s, x + Inches(0.35), top + Inches(0.9), w - Inches(0.7), h - Inches(1.1),
                 body, size=13, color=MUTED)
        x += w + gap

    slide_footer(s, 2, total, "§ 1.1 · POSTURE § 5, 18")


def slide_02_existing_stack(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "02 · THE EXISTING STACK",
                 "The evidence already exists.")

    add_text(s, MARGIN_L, Inches(2.35), Inches(11.9), Inches(0.4),
             "Every organisation already owns EDR, XDR, SIEM, cloud, identity, network and threat-intel tools.",
             size=16, color=INK)
    add_text(s, MARGIN_L, Inches(2.85), Inches(11.9), Inches(0.4),
             "They are generating exactly the evidence a great investigation needs.",
             size=16, color=INK)

    # stack grid
    cats = [
        ("SIEM",        "Splunk · Sentinel · QRadar · Elastic"),
        ("XDR",         "XSIAM · Falcon · Sentinel XDR · Cybereason"),
        ("EDR",         "CrowdStrike · SentinelOne · Defender"),
        ("Cloud audit", "CloudTrail · Azure Activity · GCP Audit"),
        ("Identity",    "Okta · Entra ID · AD-DS"),
        ("Network",     "Zeek · NetFlow · IDS/IPS · NDR"),
        ("Threat intel","VT · AbuseIPDB · URLscan · TI feeds"),
        ("Endpoint",    "Sysmon · EDR native · WMI"),
    ]
    cols = 4
    cw = Inches(3.0); ch = Inches(1.15); gx = Inches(0.03); gy = Inches(0.15)
    ox = MARGIN_L; oy = Inches(3.8)
    for i, (name, tools) in enumerate(cats):
        r, c = divmod(i, cols)
        x = ox + c * (cw + gx)
        y = oy + r * (ch + gy)
        add_rect(s, x, y, cw, ch, fill=DIVIDER)
        add_text(s, x + Inches(0.25), y + Inches(0.14), cw - Inches(0.4), Inches(0.35),
                 name, size=13, bold=True, color=INK)
        add_text(s, x + Inches(0.25), y + Inches(0.5), cw - Inches(0.4), Inches(0.6),
                 tools, size=10, color=MUTED)

    slide_footer(s, 3, total, "§ 2.2 relationships · § 4.3 wedge diagram")


def slide_03_gap(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "03 · THE GAP",
                 "Detection is solved.\nReconstruction is not.")

    add_text(s, MARGIN_L, Inches(2.35), Inches(11.9), Inches(0.4),
             "Every existing tool answers one question: DID something happen?",
             size=16, color=INK)
    add_text(s, MARGIN_L, Inches(2.85), Inches(11.9), Inches(0.4),
             "None answer the one that matters most: WHAT actually happened, in evidence?",
             size=16, color=GOLD)

    # side-by-side comparison
    x1 = MARGIN_L
    x2 = Inches(7.0)
    top = Inches(3.9); w = Inches(5.6); h = Inches(2.4)

    add_rect(s, x1, top, w, h, fill=DIVIDER)
    add_text(s, x1 + Inches(0.3), top + Inches(0.25), w - Inches(0.6), Inches(0.4),
             "What the stack does today", size=13, bold=True, color=INK)
    for i, line in enumerate([
        "Emits alerts and telemetry",
        "Runs detections against streaming logs",
        "Presents a queue of tickets",
        "Leaves the reconstruction to the analyst",
    ]):
        add_text(s, x1 + Inches(0.3), top + Inches(0.75 + i*0.36), w - Inches(0.6), Inches(0.34),
                 f"— {line}", size=12, color=MUTED)

    add_rect(s, x2, top, w, h, fill=DIVIDER, line=GOLD)
    add_text(s, x2 + Inches(0.3), top + Inches(0.25), w - Inches(0.6), Inches(0.4),
             "What analysts still do by hand", size=13, bold=True, color=GOLD)
    for i, line in enumerate([
        "Decode obfuscated payloads",
        "Pivot across SIEM / EDR / XDR / cloud",
        "Map behaviour to ATT&CK techniques",
        "Write a defensible, cited investigation report",
    ]):
        add_text(s, x2 + Inches(0.3), top + Inches(0.75 + i*0.36), w - Inches(0.6), Inches(0.34),
                 f"— {line}", size=12, color=INK)

    slide_footer(s, 4, total, "§ 3.1 pillars · § 5")


def slide_04_nivxray_today(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "04 · NIVXRAY TODAY",
                 "Evidence-Driven Security Investigation.")

    status_badge(s, Inches(11.15), Inches(0.85), "TODAY")

    add_text(s, MARGIN_L, Inches(2.0), Inches(11.9), Inches(0.45),
             "Deterministic-first.   AI-optional.",
             size=18, color=TODAY_GREEN, bold=True)
    add_text(s, MARGIN_L, Inches(2.5), Inches(11.9), Inches(0.4),
             "A deterministic investigation layer that sits on top of your existing security stack.",
             size=14, color=INK)
    add_text(s, MARGIN_L, Inches(2.9), Inches(11.9), Inches(0.4),
             "Every finding cites its evidence. Every verdict is defensible. AI is optional augmentation, never required.",
             size=14, color=GOLD, italic=True)

    # 10 verified pillars — 5 x 2 grid
    pillars = [
        ("Deterministic first", "Rules R21 · R22 · zero LLM in critical path"),
        ("Multi-language AST",  "PowerShell · CMD · Bash · Python · JS · VBS"),
        ("Recursive decode",    "12 layers · 12 codecs"),
        ("Single-pass correlation", "ICE R21 · 1385 loc · one deterministic pass"),
        ("MITRE ATT&CK",        "154 mappings · 79 display names · code-frozen"),
        ("11-field narrative",  "Executive · Analyst · Verdict · fully deterministic"),
        ("9-card + 8-tab UI",   "Analyst Brief + L4 Investigation Session"),
        ("Evidence Explorer",   "Every row cites its source input"),
        ("NIST IR export",      "MD + PDF straight from the investigation"),
        ("Wire discipline",     "Slim + allow-list + SHA-256-only IOC policy"),
    ]
    top = Inches(3.65); cols = 5
    cw = Inches(2.42); ch = Inches(1.55); gx = Inches(0.08); gy = Inches(0.12)
    ox = MARGIN_L
    for i, (t, sub) in enumerate(pillars):
        r, c = divmod(i, cols)
        x = ox + c * (cw + gx)
        y = top + r * (ch + gy)
        add_rect(s, x, y, cw, ch, fill=DIVIDER)
        add_dot(s, x + Inches(0.22), y + Inches(0.24), TODAY_GREEN, size=Inches(0.13))
        add_text(s, x + Inches(0.55), y + Inches(0.18), cw - Inches(0.8), Inches(0.45),
                 t, size=11, bold=True, color=INK)
        add_text(s, x + Inches(0.22), y + Inches(0.68), cw - Inches(0.4), Inches(0.85),
                 sub, size=9, color=MUTED)

    slide_footer(s, 5, total, "§ 3.1 · 10 pillars")


def slide_05_how_it_works(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "05 · HOW IT WORKS",
                 "The evidence spine.")

    status_badge(s, Inches(11.15), Inches(0.85), "TODAY")

    add_text(s, MARGIN_L, Inches(2.0), Inches(11.9), Inches(0.5),
             "The same deterministic spine that runs today's wedge scales into tomorrow's platform.",
             size=15, color=MUTED, italic=True)

    # 7 stages in a horizontal pipeline
    stages = [
        "Input",
        "Canonical\nEvidence",
        "Semantic\nAnalysis",
        "Correlation",
        "Investigation",
        "Attack\nReconstruction",
        "Verdict",
    ]
    top = Inches(3.3); h = Inches(1.4)
    total_w = Inches(12.0)
    n = len(stages)
    gap = Inches(0.15)
    box_w = Emu((total_w.emu - gap.emu * (n - 1)) // n)
    ox = MARGIN_L
    for i, name in enumerate(stages):
        x = ox + i * (box_w.emu + gap.emu)
        add_rect(s, Emu(int(x)), top, box_w, h,
                 fill=DIVIDER, line=GOLD if i == n - 1 else None)
        add_text(s, Emu(int(x) + Inches(0.15).emu), top + Inches(0.28),
                 Emu(box_w.emu - Inches(0.3).emu), Inches(0.85),
                 name, size=13, bold=True, color=INK, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        # arrow between boxes
        if i < n - 1:
            arrow_x = Emu(int(x) + box_w.emu + Emu(int(gap.emu * 0.1)).emu)
            ay = top + Inches(0.62)
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                     arrow_x, ay,
                                     Emu(int(gap.emu * 0.8)), Inches(0.14))
            arr.line.fill.background()
            arr.fill.solid()
            arr.fill.fore_color.rgb = MUTED

    add_text(s, MARGIN_L, Inches(5.15), Inches(11.9), Inches(0.4),
             "Every stage is deterministic. Every output is cited to its evidence.",
             size=14, color=INK)
    add_text(s, MARGIN_L, Inches(5.6), Inches(11.9), Inches(0.4),
             "LLM overlay is optional and rate-capped. It never enters the critical decision path.",
             size=13, color=MUTED)

    # footer citation
    add_text(s, MARGIN_L, Inches(6.25), Inches(11.9), Inches(0.4),
             "Verified: services/ice/correlate.py:701 (Rule R21) · services/session/summary_narrative.py (Rule R22)",
             size=9, color=MUTED, font=FONT_MONO)

    slide_footer(s, 6, total, "§ 2.6 spine · POSTURE § 3, 10")


def slide_06_why_different(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "06 · WHY DIFFERENT",
                 "Deterministic-first. AI-optional.")

    status_badge(s, Inches(11.15), Inches(0.85), "TODAY")

    add_text(s, MARGIN_L, Inches(2.05), Inches(11.9), Inches(0.5),
             "Core investigation, correlation and verdicts stay identical if the LLM overlay is removed.",
             size=16, color=INK)
    add_text(s, MARGIN_L, Inches(2.5), Inches(11.9), Inches(0.4),
             "AI/LLMs are augmentation — never foundation, never in the critical decision path.",
             size=13, color=MUTED, italic=True)

    # side-by-side: Deterministic core vs Optional AI augmentation
    top = Inches(3.15); w = Inches(5.85); h = Inches(3.2); gap = Inches(0.15)
    x1 = MARGIN_L
    x2 = x1 + w + gap

    # left column · deterministic core
    add_rect(s, x1, top, w, h, fill=DIVIDER, line=TODAY_GREEN)
    add_rect(s, x1, top, w, Inches(0.4), fill=TODAY_GREEN)
    add_text(s, x1, top + Inches(0.06), w, Inches(0.35),
             "DETERMINISTIC CORE  ·  the identity", size=12, bold=True, color=BG,
             font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for i, item in enumerate([
        "Detection reasoning",
        "Correlation (ICE Rule R21)",
        "Investigation graph",
        "ATT&CK technique attribution",
        "Verdict + confidence",
        "Evidence provenance chain",
    ]):
        add_text(s, x1 + Inches(0.35), top + Inches(0.65 + i * 0.4),
                 w - Inches(0.7), Inches(0.35),
                 f"—  {item}", size=13, color=INK)

    # right column · optional AI
    add_rect(s, x2, top, w, h, fill=DIVIDER, line=MUTED)
    add_rect(s, x2, top, w, Inches(0.4), fill=MUTED)
    add_text(s, x2, top + Inches(0.06), w, Inches(0.35),
             "OPTIONAL AI  ·  augmentation only", size=12, bold=True, color=BG,
             font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for i, item in enumerate([
        "Summarisation",
        "Analyst assistance",
        "Hunting assistance",
        "Report drafting",
        "Natural-language queries",
        "Investigation suggestions",
    ]):
        add_text(s, x2 + Inches(0.35), top + Inches(0.65 + i * 0.4),
                 w - Inches(0.7), Inches(0.35),
                 f"—  {item}", size=13, color=MUTED)

    add_text(s, MARGIN_L, Inches(6.5), Inches(11.9), Inches(0.4),
             "If the LLM hallucinates, changes behaviour, becomes unavailable or expensive — NivXRay still works, identically.",
             size=11, color=GOLD, italic=True, align=PP_ALIGN.CENTER)

    slide_footer(s, 7, total, "§ Permanent Positioning Rule · § 30")


def slide_07_wedge(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "07 · THE WEDGE",
                 "Sit above the stack. No rip-and-replace.")

    status_badge(s, Inches(11.15), Inches(0.85), "TODAY")

    add_text(s, MARGIN_L, Inches(2.0), Inches(11.9), Inches(0.55),
             "Give NivXRay the evidence from your existing tools.",
             size=18, color=INK)
    add_text(s, MARGIN_L, Inches(2.5), Inches(11.9), Inches(0.55),
             "It reconstructs and investigates what actually happened. Deterministic. ATT&CK-mapped. NIST-report-ready.",
             size=15, color=MUTED)

    # centre wedge diagram — top row: stack; middle: NivXRay; bottom: outputs
    # Top stack row (5 chips)
    stack = ["SIEM", "XDR", "EDR", "Cloud · IAM", "Network · Identity"]
    top = Inches(3.55)
    n = len(stack); cw = Inches(2.2); gap = Inches(0.1)
    total_w = n * cw.emu + (n - 1) * gap.emu
    ox = Emu((SLIDE_W.emu - total_w) // 2)
    for i, name in enumerate(stack):
        x = ox + i * (cw.emu + gap.emu)
        add_rect(s, Emu(int(x)), top, cw, Inches(0.5), fill=DIVIDER)
        add_text(s, Emu(int(x)), top + Inches(0.1), cw, Inches(0.35),
                 name, size=11, bold=True, color=INK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # arrow down
    arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                             Inches(6.5), Inches(4.15), Inches(0.3), Inches(0.35))
    arr.line.fill.background()
    arr.fill.solid()
    arr.fill.fore_color.rgb = GOLD

    # NivXRay box
    nvx_w = Inches(4.5); nvx_h = Inches(0.75)
    nvx_x = Emu((SLIDE_W.emu - nvx_w.emu) // 2)
    nvx_y = Inches(4.6)
    add_rect(s, Emu(int(nvx_x)), nvx_y, nvx_w, nvx_h, fill=GOLD)
    add_text(s, Emu(int(nvx_x)), nvx_y + Inches(0.14), nvx_w, Inches(0.5),
             "NivXRay · Investigation Layer", size=14, bold=True, color=BG,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # arrow down 2
    arr2 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                              Inches(6.5), Inches(5.45), Inches(0.3), Inches(0.35))
    arr2.line.fill.background()
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = GOLD

    # Outputs row
    outs = ["9-card Analyst Brief", "8-tab L4 Session", "NIST IR Report"]
    top2 = Inches(5.95)
    n = len(outs); cw = Inches(3.2); gap = Inches(0.2)
    total_w = n * cw.emu + (n - 1) * gap.emu
    ox = Emu((SLIDE_W.emu - total_w) // 2)
    for i, name in enumerate(outs):
        x = ox + i * (cw.emu + gap.emu)
        add_rect(s, Emu(int(x)), top2, cw, Inches(0.5), fill=DIVIDER)
        add_text(s, Emu(int(x)), top2 + Inches(0.1), cw, Inches(0.35),
                 name, size=11, bold=True, color=TODAY_GREEN,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    slide_footer(s, 8, total, "§ 4 wedge · § 4.3 diagram")


def slide_08_expansion(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "08 · EXPANSION",
                 "From wedge to platform.")

    status_badge(s, Inches(11.15), Inches(0.85), "ROADMAP")

    add_text(s, MARGIN_L, Inches(2.0), Inches(11.9), Inches(0.5),
             "The investigation spine expands upward. The deterministic-first invariant never changes.",
             size=15, color=MUTED, italic=True)

    # 5-phase roadmap horizontal
    phases = [
        ("Phase 1", "0–3 mo", "Wedge", "L4 projections · multi-tenant · RBAC · XDR JSON class."),
        ("Phase 2", "3–6 mo", "Native connectors", "First XDR connectors · Sysmon · YARA/Sigma exec · SSO"),
        ("Phase 3", "6–12 mo", "Scale + moat", "Distributed workers · cross-session graph · SOC-2 T1"),
        ("Phase 4", "12–18 mo", "Cross-domain", "SOAR-lite · detection authoring · cross-tenant intel"),
        ("Phase 5", "18–36 mo", "Platform", "SIEM + EDR + XDR + SOAR + Investigation · unified"),
    ]
    top = Inches(2.95); h = Inches(3.5)
    n = len(phases); gap = Inches(0.12)
    box_w = Emu((Inches(11.9).emu - gap.emu * (n - 1)) // n)
    ox = MARGIN_L
    for i, (ph, timing, label, body) in enumerate(phases):
        x = ox + i * (box_w.emu + gap.emu)
        # phase 1-3 are near-term (deeper green tint); 4-5 are vision blue
        accent = TODAY_GREEN if i == 0 else (ROADMAP_BLUE if i >= 3 else GOLD)
        add_rect(s, Emu(int(x)), top, box_w, h, fill=DIVIDER, line=accent)
        # phase strip
        add_rect(s, Emu(int(x)), top, box_w, Inches(0.35), fill=accent)
        add_text(s, Emu(int(x)), top + Inches(0.03), box_w, Inches(0.3),
                 ph, size=11, bold=True, color=BG, font=FONT_MONO,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Emu(int(x) + Inches(0.15).emu),
                 top + Inches(0.55),
                 Emu(box_w.emu - Inches(0.3).emu), Inches(0.35),
                 timing, size=9, color=MUTED, font=FONT_MONO)
        add_text(s, Emu(int(x) + Inches(0.15).emu),
                 top + Inches(0.9),
                 Emu(box_w.emu - Inches(0.3).emu), Inches(0.5),
                 label, size=13, bold=True, color=INK)
        add_text(s, Emu(int(x) + Inches(0.15).emu),
                 top + Inches(1.5),
                 Emu(box_w.emu - Inches(0.3).emu), h - Inches(1.7),
                 body, size=10, color=MUTED)

    add_text(s, MARGIN_L, Inches(6.55), Inches(11.9), Inches(0.4),
             "Every future stage inherits the deterministic-first + evidence-provenance invariants.",
             size=11, color=ROADMAP_BLUE, italic=True, align=PP_ALIGN.CENTER)

    slide_footer(s, 9, total, "§ 6 roadmap · § 2.4 Product Loop")


def slide_09_moat(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "09 · MOAT",
                 "The combination is the moat.")

    add_text(s, MARGIN_L, Inches(2.05), Inches(11.9), Inches(0.5),
             "Not any single pillar. The five together — reproduced retroactively — is prohibitively expensive.",
             size=15, color=MUTED, italic=True)

    # 5 pillars stacked
    pillars = [
        ("Deterministic-first architecture",
         "Foundational design choice. Reproducibility · traceability · explainability."),
        ("Evidence-provenance discipline",
         "Codified from wire boundary to UI. Every claim cites its evidence."),
        ("Investigation Knowledge Graph",
         "Compounding across investigations. Network effect within MSSPs (Phase 3+)."),
        ("Curated AST / decode / ATT&CK corpus",
         "Years of fixture-locked reasoning · 154 mappings · 6 AST engines."),
        ("Accumulated investigation knowledge",
         "Governance harness enforces zero-drift. 56-file equivalence suite."),
    ]
    top = Inches(2.85); w = Inches(11.9); h = Inches(0.75); gy = Inches(0.13)
    for i, (t, body) in enumerate(pillars):
        y = top + i * (h.emu + gy.emu)
        add_rect(s, MARGIN_L, Emu(int(y)), w, h, fill=DIVIDER)
        add_rect(s, MARGIN_L, Emu(int(y)), Inches(0.15), h, fill=GOLD)
        add_text(s, MARGIN_L + Inches(0.4), Emu(int(y) + Inches(0.13).emu),
                 Inches(4.0), Inches(0.5),
                 t, size=13, bold=True, color=INK)
        add_text(s, MARGIN_L + Inches(4.5), Emu(int(y) + Inches(0.16).emu),
                 Inches(7.2), Inches(0.5),
                 body, size=12, color=MUTED)

    add_text(s, MARGIN_L, Inches(6.35), Inches(11.9), Inches(0.4),
             "Defensibility hypothesis: retrofitting all five into an LLM-first codebase requires re-architecture, not a feature bolt-on.",
             size=11, color=GOLD, italic=True, align=PP_ALIGN.CENTER)

    slide_footer(s, 10, total, "§ 7 moat refinement")


def slide_10_vision(prs, total):
    s = new_slide(prs); paint_background(s)
    slide_header(s, "10 · VISION + INVESTMENT",
                 "Build the Evidence-Driven\nSecurity Operations Platform.")

    status_badge(s, Inches(11.15), Inches(0.85), "ROADMAP")

    # dual-line vision
    add_text(s, MARGIN_L, Inches(2.5), Inches(11.9), Inches(0.5),
             "Year 1–2  ·  own the Evidence-Driven AI SOC Investigation wedge.",
             size=14, color=TODAY_GREEN)
    add_text(s, MARGIN_L, Inches(2.9), Inches(11.9), Inches(0.5),
             "Year 3+   ·  expand the same evidence spine into a unified Security Operations Platform.",
             size=14, color=ROADMAP_BLUE)

    # target scoreboard 3 cols
    cols = [
        ("2027 · Wedge",   ["10 paying customers", "3 MSSP partners", "2 native connectors", "$500k–$1M ARR", "SOC-2 T1 started"]),
        ("2028 · Expansion",["30 customers", "10 MSSP partners", "6 native connectors", "$2M–$5M ARR", "SOC-2 T2"]),
        ("2029 · Platform",["80 customers", "30 MSSP partners", "15 native connectors", "$10M–$25M ARR", "ISO-27001 · HIPAA-ready"]),
    ]
    top = Inches(3.6); n = 3; gap = Inches(0.2)
    cw = Emu((Inches(11.9).emu - gap.emu * (n - 1)) // n)
    ox = MARGIN_L
    for i, (title, items) in enumerate(cols):
        x = ox + i * (cw.emu + gap.emu)
        add_rect(s, Emu(int(x)), top, cw, Inches(2.7), fill=DIVIDER)
        add_rect(s, Emu(int(x)), top, cw, Inches(0.4),
                 fill=(TODAY_GREEN if i == 0 else (GOLD if i == 1 else ROADMAP_BLUE)))
        add_text(s, Emu(int(x)), top + Inches(0.06), cw, Inches(0.35),
                 title, size=12, bold=True, color=BG, font=FONT_MONO,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        for j, item in enumerate(items):
            add_text(s, Emu(int(x) + Inches(0.35).emu),
                     top + Inches(0.6 + j * 0.4),
                     Emu(cw.emu - Inches(0.7).emu), Inches(0.35),
                     f"—  {item}", size=11, color=INK)

    # ask + closing
    add_text(s, MARGIN_L, Inches(6.5), Inches(11.9), Inches(0.4),
             "Raising seed round to fund Phase 1 + Phase 2 execution.",
             size=13, color=GOLD, align=PP_ALIGN.CENTER)

    slide_footer(s, 11, total, "§ 9 dual-lens vision")


def slide_11_close(prs, total):
    s = new_slide(prs); paint_background(s)

    add_rect(s, Inches(0.7), Inches(0.55), Inches(0.6), Inches(0.04), fill=GOLD)
    add_text(s, Inches(0.7), Inches(0.68), Inches(6), Inches(0.35),
             "NIVXRAY", size=11, color=GOLD, font=FONT_MONO, bold=True)

    add_text(s, Inches(0.7), Inches(2.5), Inches(11.9), Inches(1.3),
             '"Verdict, cited.',
             size=76, bold=True, color=INK, font=FONT_HEAD)
    add_text(s, Inches(0.7), Inches(3.5), Inches(11.9), Inches(1.3),
             'Every time."',
             size=76, bold=True, color=GOLD, font=FONT_HEAD)

    add_text(s, Inches(0.7), Inches(5.1), Inches(11.9), Inches(0.5),
             "Thank you.",
             size=22, color=INK, italic=True)
    add_text(s, Inches(0.7), Inches(5.65), Inches(11.9), Inches(0.5),
             "Deterministic-first · Evidence-cited · Investigation-ready.",
             size=14, color=MUTED)

    slide_footer(s, 12, total, "§ 1 battle-cry")


# ══════════════════════════════════════════════════════════════════════════════
# Build
# ══════════════════════════════════════════════════════════════════════════════
def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = 12
    slide_00_title(prs, total)
    slide_01_problem(prs, total)
    slide_02_existing_stack(prs, total)
    slide_03_gap(prs, total)
    slide_04_nivxray_today(prs, total)
    slide_05_how_it_works(prs, total)
    slide_06_why_different(prs, total)
    slide_07_wedge(prs, total)
    slide_08_expansion(prs, total)
    slide_09_moat(prs, total)
    slide_10_vision(prs, total)
    slide_11_close(prs, total)
    prs.save(str(OUT_PATH))
    return OUT_PATH


if __name__ == "__main__":
    p = build()
    print(f"OK · wrote {p} · {p.stat().st_size:,} bytes")
