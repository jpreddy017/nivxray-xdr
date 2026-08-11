#!/usr/bin/env python3
"""Build the NivXRay customer-facing technical deck.

Produces:
  /app/frontend/public/downloads/NivXRay-Customer-Technical-Deck.pptx

Source-of-truth grounding:
  - /app/memory/adr/0010-nivxray-product-blueprint.md    (Blueprint)
  - /app/memory/adr/0007-current-state-master-snapshot.md (Audit)

Audience         : mixed — CISO/procurement (slides 1-3) + SOC/technical (4-17)
Maturity policy  : full transparency — LIVE TODAY vs UNDER-VALIDATION / ROADMAP
Design           : dark theme, 16:9, 1920×1080, native PPTX (editable)
"""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ─── Theme ─────────────────────────────────────────────────────────────
BG          = RGBColor(0x0F, 0x12, 0x18)   # near-black slate
CARD_BG     = RGBColor(0x16, 0x1B, 0x22)
BORDER      = RGBColor(0x30, 0x36, 0x3D)
FG          = RGBColor(0xE6, 0xED, 0xF3)   # off-white body
MUTED       = RGBColor(0x8B, 0x96, 0xA5)
ACCENT      = RGBColor(0x7E, 0xE7, 0x87)   # signature green
ACCENT_2    = RGBColor(0x79, 0xC0, 0xFF)   # info blue
WARN        = RGBColor(0xF0, 0xD0, 0x5E)   # amber
DANGER      = RGBColor(0xFF, 0x7B, 0x72)   # coral

FONT_HEAD   = "Segoe UI Semibold"
FONT_BODY   = "Segoe UI"
FONT_MONO   = "Consolas"

# 16:9 · 1920 × 1080 in EMU
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK = prs.slide_layouts[6]  # blank layout

# ─── Helpers ───────────────────────────────────────────────────────────
def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    return s

def text(slide, text, left, top, width, height, *,
         size=18, bold=False, color=FG, font=FONT_BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = ln
        run.font.name = font; run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color
    return tb

def rect(slide, left, top, width, height, *, fill=CARD_BG, line=BORDER, line_w=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.adjustments[0] = 0.08
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh

def add_header(slide, section, title, subtitle=None):
    # Accent tick
    tick = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.5),
                                  Inches(0.06), Inches(0.5))
    tick.line.fill.background(); tick.fill.solid(); tick.fill.fore_color.rgb = ACCENT
    text(slide, section.upper(), Inches(0.75), Inches(0.5),
         Inches(6), Inches(0.3), size=11, bold=True, color=MUTED, font=FONT_HEAD)
    text(slide, title, Inches(0.75), Inches(0.8), Inches(12), Inches(0.7),
         size=30, bold=True, color=FG, font=FONT_HEAD)
    if subtitle:
        text(slide, subtitle, Inches(0.75), Inches(1.45), Inches(12), Inches(0.5),
             size=14, color=MUTED, font=FONT_BODY)

def add_footer(slide, idx):
    text(slide, "NivXRay · Customer Technical Deck", Inches(0.75), Inches(7.05),
         Inches(6), Inches(0.3), size=9, color=MUTED)
    text(slide, f"{idx}", Inches(12.4), Inches(7.05),
         Inches(0.5), Inches(0.3), size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    text(slide, "Confidential · Draft · 2026-08-11", Inches(6.6), Inches(7.05),
         Inches(5.5), Inches(0.3), size=9, color=MUTED, align=PP_ALIGN.RIGHT)

def bullets(slide, items, left, top, width, height, *, size=14, color=FG,
            indent=0.15, gap_pt=6):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap_pt)
        if isinstance(item, tuple):
            marker, txt = item
            run1 = p.add_run(); run1.text = marker + " "
            run1.font.name = FONT_HEAD; run1.font.size = Pt(size); run1.font.bold = True
            run1.font.color.rgb = ACCENT
            run2 = p.add_run(); run2.text = txt
            run2.font.name = FONT_BODY; run2.font.size = Pt(size)
            run2.font.color.rgb = color
        else:
            run = p.add_run(); run.text = "• " + item
            run.font.name = FONT_BODY; run.font.size = Pt(size)
            run.font.color.rgb = color
    return tb

def status_pill(slide, label, left, top, *, color=ACCENT, width=Inches(1.4)):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                  width, Inches(0.32))
    pill.adjustments[0] = 0.5
    pill.line.color.rgb = color; pill.line.width = Pt(0.75)
    pill.fill.solid(); pill.fill.fore_color.rgb = BG
    tf = pill.text_frame; tf.margin_left = Emu(60000); tf.margin_right = Emu(60000)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = label
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]; r.font.name = FONT_HEAD; r.font.size = Pt(10)
    r.font.bold = True; r.font.color.rgb = color

def code_block(slide, lines, left, top, width, height, size=11):
    box = rect(slide, left, top, width, height, fill=RGBColor(0x0B, 0x10, 0x17))
    tb = slide.shapes.add_textbox(left + Emu(80000), top + Emu(80000),
                                  width - Emu(160000), height - Emu(160000))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        r = p.add_run(); r.text = ln
        r.font.name = FONT_MONO; r.font.size = Pt(size)
        r.font.color.rgb = RGBColor(0xC9, 0xD1, 0xD9)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 · Cover
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
# Left accent bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SLIDE_H)
bar.line.fill.background(); bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT

text(s, "NivXRay", Inches(0.85), Inches(2.2), Inches(11), Inches(1.4),
     size=72, bold=True, color=FG, font=FONT_HEAD)
text(s, "Evidence-Provenanced Security Investigation Platform",
     Inches(0.9), Inches(3.5), Inches(11), Inches(0.7),
     size=26, color=ACCENT, font=FONT_HEAD)
text(s, "Turn security inputs into evidence-provenanced investigations and explainable judgments.",
     Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.6),
     size=18, color=MUTED)
text(s, "Customer-Deployed · Technical Overview",
     Inches(0.9), Inches(5.8), Inches(6), Inches(0.4),
     size=13, color=MUTED, font=FONT_HEAD, bold=True)
text(s, "Version 1.0 · 2026-08-11",
     Inches(0.9), Inches(6.15), Inches(6), Inches(0.4),
     size=11, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 · Executive Summary
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "Executive", "What NivXRay Delivers",
           "For customer-hosted security operations · single platform, three pillars of value.")

cards = [
    ("Evidence over inference",
     "Every ATT&CK technique NivXRay emits is traceable to the exact rule, "
     "field, and observed value that triggered it. No unexplainable AI verdicts. "
     "No hallucinated attributions."),
    ("Analyst-centric investigation",
     "One Workspace: paste an artifact, get a decoded chain, Timeline, "
     "Query/Hunt, MITRE Attack Chain, Verdict, and a signed report. "
     "Purpose-built for a Level-2 analyst, not another SIEM console."),
    ("Deploys inside your environment",
     "Runs entirely in your infrastructure. Your data, your MongoDB, "
     "your LLM key (or none). No telemetry to any vendor."),
]
for i, (title, body) in enumerate(cards):
    left = Inches(0.75 + i * 4.2)
    top  = Inches(2.3)
    rect(s, left, top, Inches(4.0), Inches(4.2))
    text(s, title, left + Inches(0.3), top + Inches(0.3),
         Inches(3.5), Inches(0.7), size=17, bold=True, color=ACCENT, font=FONT_HEAD)
    text(s, body, left + Inches(0.3), top + Inches(1.1),
         Inches(3.5), Inches(3.0), size=13, color=FG)

text(s, "Position: NivXRay is not a SIEM, EDR, or XDR. It is the evidence-provenanced "
        "investigation layer that consumes analyst inputs and produces defensible judgments.",
     Inches(0.75), Inches(6.7), Inches(12), Inches(0.4), size=12, color=MUTED)
add_footer(s, 2)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 · The Problem
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "Executive", "The Problem We Address",
           "Modern SOC tools generate more claims than they can defend.")

problems = [
    ("Hallucinated attribution",
     "Analyst tools emit 'this is malware' without showing why. Post-incident "
     "review cannot reproduce or defend the finding."),
    ("Analyst context loss",
     "Investigators paste artifacts across 6-8 tools; evidence, hypotheses, "
     "and decode chains never live in one place."),
    ("SIEM ≠ investigation",
     "SIEMs surface alerts. They do not decode payloads, understand ATT&CK "
     "context, or produce defensible investigation reports."),
    ("Compliance-grade explainability gap",
     "Regulators and incident-response leads increasingly ask 'show me the "
     "evidence chain.' Most detection tools cannot answer."),
]
for i, (title, body) in enumerate(problems):
    row, col = divmod(i, 2)
    left = Inches(0.75 + col * 6.15)
    top  = Inches(2.3 + row * 2.3)
    rect(s, left, top, Inches(5.9), Inches(2.05))
    text(s, title, left + Inches(0.3), top + Inches(0.25),
         Inches(5.3), Inches(0.5), size=16, bold=True, color=ACCENT_2, font=FONT_HEAD)
    text(s, body, left + Inches(0.3), top + Inches(0.85),
         Inches(5.3), Inches(1.1), size=12.5, color=FG)
add_footer(s, 3)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 · What NivXRay Is (plain English)
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "Overview", "What NivXRay Is, In One Sentence")

text(s, "\u201CNivXRay is a browser-based, evidence-provenanced security investigation Workspace that "
        "transforms pasted or uploaded analyst inputs into evidence-gated MITRE ATT&CK mappings, "
        "Timelines, Query/Hunt views, Attack Chain diagrams, and signed downloadable reports.\u201D",
     Inches(0.75), Inches(2.2), Inches(11.8), Inches(2.0),
     size=20, color=FG, font=FONT_HEAD)

text(s, "The six pillars of the product", Inches(0.75), Inches(4.4),
     Inches(6), Inches(0.4), size=13, bold=True, color=ACCENT, font=FONT_HEAD)

pillars = [
    ("1", "Input Understanding", "What did the analyst give us?"),
    ("2", "Analysis",             "What does the input contain?"),
    ("3", "Evidence",             "What did we actually observe?"),
    ("4", "Investigation",        "How are those observations related?"),
    ("5", "Judgment",             "What does the evidence mean?"),
    ("6", "Analyst Experience",   "How is it consumed?"),
]
for i, (num, name, desc) in enumerate(pillars):
    left = Inches(0.75 + i * 2.05)
    top  = Inches(4.9)
    rect(s, left, top, Inches(1.9), Inches(1.7))
    text(s, num, left, top + Inches(0.15), Inches(1.9), Inches(0.6),
         size=28, bold=True, color=ACCENT, font=FONT_HEAD, align=PP_ALIGN.CENTER)
    text(s, name, left + Inches(0.1), top + Inches(0.75),
         Inches(1.7), Inches(0.4), size=12.5, bold=True, color=FG,
         align=PP_ALIGN.CENTER, font=FONT_HEAD)
    text(s, desc, left + Inches(0.1), top + Inches(1.1),
         Inches(1.7), Inches(0.5), size=10, color=MUTED, align=PP_ALIGN.CENTER)
add_footer(s, 4)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 · What NivXRay Can Analyze Today
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "Capabilities", "What NivXRay Can Analyze Today",
           "Full transparency: LIVE capabilities on the left; roadmap on the right.")

# LIVE column
rect(s, Inches(0.75), Inches(2.1), Inches(6.0), Inches(4.7))
text(s, "LIVE TODAY", Inches(0.9), Inches(2.25), Inches(3), Inches(0.4),
     size=13, bold=True, color=ACCENT, font=FONT_HEAD)
status_pill(s, "🟢 SHIPPING", Inches(4.9), Inches(2.25), color=ACCENT, width=Inches(1.6))

live_items = [
    "Command lines (PowerShell, cmd, bash, VBS, JS, Python)",
    "Multi-layer decoding (base64/hex/URL/rot13/brotli/lzma/AES/RC4 · 200+ decoders)",
    "PowerShell AST + deobfuscation (encoded, alias, backtick unwrap)",
    "Documents (PDF, DOCX, PPTX, XLSX, EML) — text extraction",
    "PE / shellcode static analysis (pefile · capstone disassembly)",
    "IOC extraction & reputation (URL, IP, domain, hash, email)",
    "Small tabular EDR CSVs (Symantec SEP shape)",
    "Evidence-gated MITRE ATT&CK mapping (P0.2)",
    "14-lane Attack Chain diagram",
    "Timeline projection & Query/Hunt filtering",
    "Deterministic verdict + Attack Story + Reports (MD/STIX/PDF)",
    "8 OSINT feeds (65K+ IOCs) · TAXII 2.1 push",
]
bullets(s, live_items, Inches(0.9), Inches(2.75), Inches(5.7), Inches(4.0), size=11.5)

# ROADMAP column
rect(s, Inches(6.85), Inches(2.1), Inches(6.0), Inches(4.7))
text(s, "UNDER VALIDATION / ROADMAP", Inches(7.0), Inches(2.25),
     Inches(4), Inches(0.4), size=13, bold=True, color=WARN, font=FONT_HEAD)
status_pill(s, "🟡 ARCHITECTURE", Inches(11.05), Inches(2.25),
            color=WARN, width=Inches(1.7))

roadmap_items = [
    "Server-side large-file ingestion (removes 256 KB cap)",
    "Sysmon / EVTX telemetry adapter",
    "Broader EDR/XDR ingest (JSON/NDJSON/CSV/Syslog/Webhook)",
    "Investigation Knowledge Graph (13 nodes · 14 edges)",
    "Verdict Engine v3 with adaptive weight profiles",
    "Event → Process → Chain → Device → Incident correlation",
    "Cross-event & cross-case correlation",
    "AI-clustered TI campaigns (e.g. via TweetFeed)",
    "Deterministic PDF re-render CI gate",
    "Air-gapped LLM deployment path",
]
bullets(s, roadmap_items, Inches(7.0), Inches(2.75), Inches(5.7),
        Inches(4.0), size=11.5, color=RGBColor(0xC9, 0xD1, 0xD9))

text(s, "Everything in the right column is architecturally coded and shadow-observed today; "
        "promotion follows an evidence-driven replay/validation gate — never a flag flip alone.",
     Inches(0.75), Inches(6.9), Inches(12), Inches(0.4), size=11, color=MUTED)
add_footer(s, 5)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 · How an Investigation Works
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "Workflow", "How an Investigation Works",
           "A single analyst paste flows through the six pillars to a defensible verdict.")

# Vertical flow
steps = [
    ("Analyst pastes / uploads",   "Command line · script · document · CSV · IOC"),
    ("Input Understanding",         "Classifies language, encoding, decoding decision"),
    ("Decode & analyze",            "200+ deterministic decoders + AST + semantic engine"),
    ("Extract evidence",            "IOCs · LOLBAS · behaviors · artifacts — with provenance"),
    ("Map to ATT&CK (P0.2 gate)",   "Emit MITRE only when {source, rule, field, value, ref} exist"),
    ("Project into Workspace",      "Timeline · Attack Chain · Query · Verdict · Report"),
    ("Analyst reviews & saves",     "Case Vault · corrections · exports · TAXII push"),
]
for i, (title, sub) in enumerate(steps):
    top = Inches(2.15 + i * 0.68)
    # number bubble
    b = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), top,
                           Inches(0.5), Inches(0.5))
    b.fill.solid(); b.fill.fore_color.rgb = ACCENT
    b.line.fill.background()
    tf = b.text_frame; tf.text = str(i + 1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]; r.font.name = FONT_HEAD; r.font.size = Pt(14)
    r.font.bold = True; r.font.color.rgb = BG
    # title
    text(s, title, Inches(1.55), top + Inches(0.02), Inches(4.5), Inches(0.35),
         size=14.5, bold=True, color=FG, font=FONT_HEAD)
    # subtitle
    text(s, sub, Inches(1.55), top + Inches(0.32), Inches(11), Inches(0.35),
         size=12, color=MUTED)

# right callout
rect(s, Inches(9.0), Inches(2.15), Inches(3.85), Inches(4.6))
text(s, "Key property", Inches(9.15), Inches(2.3), Inches(3.5), Inches(0.4),
     size=12, bold=True, color=ACCENT, font=FONT_HEAD)
text(s, "The evidence chain is the audit trail.",
     Inches(9.15), Inches(2.7), Inches(3.55), Inches(1.0),
     size=15, bold=True, color=FG, font=FONT_HEAD)
text(s, "Every MITRE technique on the analyst's screen carries the exact rule, "
        "field, and observed value that triggered it. Absence of provenance ⇒ the "
        "technique is silently dropped, not fabricated.",
     Inches(9.15), Inches(3.9), Inches(3.55), Inches(2.6),
     size=11.5, color=RGBColor(0xC9, 0xD1, 0xD9))
add_footer(s, 6)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 · Evidence-Provenance Differentiator
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "Differentiator", "Evidence-Provenance is Built In, Not Bolted On",
           "The P0.2 evidence chain: 30 CI-locked tests enforce no evidence → no MITRE.")

code_block(s, [
    "# Every MITRE technique NivXRay emits carries:",
    "{",
    "  \"technique\":     \"T1059.001 · PowerShell\",",
    "  \"source\":        \"die.narrative\",",
    "  \"event_or_rule\": \"powershell_encodedcommand_v3\",",
    "  \"field\":         \"decoded_layer_2\",",
    "  \"observed_value\": \"IEX (New-Object Net.WebClient).DownloadString(...)\",",
    "  \"evidence_ref\":  \"artifact://abc123#offset=42-118\"",
    "}",
    "",
    "# If provenance is absent, the technique is REJECTED.",
    "# No exceptions. Regression-locked by 30 tests.",
], Inches(0.75), Inches(2.15), Inches(7.5), Inches(4.5), size=12)

# right column benefits
benefits = [
    ("Defensibility",
     "Every claim on the screen can be replayed back to the exact input span."),
    ("Compliance-ready",
     "Auditors and regulators can walk the chain without vendor cooperation."),
    ("No hallucination",
     "LLM narrate is additive prose; it never introduces new techniques."),
    ("Deterministic reports",
     "Same envelope → same bytes → same SHA-256. Markdown & STIX CI-locked."),
]
for i, (t, b) in enumerate(benefits):
    top = Inches(2.15 + i * 1.15)
    rect(s, Inches(8.5), top, Inches(4.35), Inches(1.0))
    text(s, t, Inches(8.65), top + Inches(0.12), Inches(4.0), Inches(0.35),
         size=13, bold=True, color=ACCENT, font=FONT_HEAD)
    text(s, b, Inches(8.65), top + Inches(0.45), Inches(4.0), Inches(0.55),
         size=11, color=RGBColor(0xC9, 0xD1, 0xD9))
add_footer(s, 7)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 · Reference Architecture (LIVE)
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "Architecture", "Reference Architecture — LIVE Today",
           "Everything below runs inside your environment.")

# Analyst
def band(label, sub, left, top, width, height, *, color=ACCENT):
    rect(s, left, top, width, height)
    text(s, label, left + Inches(0.2), top + Inches(0.15),
         width - Inches(0.4), Inches(0.35), size=13, bold=True,
         color=color, font=FONT_HEAD)
    text(s, sub, left + Inches(0.2), top + Inches(0.5),
         width - Inches(0.4), Inches(0.8), size=10.5, color=RGBColor(0xC9, 0xD1, 0xD9))

# Three-tier stack, horizontal
band("Analyst UI (React 19)",
     "WorkspacePage · Timeline · Query · Attack Chain · Verdict · Reports",
     Inches(0.75), Inches(2.1), Inches(12.1), Inches(0.95))
band("Backend API (FastAPI · 77 routers)",
     "Authentication · Investigation · DIE · Analyze · Threat-Intel · Reports · Lab",
     Inches(0.75), Inches(3.2), Inches(12.1), Inches(0.95), color=ACCENT_2)

# Two service columns
band("Analysis Services",
     "IUE · Decoders · Semantic (PS AST) · MITRE evidence chain · IOC extractor · IDA · LOLBAS",
     Inches(0.75), Inches(4.3), Inches(6.0), Inches(1.1), color=ACCENT_2)
band("Projection Services",
     "Canonical SSOT · Verdict · Attack Chain · Attack Story · Timeline · Query · Reports (Markdown/STIX/PDF)",
     Inches(6.85), Inches(4.3), Inches(6.0), Inches(1.1), color=ACCENT_2)

# Storage
band("MongoDB (single database, inside your environment)",
     "Cases · Investigations · IOCs (65K+) · Corrections · Documents (GridFS) · Feed sync history · Lab",
     Inches(0.75), Inches(5.55), Inches(12.1), Inches(0.95), color=WARN)

# Outbound (small strip)
band("Outbound (analyst-controlled)",
     "OSINT feed pull (URLhaus · OTX · AbuseIPDB · ThreatFox · MB · Talos · CINS · MalwareBytes) · "
     "LLM narrate (your key) · TAXII push · URL enrichment (SSRF-guarded)",
     Inches(0.75), Inches(6.55), Inches(12.1), Inches(0.75), color=DANGER)
add_footer(s, 8)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 · Customer Data & Privacy
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "Data & Privacy", "Customer Data Flow & Data Residency",
           "Your environment is the perimeter. Nothing leaves without an analyst-initiated action.")

# Big centered box: customer env
env = rect(s, Inches(0.75), Inches(2.15), Inches(8.0), Inches(4.85),
           line=ACCENT, line_w=1.5)
text(s, "YOUR ENVIRONMENT", Inches(0.9), Inches(2.28), Inches(4.5),
     Inches(0.35), size=11, bold=True, color=ACCENT, font=FONT_HEAD)

inside = [
    "Analyst inputs (pasted commands, uploaded files, IOCs)",
    "MongoDB (cases, events, IOCs, corrections, GridFS docs)",
    "Decoded artifacts, PowerShell ASTs, decode chains",
    "Evidence graph, canonical SSOT, verdicts, reports",
    "Analyst corrections, sessions, benchmark history",
    "TI feed cache (IOCs pulled from OSINT providers)",
    "All logs, audit records, and application state",
]
bullets(s, [(f"▪", i) for i in inside], Inches(1.0), Inches(2.75),
        Inches(7.5), Inches(4.0), size=12.5)

# Outbound arrows
text(s, "OUTBOUND — only when the analyst or operator triggers it",
     Inches(9.1), Inches(2.28), Inches(4), Inches(0.35),
     size=11, bold=True, color=DANGER, font=FONT_HEAD)

outbound = [
    ("→", "OSINT feed pull",     "URLhaus · OTX · AbuseIPDB · ThreatFox · MalwareBazaar · MalwareBytes · Talos · CINS Army"),
    ("→", "LLM narrate (opt)",   "Prompt sent to your configured LLM provider. Off if not configured."),
    ("→", "URL enrichment",      "Analyst-triggered fetch. Private-IP / SSRF blocked."),
    ("→", "TAXII 2.1 push",      "STIX bundle to your configured TAXII endpoint. Admin-triggered."),
    ("→", "Report download",     "Analyst downloads to their browser."),
]
top = Inches(2.75)
for arrow, name, desc in outbound:
    rect(s, Inches(9.1), top, Inches(3.75), Inches(0.75), line=BORDER)
    text(s, f"{arrow}  {name}", Inches(9.25), top + Inches(0.05),
         Inches(3.5), Inches(0.3), size=11.5, bold=True, color=ACCENT_2, font=FONT_HEAD)
    text(s, desc, Inches(9.25), top + Inches(0.32),
         Inches(3.5), Inches(0.4), size=9.5, color=RGBColor(0xC9, 0xD1, 0xD9))
    top += Inches(0.83)

text(s, "The Emergent-hosted mode is optional. In customer-deployed mode, no data flows to Emergent Labs.",
     Inches(0.75), Inches(7.05), Inches(12), Inches(0.35), size=11, color=MUTED)
add_footer(s, 9)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 · Deployment Options & Maturity
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "Deployment", "Deployment Options & Current Maturity",
           "Honest state today. Production packaging is on the near-term roadmap.")

# Table headers
cols = ["Deployment Model", "Feasibility Today", "What's Needed"]
rows = [
    ("Single-tenant SaaS · Emergent-hosted",
     "🟢 Live",
     "Currently operational; ideal for evaluation only."),
    ("Customer-hosted single-node (VM)",
     "🟡 Under Validation",
     "Docker Compose packaging + `.env` per customer + operator runbook."),
    ("Customer-hosted Kubernetes",
     "🟡 Under Validation",
     "Helm chart, image registry, MongoDB StatefulSet, secret management."),
    ("Air-gapped enterprise",
     "🟠 Roadmap",
     "Offline LLM path documented; OSINT feed sync must be operator-toggled off."),
    ("Multi-tenant SaaS",
     "⚪ Not Supported",
     "Requires tenant model + per-tenant DB + SSO. Explicitly deferred."),
    ("Bring-Your-Own LLM (customer's OpenAI/Anthropic/Gemini)",
     "🟡 Under Validation",
     "LiteLLM shim already routes any provider; needs config-management flow."),
]

# Draw table
y = Inches(2.15)
row_h = Inches(0.65)
widths = [Inches(4.5), Inches(2.5), Inches(5.35)]
# header
x = Inches(0.75)
for i, c in enumerate(cols):
    rect(s, x, y, widths[i], row_h, fill=RGBColor(0x1F, 0x27, 0x30))
    text(s, c, x + Inches(0.15), y + Inches(0.17), widths[i] - Inches(0.3),
         row_h, size=12.5, bold=True, color=ACCENT, font=FONT_HEAD)
    x += widths[i]
y += row_h
for r in rows:
    x = Inches(0.75)
    for i, cell in enumerate(r):
        rect(s, x, y, widths[i], row_h)
        color = ACCENT if "🟢" in cell else (WARN if "🟡" in cell else
                (RGBColor(0xFF, 0xA6, 0x57) if "🟠" in cell else MUTED))
        text(s, cell, x + Inches(0.15), y + Inches(0.17),
             widths[i] - Inches(0.3), row_h, size=11.5,
             color=color if i == 1 else FG, font=FONT_BODY,
             bold=(i == 1))
        x += widths[i]
    y += row_h

text(s, "Standard deployment sizing (single-node): 2 vCPU · 4 GB RAM · 20 GB storage · outbound HTTPS to OSINT/LLM as configured.",
     Inches(0.75), Inches(6.95), Inches(12), Inches(0.35), size=10.5, color=MUTED)
add_footer(s, 10)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 · Component Functionality
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "Components", "Component-by-Component Functionality",
           "The engines that make up the LIVE analysis pipeline.")

comps = [
    ("Input Understanding Engine",
     "Classifies each analyst input — language, encoding, decode plan. "
     "Distinguishes prose vs script vs binary vs structured CSV."),
    ("Multi-Layer Decoder",
     "200+ deterministic decoders. Peels base64/hex/URL/rot13/gzip/lzma/brotli/AES/RC4 "
     "chains recursively to a fixed point."),
    ("Semantic Engine (PowerShell AST)",
     "AST-level deobfuscation — alias resolution, backtick unwrap, variable "
     "substitution, string join reconstruction."),
    ("MITRE Evidence Chain (P0.2)",
     "Enforces that every emitted ATT&CK technique carries {source, rule, "
     "field, observed_value, evidence_ref}."),
    ("Intelligent Document Analyzer",
     "PE/PDF/DOCX/PPTX/XLSX/EML/image parsers. Recursive artifact discovery "
     "with fixed-point termination."),
    ("Canonical SSOT & Projections",
     "Append-only truth graph plus pure-function projections (verdict · "
     "attack chain · attack story · IOC · timeline · reports)."),
    ("Threat-Intel Feed System",
     "8 providers · 65 K+ IOCs · scheduled sync + on-demand enrichment + "
     "TAXII 2.1 push."),
    ("Report Engine (deterministic)",
     "Markdown · STIX 2.1 bundle · PDF · ZIP. SHA-256 signed envelope, "
     "byte-identical re-render (Markdown & STIX CI-locked)."),
]
for i, (t, b) in enumerate(comps):
    row, col = divmod(i, 2)
    left = Inches(0.75 + col * 6.1)
    top  = Inches(2.1 + row * 1.25)
    rect(s, left, top, Inches(5.9), Inches(1.1))
    text(s, t, left + Inches(0.25), top + Inches(0.12),
         Inches(5.4), Inches(0.4), size=13, bold=True, color=ACCENT_2, font=FONT_HEAD)
    text(s, b, left + Inches(0.25), top + Inches(0.5),
         Inches(5.4), Inches(0.7), size=11, color=RGBColor(0xC9, 0xD1, 0xD9))
add_footer(s, 11)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 · Security Posture
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "Security", "NivXRay Security Posture",
           "Honest inventory: what is in place, what is being hardened next.")

# In place
rect(s, Inches(0.75), Inches(2.15), Inches(6.0), Inches(4.85))
text(s, "IN PLACE TODAY", Inches(0.9), Inches(2.28), Inches(4), Inches(0.4),
     size=13, bold=True, color=ACCENT, font=FONT_HEAD)
in_place = [
    "JWT authentication with bcrypt-hashed credentials",
    "Idempotent admin seed; forced password change (configurable)",
    "SSRF protection (loopback / RFC1918 / link-local blocked)",
    "Request body-size cap (512 KB default / 50 MB whitelisted)",
    "Per-path request timeout with X-Request-ID",
    "GZip response compression",
    "LLM prompt / response telemetry hook",
    "Motor + Pydantic model validation (NoSQL injection guard)",
    "MongoDB bound to localhost inside the pod / VM",
    "P0.3 payload firewall (10-key allowlist)",
    "3,621 test functions across 389 backend test files",
    "Determinism CI gate (Markdown + STIX byte-identical)",
]
bullets(s, [(f"✓", i) for i in in_place], Inches(0.95),
        Inches(2.75), Inches(5.7), Inches(4.2), size=11)

# Hardening in progress (P0 gate)
rect(s, Inches(6.85), Inches(2.15), Inches(6.0), Inches(4.85))
text(s, "P0 SECURITY HARDENING (IN PROGRESS)", Inches(7.0),
     Inches(2.28), Inches(5), Inches(0.4), size=13, bold=True,
     color=WARN, font=FONT_HEAD)
progress = [
    "Explicit CORS allow-list (replaces permissive wildcard)",
    "Login rate-limit / lockout on `/api/auth/login`",
    "Archive-bomb / recursion / expanded-size / file-count guards",
    "Safe failure for malformed archives (fail-loud)",
    "Regression tests locking every new guard",
    "P0 documented as customer-facing readiness gate",
]
bullets(s, [(f"→", i) for i in progress], Inches(7.05),
        Inches(2.75), Inches(5.7), Inches(2.5), size=11)

text(s, "PLANNED NEXT (P1-P2)", Inches(7.0), Inches(5.15),
     Inches(5), Inches(0.4), size=12, bold=True, color=MUTED, font=FONT_HEAD)
planned = [
    "Server-side file store (removes 256 KB paste ceiling)",
    "Subprocess / sandbox isolation for hostile-input parsers",
    "Audit log surfacing (v2_audit_log)",
    "Multi-tenant model (deferred until customer requires)",
    "SSO / OAuth (deferred until customer requires)",
]
bullets(s, [(f"·", i) for i in planned], Inches(7.05),
        Inches(5.5), Inches(5.7), Inches(1.5), size=10.5,
        color=RGBColor(0xC9, 0xD1, 0xD9))

add_footer(s, 12)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 · Integrations & TI
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "Integrations", "Threat Intelligence, Standards, and LLM",
           "How NivXRay connects to your existing ecosystem.")

# Left · TI providers (live)
rect(s, Inches(0.75), Inches(2.15), Inches(4.0), Inches(4.85))
text(s, "TI PROVIDERS — LIVE", Inches(0.9), Inches(2.28), Inches(3.5),
     Inches(0.4), size=12, bold=True, color=ACCENT, font=FONT_HEAD)
prov = [
    "AlienVault OTX (API key)",
    "AbuseIPDB (API key)",
    "URLhaus (abuse.ch)",
    "ThreatFox (abuse.ch)",
    "MalwareBazaar (abuse.ch)",
    "Feodo Tracker (abuse.ch)",
    "MalwareBytes",
    "Talos Reputation Center",
    "CINS Army",
    "SANS DShield",
    "Blocklist.de",
    "CISA KEV catalogue",
]
bullets(s, prov, Inches(0.95), Inches(2.75), Inches(3.7), Inches(4.2), size=11)

# Middle · Standards
rect(s, Inches(4.85), Inches(2.15), Inches(4.0), Inches(4.85))
text(s, "STANDARDS — LIVE", Inches(5.0), Inches(2.28),
     Inches(3.5), Inches(0.4), size=12, bold=True, color=ACCENT_2, font=FONT_HEAD)
std = [
    "STIX 2.1 export (indicators, TLP:CLEAR)",
    "TAXII 2.1 push (admin, real-world usage)",
    "MITRE ATT&CK framework mapping",
    "Sigma rule emission",
    "YARA rule emission",
    "MITRE Navigator v4.5 layer JSON",
]
bullets(s, std, Inches(5.05), Inches(2.75), Inches(3.7), Inches(4.2), size=11)
text(s, "TI PULL — ROADMAP", Inches(5.0), Inches(5.4),
     Inches(3.5), Inches(0.4), size=12, bold=True, color=WARN, font=FONT_HEAD)
pull = [
    "STIX/TAXII 2.1 poll (after determinism CI)",
    "MISP integration (bidirectional)",
    "AI-clustered campaigns (TweetFeed)",
]
bullets(s, pull, Inches(5.05), Inches(5.85),
        Inches(3.7), Inches(1.15), size=10.5, color=RGBColor(0xC9, 0xD1, 0xD9))

# Right · LLM
rect(s, Inches(8.95), Inches(2.15), Inches(3.9), Inches(4.85))
text(s, "LLM OPTIONS", Inches(9.1), Inches(2.28),
     Inches(3.5), Inches(0.4), size=12, bold=True, color=ACCENT_2, font=FONT_HEAD)
llm = [
    "Customer's Anthropic Claude key",
    "Customer's OpenAI key",
    "Customer's Google Gemini key",
    "Any provider via LiteLLM shim",
    "LLM narrate can be disabled",
    "Air-gapped LLM path (roadmap)",
]
bullets(s, llm, Inches(9.15), Inches(2.75), Inches(3.6), Inches(3.0), size=11)
text(s, "Property", Inches(9.1), Inches(5.4),
     Inches(3.5), Inches(0.3), size=11, bold=True, color=ACCENT, font=FONT_HEAD)
text(s, "LLM is used only for narrative generation. It NEVER introduces new "
        "ATT&CK techniques or verdicts. Deterministic evidence chain is "
        "independent of the LLM.",
     Inches(9.15), Inches(5.7), Inches(3.55), Inches(1.3),
     size=10, color=RGBColor(0xC9, 0xD1, 0xD9))
add_footer(s, 13)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 · Current vs Next-Generation Architecture
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "Architecture", "Current LIVE Path vs. Next-Generation (Under Validation)",
           "The next-gen architecture is coded, shadow-observed, and gated by evidence-driven promotion.")

# Left: current live
rect(s, Inches(0.75), Inches(2.15), Inches(6.0), Inches(4.85), line=ACCENT, line_w=1.25)
text(s, "CURRENT · LIVE", Inches(0.9), Inches(2.28),
     Inches(3), Inches(0.4), size=13, bold=True, color=ACCENT, font=FONT_HEAD)
status_pill(s, "🟢 SHIPPING", Inches(5.05), Inches(2.28),
            color=ACCENT, width=Inches(1.5))
cur_flow = [
    "Input (paste / upload ≤ 256 KB)",
    "  → DIE Analysis Pipeline",
    "     → Canonical SSOT (per-request, in-memory)",
    "        → Deterministic Projections",
    "           → Workspace panels + Reports",
]
bullets(s, cur_flow, Inches(0.95), Inches(2.85),
        Inches(5.7), Inches(3.0), size=12.5)
text(s, "Regression-locked · 114 canonical API tests · P0.2 evidence chain · "
        "P0.3 payload firewall · determinism CI (Markdown + STIX).",
     Inches(0.95), Inches(5.6), Inches(5.7), Inches(1.3),
     size=11, color=MUTED)

# Right: next-gen
rect(s, Inches(6.85), Inches(2.15), Inches(6.0), Inches(4.85), line=WARN, line_w=1.25)
text(s, "NEXT-GEN · UNDER VALIDATION", Inches(7.0),
     Inches(2.28), Inches(5), Inches(0.4), size=13, bold=True,
     color=WARN, font=FONT_HEAD)
status_pill(s, "🟡 SHADOW", Inches(11.35), Inches(2.28),
            color=WARN, width=Inches(1.4))
ng_flow = [
    "Universal Input Router",
    "  → Adapter Tier (Sysmon · EVTX · JSON · CSV · Syslog · Webhook)",
    "     → Canonical Event Bag (persisted)",
    "        → Investigation Knowledge Graph (IKG · 13 node · 14 edge)",
    "           → Correlation Engine (Event→Process→Chain→Device→Incident)",
    "              → Verdict Engine v3 (Adaptive Weight Profiles)",
    "                 → Attack Story · Mitigation · Reports",
]
bullets(s, ng_flow, Inches(7.05), Inches(2.85),
        Inches(5.7), Inches(3.5), size=11.5, color=RGBColor(0xC9, 0xD1, 0xD9))
text(s, "Promoted from shadow → live only via replay-parity gate. "
        "Never a flag-flip in isolation.",
     Inches(7.05), Inches(6.15), Inches(5.7), Inches(0.7),
     size=11, color=MUTED)
add_footer(s, 14)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15 · Roadmap
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "Roadmap", "The Delivery Sequence",
           "Ordered so security and ingestion foundations land before capability expansion.")

phases = [
    ("P0", "Security Hardening Gate",
     "CORS allow-list · login rate-limit · archive-bomb guards · fail-loud archive handling.",
     ACCENT),
    ("P1", "Server-Side File Mode",
     "Backend file store + provenance envelope + Input Router. Removes 256 KB ceiling.",
     ACCENT_2),
    ("P2", "Sysmon / EVTX Adapter",
     "Real telemetry into the canonical event bag; Timeline & Query consume automatically.",
     ACCENT_2),
    ("P3", "Shadow-Pipeline Replay & Promotion",
     "IKG · Verdict v3 · Correlation · Case Engine · Artifact Store lit under replay parity.",
     WARN),
    ("P4", "Broader EDR/XDR ingest + TweetFeed campaigns",
     "CrowdStrike/Defender/SentinelOne/Splunk connectors; AI-clustered TI campaigns.",
     WARN),
    ("P5", "Enterprise readiness",
     "Multi-tenant · SSO/SAML · audit trail · air-gapped installer.",
     MUTED),
]
top = Inches(2.15)
for i, (tag, title, desc, color) in enumerate(phases):
    left = Inches(0.75)
    # Tag
    tag_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                 Inches(0.9), Inches(0.75))
    tag_box.adjustments[0] = 0.3
    tag_box.fill.solid(); tag_box.fill.fore_color.rgb = color
    tag_box.line.fill.background()
    tf = tag_box.text_frame; tf.text = tag
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]; r.font.name = FONT_HEAD; r.font.size = Pt(16)
    r.font.bold = True; r.font.color.rgb = BG
    # Title + desc
    text(s, title, Inches(1.85), top + Inches(0.05),
         Inches(11.0), Inches(0.35), size=14, bold=True,
         color=color, font=FONT_HEAD)
    text(s, desc, Inches(1.85), top + Inches(0.4),
         Inches(11.0), Inches(0.4), size=11.5,
         color=RGBColor(0xC9, 0xD1, 0xD9))
    top += Inches(0.82)

text(s, "Timeline is capability-gated, not date-gated. Every phase ships with its own "
        "regression tests before the next opens.",
     Inches(0.75), Inches(7.05), Inches(12), Inches(0.35), size=10.5, color=MUTED)
add_footer(s, 15)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 16 · What NivXRay Does Not Claim
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "Boundaries", "What NivXRay Does NOT Claim",
           "Explicit boundaries. Trust is built by being clear about what we don't do.")

not_claims = [
    ("NivXRay is not a SIEM",
     "We do not ingest, retain, or query large-scale event streams. Bring your SIEM."),
    ("NivXRay is not an EDR or XDR",
     "We do not deploy endpoint agents. We do not collect telemetry ourselves."),
    ("NivXRay is not a sandbox / detonation service",
     "We do not execute suspicious binaries. Static analysis + parsers only."),
    ("NivXRay is not multi-tenant today",
     "One deployment = one customer. SaaS multi-tenant is on the roadmap only."),
    ("NivXRay does not depend on an LLM for verdicts",
     "The evidence chain is deterministic. LLM is prose narrative only, and optional."),
    ("NivXRay does not send your data to Emergent Labs in customer-hosted mode",
     "In your environment: your MongoDB, your keys, your outbound feeds. No vendor callbacks."),
    ("NivXRay does not certify compliance without your validation",
     "Explainability supports SOC2 / ISO / NIS2 evidence — your GRC team validates the fit."),
    ("NivXRay does not overpromise the shadow architecture",
     "IKG · Verdict v3 · Correlation · Adapters are labelled 'under validation' until promoted."),
]
for i, (t, b) in enumerate(not_claims):
    row, col = divmod(i, 2)
    left = Inches(0.75 + col * 6.1)
    top  = Inches(2.1 + row * 1.2)
    rect(s, left, top, Inches(5.9), Inches(1.05))
    text(s, "✗  " + t, left + Inches(0.25), top + Inches(0.12),
         Inches(5.4), Inches(0.4), size=13, bold=True,
         color=DANGER, font=FONT_HEAD)
    text(s, b, left + Inches(0.6), top + Inches(0.5),
         Inches(5.1), Inches(0.55), size=10.5,
         color=RGBColor(0xC9, 0xD1, 0xD9))
add_footer(s, 16)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 17 · Evaluation / POC Workflow
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
add_header(s, "POC", "Customer Evaluation & POC Workflow",
           "A four-stage path from first look to production deployment.")

stages = [
    ("Week 1", "Guided walkthrough",
     "Live session on the Emergent-hosted instance. Analysts paste real "
     "artifacts and see the evidence chain end-to-end."),
    ("Week 2", "Customer-hosted PoC",
     "Docker Compose install in your DMZ or evaluation VLAN. Bring 3-5 "
     "historical investigations; we reproduce them under evidence-provenance."),
    ("Week 3", "Integration & workflow validation",
     "Connect TI feeds, LLM (or none), TAXII target. Analyst team runs "
     "10-20 live investigations. Corrections and feedback captured."),
    ("Week 4", "Assessment & go/no-go",
     "Joint review: capability coverage · security posture · deployment fit · "
     "roadmap priorities. Signed evidence report per investigation."),
]
top = Inches(2.15)
for i, (when, name, desc) in enumerate(stages):
    left = Inches(0.75 + i * 3.15)
    rect(s, left, top, Inches(3.0), Inches(4.4))
    # when badge
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.3),
                               top + Inches(0.3), Inches(1.4), Inches(0.4))
    badge.adjustments[0] = 0.4
    badge.fill.solid(); badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    tf = badge.text_frame; tf.text = when
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]; r.font.name = FONT_HEAD; r.font.size = Pt(11)
    r.font.bold = True; r.font.color.rgb = BG
    text(s, name, left + Inches(0.25), top + Inches(0.9),
         Inches(2.6), Inches(0.6), size=15, bold=True, color=FG, font=FONT_HEAD)
    text(s, desc, left + Inches(0.25), top + Inches(1.6),
         Inches(2.6), Inches(2.6), size=11.5,
         color=RGBColor(0xC9, 0xD1, 0xD9))

text(s, "Every stage produces reproducible artefacts you keep: decoded chains, "
        "verdict cards, MITRE mappings, signed reports, and evidence graphs.",
     Inches(0.75), Inches(6.95), Inches(12), Inches(0.4), size=11, color=MUTED)
add_footer(s, 17)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 18 · Closing / Contact
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
# accent bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SLIDE_H)
bar.line.fill.background(); bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT

text(s, "Evidence over inference.",
     Inches(0.85), Inches(2.3), Inches(11), Inches(1.0),
     size=54, bold=True, color=FG, font=FONT_HEAD)
text(s, "Explainability over opacity.",
     Inches(0.85), Inches(3.2), Inches(11), Inches(1.0),
     size=54, bold=True, color=ACCENT, font=FONT_HEAD)
text(s, "Your environment. Your data. Your control.",
     Inches(0.85), Inches(4.4), Inches(11), Inches(0.6),
     size=22, color=MUTED, font=FONT_HEAD)

text(s, "Next steps",
     Inches(0.85), Inches(5.4), Inches(6), Inches(0.4),
     size=13, bold=True, color=ACCENT, font=FONT_HEAD)
bullets(s, [
    "Schedule guided walkthrough on Emergent-hosted evaluation instance",
    "Confirm PoC deployment target (VM / Kubernetes / evaluation VLAN)",
    "Identify 3-5 historical investigations for replay validation",
    "Nominate technical evaluator and security architect for weekly cadence",
], Inches(0.85), Inches(5.8), Inches(11), Inches(1.5), size=13)

text(s, "NivXRay · Customer Technical Deck · Draft · 2026-08-11",
     Inches(0.75), Inches(7.15), Inches(12), Inches(0.3),
     size=9, color=MUTED, align=PP_ALIGN.RIGHT)

# ── Save ──
OUT = Path("/app/frontend/public/downloads/NivXRay-Customer-Technical-Deck.pptx")
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"WROTE {OUT}  ({OUT.stat().st_size:,} bytes · {len(prs.slides)} slides)")
