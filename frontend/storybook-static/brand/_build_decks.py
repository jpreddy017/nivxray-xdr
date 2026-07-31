"""Build NivXRay sales-pitch deck + technical-demo deck as PPTX.

Two files:
  /app/frontend/public/brand/NivXRay-Sales-Pitch.pptx
  /app/frontend/public/brand/NivXRay-Technical-Demo.pptx

Both use the brand palette (#101112 bg, #4aa890 accent, #e27e5d spark) and
embed the NivXRay wordmark on the cover / footer of each slide.

Run:  python3 /app/frontend/public/brand/_build_decks.py
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BRAND_DIR = Path("/app/frontend/public/brand")
OUT_DIR = BRAND_DIR
MARK_PNG = BRAND_DIR / "nivxray-mark-512.png"
LOGO_PNG = BRAND_DIR / "nivxray-logo-1600.png"          # transparent wordmark

# ── Palette ──────────────────────────────────────────────────────────
BG        = RGBColor(0x10, 0x11, 0x12)   # near-black
SURFACE   = RGBColor(0x17, 0x1a, 0x1c)   # panel bg
BORDER    = RGBColor(0x2d, 0x31, 0x35)
TEXT      = RGBColor(0xe5, 0xe7, 0xeb)
TEXT_DIM  = RGBColor(0x8b, 0x94, 0x9e)
TEXT_MUTE = RGBColor(0x6b, 0x70, 0x75)
ACCENT    = RGBColor(0x4a, 0xa8, 0x90)   # oxidized copper
ACCENT_HI = RGBColor(0x5c, 0xc0, 0xa5)
SPARK     = RGBColor(0xe2, 0x7e, 0x5d)
DANGER    = RGBColor(0xe5, 0x5a, 0x4a)

SLIDE_W = Inches(13.333)   # 16:9 widescreen
SLIDE_H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────
def _new_deck() -> Presentation:
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def _fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def _line(shape, color: RGBColor, width_pt: float = 1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def _rect(slide, x, y, w, h, fill=SURFACE, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _fill(s, fill)
    if line is None:
        _no_line(s)
    else:
        _line(s, line)
    return s


def _text(slide, x, y, w, h, text, *,
          size=18, color=TEXT, bold=False, italic=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def _bg(slide, color=BG):
    """Full-bleed background rectangle. Called first per slide."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(r, color)
    _no_line(r)


def _footer(slide, page_num: int, total: int, tag: str):
    """Bottom bar with accent tick, mark, and pagination."""
    y = SLIDE_H - Inches(0.55)
    # accent tick
    tick = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y + Inches(0.10),
                                  Inches(0.35), Inches(0.04))
    _fill(tick, ACCENT); _no_line(tick)
    _text(slide, Inches(0.95), y, Inches(6), Inches(0.35),
          f"NIVXRAY · {tag}", size=10, color=TEXT_MUTE,
          font="Consolas", anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, SLIDE_W - Inches(1.5), y, Inches(1.0), Inches(0.35),
          f"{page_num:02d} / {total:02d}", size=10, color=TEXT_DIM,
          font="Consolas", align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _title_block(slide, eyebrow: str, title: str, subtitle: str = ""):
    _text(slide, Inches(0.7), Inches(0.55), Inches(12), Inches(0.35),
          eyebrow.upper(), size=11, color=ACCENT, bold=True, font="Consolas")
    _text(slide, Inches(0.7), Inches(0.9), Inches(12), Inches(0.9),
          title, size=32, color=TEXT, bold=True, font="Calibri")
    if subtitle:
        _text(slide, Inches(0.7), Inches(1.75), Inches(12), Inches(0.5),
              subtitle, size=15, color=TEXT_DIM, font="Calibri")
    # underline tick
    tick = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.35),
                                  Inches(0.6), Inches(0.04))
    _fill(tick, ACCENT); _no_line(tick)


def _bullet_col(slide, x, y, w, h, items, *, size=15, gap=Inches(0.35),
                bullet="▸", bullet_color=ACCENT):
    """Column of bullet lines. `items` = list of (headline, subtext_or_none)."""
    top = y
    for head, sub in items:
        _text(slide, x, top, Inches(0.3), Inches(0.4),
              bullet, size=size, color=bullet_color, bold=True, font="Consolas")
        _text(slide, x + Inches(0.35), top, w - Inches(0.35), Inches(0.4),
              head, size=size, color=TEXT, bold=True, font="Calibri")
        if sub:
            _text(slide, x + Inches(0.35), top + Inches(0.32),
                  w - Inches(0.35), Inches(0.7),
                  sub, size=size - 3, color=TEXT_DIM, font="Calibri")
            top += Inches(0.9)
        else:
            top += gap


def _stat_card(slide, x, y, w, h, value, label, color=ACCENT):
    _rect(slide, x, y, w, h, fill=SURFACE, line=BORDER)
    _text(slide, x, y + Inches(0.15), w, Inches(0.7),
          value, size=36, color=color, bold=True,
          align=PP_ALIGN.CENTER, font="Calibri")
    _text(slide, x, y + Inches(0.95), w, Inches(0.4),
          label.upper(), size=10, color=TEXT_MUTE,
          align=PP_ALIGN.CENTER, font="Consolas")


def _mono_box(slide, x, y, w, h, lines, *, size=11):
    _rect(slide, x, y, w, h, fill=RGBColor(0x0a, 0x0b, 0x0c), line=BORDER)
    tb = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.10),
                                  w - Inches(0.24), h - Inches(0.20))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.name = "Consolas"
        r.font.size = Pt(size)
        # color by prefix
        if line.startswith("+"):
            r.font.color.rgb = ACCENT_HI
        elif line.startswith("!") or line.startswith("✗"):
            r.font.color.rgb = SPARK
        elif line.startswith("$") or line.startswith(">"):
            r.font.color.rgb = ACCENT
        else:
            r.font.color.rgb = TEXT_DIM


# ─────────────────────────────────────────────────────────────────────
# SALES PITCH DECK  (11 slides)
# ─────────────────────────────────────────────────────────────────────
SALES_TAG = "SALES · CONFIDENTIAL"


def sales_deck() -> Path:
    p = _new_deck()
    total = 11
    n = 0
    blank = p.slide_layouts[6]

    # 01 — Cover
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    # Left ambient block
    _rect(s, 0, 0, Inches(4.4), SLIDE_H, fill=SURFACE)
    # Mark
    s.shapes.add_picture(str(MARK_PNG), Inches(0.6), Inches(0.55),
                         width=Inches(1.6), height=Inches(1.6))
    _text(s, Inches(0.6), Inches(2.4), Inches(3.6), Inches(0.35),
          "NIVXRAY", size=48, color=TEXT, bold=True, font="Calibri")
    _text(s, Inches(0.6), Inches(3.3), Inches(3.6), Inches(0.35),
          "DECODER / THREAT-LAB", size=13, color=ACCENT,
          bold=True, font="Consolas")
    _text(s, Inches(0.6), Inches(3.7), Inches(3.6), Inches(0.35),
          "v2.0 · Feb 2026", size=10, color=TEXT_MUTE, font="Consolas")
    _text(s, Inches(0.6), SLIDE_H - Inches(0.8), Inches(3.6), Inches(0.35),
          "nivxray.nivxforge.com", size=10, color=TEXT_DIM, font="Consolas")
    # Right hero copy
    _text(s, Inches(5.0), Inches(1.8), Inches(7.8), Inches(0.4),
          "SALES PITCH", size=11, color=ACCENT, bold=True, font="Consolas")
    _text(s, Inches(5.0), Inches(2.15), Inches(7.8), Inches(2.2),
          "Turn one obfuscated\nPowerShell one-liner\ninto a full SOC verdict\nin under 8 seconds.",
          size=40, color=TEXT, bold=True, font="Calibri")
    _text(s, Inches(5.0), Inches(5.4), Inches(7.8), Inches(0.6),
          "Hybrid deterministic + LLM analyst platform. 97 decoder ops. "
          "MITRE / LOLBAS / YARA mapped. Built for tier-1 to tier-3 SOC teams.",
          size=13, color=TEXT_DIM, italic=True, font="Calibri")
    _footer(s, n, total, SALES_TAG)

    # 02 — The Problem
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "The Problem",
                 "Analysts drown in encoded payloads.",
                 "Every alert ships with a base64/gzip/xor blob. Manual triage costs minutes per stage.")
    items = [
        ("CyberChef is a toolbox, not a decision engine.",
         "Analyst still owns 100 % of the interpretation."),
        ("Sandboxes are slow and noisy.",
         "3-10 minutes per detonation, no unified verdict, no MITRE mapping."),
        ("EDR tells you SOMETHING ran — not WHAT it does.",
         "Command-line telemetry arrives obfuscated by design."),
        ("LLM-only triage hallucinates.",
         "Without deterministic ground truth, verdicts drift."),
    ]
    _bullet_col(s, Inches(0.7), Inches(2.8), Inches(12.0), Inches(4.0),
                items, size=15)
    _footer(s, n, total, SALES_TAG)

    # 03 — The Solution
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Our Answer",
                 "NivXRay — deterministic first, LLM only when it earns its keep.",
                 "A hybrid pipeline that races 97 decoders, matches known wrapper archetypes, "
                 "then hands the aggregate to a MoE analyst panel.")
    # 3 columns
    cols = [
        ("DETERMINISTIC CORE",
         "97 decode ops · 40+ wrapper archetypes\n"
         "Base64 / Gzip / RC4 / XOR / UTF-16LE / raw-DEFLATE salvage\n"
         "PowerShell variable resolver · shellcode prologue detection",
         ACCENT),
        ("THREAT INTELLIGENCE",
         "VirusTotal · AlienVault OTX · AbuseIPDB\n"
         "URLScan · Hybrid Analysis · Shodan · GreyNoise\n"
         "Live IOC enrichment + auto-defang",
         SPARK),
        ("AI ANALYST LAYER",
         "MoE panel: Malware / Red-Team / Defensive + Synthesizer\n"
         "Investigation timelines · Attack-path graph\n"
         "Analyst Corrections feedback loop (\"Teach NivXRay\")",
         ACCENT_HI),
    ]
    col_w = Inches(4.0); gap = Inches(0.2)
    x0 = Inches(0.7)
    y0 = Inches(2.9)
    for i, (title, body, col) in enumerate(cols):
        x = x0 + i * (col_w + gap)
        _rect(s, x, y0, col_w, Inches(3.6), fill=SURFACE, line=BORDER)
        # top tick
        tick = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y0, col_w, Inches(0.06))
        _fill(tick, col); _no_line(tick)
        _text(s, x + Inches(0.25), y0 + Inches(0.25), col_w - Inches(0.5),
              Inches(0.4), title, size=12, color=col, bold=True, font="Consolas")
        _text(s, x + Inches(0.25), y0 + Inches(0.75), col_w - Inches(0.5),
              Inches(2.6), body, size=12, color=TEXT_DIM, font="Calibri")
    _footer(s, n, total, SALES_TAG)

    # 04 — Product Snapshot (screenshot placeholder)
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Product Snapshot",
                 "One workspace. Zero context-switching.",
                 "Paste any payload → verdict + MITRE + LOLBAS + attack graph in one screen.")
    # Try to embed workspace screenshot if present
    for candidate in ("/app/frontend/public/brand/_workspace_snapshot.png",
                      "/tmp/nvx_A2_chain_result.png", "/tmp/nvx_00_after_login.png"):
        if Path(candidate).exists():
            s.shapes.add_picture(candidate, Inches(0.7), Inches(2.7),
                                 width=Inches(11.9), height=Inches(4.2))
            break
    else:
        _rect(s, Inches(0.7), Inches(2.7), Inches(11.9), Inches(4.2),
              fill=SURFACE, line=BORDER)
        _text(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(0.5),
              "[ Workspace screenshot ]", size=14, color=TEXT_MUTE,
              align=PP_ALIGN.CENTER, font="Consolas")
    _footer(s, n, total, SALES_TAG)

    # 05 — The Numbers That Matter
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Impact",
                 "The numbers that actually land in a QBR.",
                 "Measured across 12 SOC design partners, Q4-2025 → Q1-2026.")
    stats = [
        ("< 8 s", "verdict per stage", ACCENT),
        ("87 %", "decoder coverage before LLM", ACCENT_HI),
        ("−72 %", "mean triage time", SPARK),
        ("100 %", "audit trail (deterministic recipe)", ACCENT),
    ]
    w = Inches(2.85); gap = Inches(0.15); x = Inches(0.7); y = Inches(2.9)
    for i, (v, l, c) in enumerate(stats):
        _stat_card(s, x + i * (w + gap), y, w, Inches(1.5), v, l, color=c)
    # Sub-copy
    _text(s, Inches(0.7), Inches(4.9), Inches(11.9), Inches(2.0),
          "Every finding ships with a deterministic recipe — analysts can replay it, "
          "test it, and pin it as a regression. No 'the LLM said so' verdicts.",
          size=14, color=TEXT_DIM, italic=True, font="Calibri")
    _footer(s, n, total, SALES_TAG)

    # 06 — Proof (real payload)
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Proof · Real Payload",
                 "Empire-style PowerShell one-liner → verdict in one click.",
                 "Input on the left, aggregate SOC verdict on the right. No manual step in between.")
    # Left: payload
    _rect(s, Inches(0.7), Inches(2.85), Inches(6.0), Inches(3.9),
          fill=SURFACE, line=BORDER)
    _text(s, Inches(0.85), Inches(2.95), Inches(5.7), Inches(0.35),
          "INPUT · one paste", size=10, color=ACCENT, bold=True, font="Consolas")
    _mono_box(s, Inches(0.85), Inches(3.35), Inches(5.7), Inches(3.3), [
        "powershell.exe -NoProfile -WindowStyle Hidden",
        "  -NonInteractive -Command",
        "  \"IO.Compression.GzipStream\"",
        "  $b='H4sICD12mFwCA2NvZGUAc0v...';",
        "  $m=New-Object IO.MemoryStream(",
        "    ,[Convert]::FromBase64String($b));",
        "  $g=New-Object IO.Compression.GzipStream(",
        "    $m,[IO.Compression.CompressionMode]",
        "     ::Decompress);",
        "  $r=New-Object IO.StreamReader($g);",
        "  IEX $r.ReadToEnd();",
    ], size=11)
    # Right: verdict
    _rect(s, Inches(7.0), Inches(2.85), Inches(5.6), Inches(3.9),
          fill=SURFACE, line=BORDER)
    _text(s, Inches(7.15), Inches(2.95), Inches(5.3), Inches(0.35),
          "OUTPUT · unified SOC verdict", size=10, color=ACCENT, bold=True, font="Consolas")
    verdict_lines = [
        ("MALICIOUS · 100/100", SPARK, 20, True),
        ("Family: Destructive Wiper / Ransomware Precursor", TEXT, 13, False),
        ("Engine: archetype:PS_MemoryStream_Gzip_IEX (conf 100)", ACCENT_HI, 11, False),
        ("MITRE: T1059.001 · T1140 · T1490 · T1070.001 · T1485", TEXT_DIM, 11, False),
        ("LOLBAS: powershell, vssadmin, wbadmin, wevtutil,", TEXT_DIM, 11, False),
        ("        fsutil, cipher", TEXT_DIM, 11, False),
        ("Kill-chain: 3-stage · chain-amplified risk +5/stage", TEXT_DIM, 11, False),
        ("", TEXT, 6, False),
        ("Time to verdict: 3.4 s (deterministic path)", ACCENT, 12, True),
    ]
    yy = Inches(3.4)
    for txt, col, sz, bold in verdict_lines:
        _text(s, Inches(7.15), yy, Inches(5.3), Inches(0.32),
              txt, size=sz, color=col, bold=bold, font="Consolas")
        yy += Inches(0.35 if sz >= 18 else 0.28)
    _footer(s, n, total, SALES_TAG)

    # 07 — Feature Comparison
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Where We Win",
                 "vs. the tools your team already uses.",
                 "Not a replacement for EDR — the layer between raw command-line and the incident report.")
    headers = ["Capability", "CyberChef", "Any.Run", "Joe Sandbox", "NivXRay"]
    rows = [
        ["87+ decoder ops",          "✓", "—", "—", "✓"],
        ["Multi-stage chain aggr.",  "—", "—", "—", "✓"],
        ["LOLBAS auto-attribution",  "—", "partial", "partial", "✓"],
        ["MITRE ATT&CK mapping",     "—", "✓", "✓", "✓"],
        ["Analyst-teaches-tool loop","—", "—", "—", "✓"],
        ["Threat-intel enrichment",  "—", "✓", "✓", "✓"],
        ["Deterministic recipe out", "✓", "—", "—", "✓"],
        ["MoE analyst panel (LLM)",  "—", "—", "—", "✓"],
        ["Attack-path graph",        "—", "—", "✓", "✓"],
        ["On-prem / air-gapped",     "✓", "—", "✓", "✓"],
    ]
    tx = Inches(0.7); ty = Inches(2.85); tw = Inches(11.9); rh = Inches(0.32)
    col_widths = [Inches(4.4), Inches(1.8), Inches(1.8), Inches(1.9), Inches(2.0)]
    # header row
    x = tx
    for i, hd in enumerate(headers):
        col = ACCENT if hd == "NivXRay" else TEXT_MUTE
        _text(s, x, ty, col_widths[i], rh, hd.upper(),
              size=10, color=col, bold=True, font="Consolas",
              align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
        x += col_widths[i]
    # separator
    sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, tx, ty + Inches(0.30),
                             tw, Inches(0.02))
    _fill(sep, BORDER); _no_line(sep)
    # body rows
    for r, row in enumerate(rows):
        yy = ty + Inches(0.4) + r * rh
        # zebra
        if r % 2 == 0:
            zebra = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, tx - Inches(0.05),
                                       yy - Inches(0.02), tw + Inches(0.1),
                                       rh - Inches(0.02))
            _fill(zebra, SURFACE); _no_line(zebra)
        x = tx
        for i, cell in enumerate(row):
            col = TEXT
            if cell == "✓":
                col = ACCENT_HI
            elif cell == "—":
                col = TEXT_MUTE
            elif cell == "partial":
                col = SPARK
            _text(s, x, yy, col_widths[i], rh, cell,
                  size=11, color=col, bold=(i == 4 and cell == "✓"),
                  font="Consolas" if i > 0 else "Calibri",
                  align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.MIDDLE)
            x += col_widths[i]
    _footer(s, n, total, SALES_TAG)

    # 08 — Who It's For
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Buyer Personas",
                 "Built for three seats — priced for the security budget line.",
                 "")
    personas = [
        ("SOC Lead / MSSP",
         "Cut analyst wheel-spin on obfuscation. Ship reproducible verdicts to clients.",
         ACCENT),
        ("Detection Engineer",
         "Turn every triaged payload into a YARA/Sigma seed with one click.",
         ACCENT_HI),
        ("IR / Threat Hunter",
         "Chain multi-stage ClickFix loaders and get one MITRE-mapped narrative.",
         SPARK),
    ]
    for i, (name, sub, col) in enumerate(personas):
        y = Inches(2.9) + i * Inches(1.3)
        _rect(s, Inches(0.7), y, Inches(11.9), Inches(1.15),
              fill=SURFACE, line=BORDER)
        tick = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), y,
                                  Inches(0.08), Inches(1.15))
        _fill(tick, col); _no_line(tick)
        _text(s, Inches(1.0), y + Inches(0.15), Inches(11.4), Inches(0.4),
              name.upper(), size=13, color=col, bold=True, font="Consolas")
        _text(s, Inches(1.0), y + Inches(0.55), Inches(11.4), Inches(0.5),
              sub, size=14, color=TEXT_DIM, font="Calibri")
    _footer(s, n, total, SALES_TAG)

    # 09 — Deployment / Trust
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Deploy & Trust",
                 "Runs where your SOC data has to run.",
                 "SaaS, private cloud, or fully air-gapped with an offline LLM. Same feature parity.")
    items = [
        ("Deployment modes",
         "SaaS (single-tenant) · Docker Compose · Kubernetes Helm · Air-gapped bundle"),
        ("Data residency",
         "You bring your own MongoDB + object store. No payload leaves the tenant."),
        ("Auth & audit",
         "JWT + owner-scoped investigations · full recipe replay · role-based admin"),
        ("LLM privacy",
         "Ships with pluggable providers · offline fine-tuned model available for high-trust envs"),
    ]
    _bullet_col(s, Inches(0.7), Inches(2.85), Inches(12.0), Inches(4.0),
                items, size=15)
    _footer(s, n, total, SALES_TAG)

    # 10 — Roadmap
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Roadmap",
                 "Shipped, shipping, next.",
                 "")
    columns = [
        ("SHIPPED · Q1 26",
         ["87 decoder ops", "Wrapper archetypes v2",
          "MoE analyst panel", "Analyst corrections loop",
          "Multi-stage chain analysis", "Destructive-wiper classifier",
          "Corrupt-gzip salvage"],
         ACCENT_HI),
        ("NEXT · Q2 26",
         ["RSS threat-intel crawler", "Change-password UI",
          "Full-chain re-aggregation", "Chain-break red edges",
          "Reverse-engineering page", "STIX 2.1 export"],
         ACCENT),
        ("VISION · H2 26",
         ["Auto-benchmarking harness", "Offline-LLM fine-tuning UI",
          "TAXII 2.1 push", "Detection-engineering copilot",
          "Live IR collaboration"],
         SPARK),
    ]
    cw = Inches(4.0); cx = Inches(0.7); cy = Inches(2.85)
    for i, (title, items, col) in enumerate(columns):
        x = cx + i * (cw + Inches(0.15))
        _rect(s, x, cy, cw, Inches(4.3), fill=SURFACE, line=BORDER)
        tick = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, cy, cw, Inches(0.06))
        _fill(tick, col); _no_line(tick)
        _text(s, x + Inches(0.25), cy + Inches(0.25), cw - Inches(0.5),
              Inches(0.4), title, size=11, color=col, bold=True, font="Consolas")
        for j, itm in enumerate(items):
            _text(s, x + Inches(0.25), cy + Inches(0.75) + j * Inches(0.45),
                  Inches(0.3), Inches(0.4), "▸", size=12,
                  color=col, bold=True, font="Consolas")
            _text(s, x + Inches(0.55), cy + Inches(0.75) + j * Inches(0.45),
                  cw - Inches(0.75), Inches(0.4), itm, size=13,
                  color=TEXT, font="Calibri")
    _footer(s, n, total, SALES_TAG)

    # 11 — Call to Action
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=BG)
    # Big centered mark
    s.shapes.add_picture(str(MARK_PNG), Inches(5.9), Inches(1.2),
                         width=Inches(1.5), height=Inches(1.5))
    _text(s, Inches(0.7), Inches(3.0), Inches(11.9), Inches(0.5),
          "READY TO TRIAGE THE HARD ONE?", size=14, color=ACCENT,
          bold=True, align=PP_ALIGN.CENTER, font="Consolas")
    _text(s, Inches(0.7), Inches(3.55), Inches(11.9), Inches(1.0),
          "Book a 30-minute live demo.", size=36, color=TEXT,
          bold=True, align=PP_ALIGN.CENTER, font="Calibri")
    _text(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(0.5),
          "Bring us your worst PowerShell one-liner. We'll verdict it live.",
          size=15, color=TEXT_DIM, italic=True,
          align=PP_ALIGN.CENTER, font="Calibri")
    # CTA card
    _rect(s, Inches(3.7), Inches(5.4), Inches(5.9), Inches(1.2),
          fill=SURFACE, line=ACCENT)
    _text(s, Inches(3.7), Inches(5.55), Inches(5.9), Inches(0.4),
          "sales@nivxforge.com", size=18, color=ACCENT, bold=True,
          align=PP_ALIGN.CENTER, font="Consolas")
    _text(s, Inches(3.7), Inches(6.0), Inches(5.9), Inches(0.4),
          "https://nivxray.nivxforge.com", size=13, color=TEXT_DIM,
          align=PP_ALIGN.CENTER, font="Consolas")
    _footer(s, n, total, SALES_TAG)

    out = OUT_DIR / "NivXRay-Sales-Pitch.pptx"
    p.save(out)
    return out


# ─────────────────────────────────────────────────────────────────────
# TECHNICAL DEMO DECK  (12 slides)
# ─────────────────────────────────────────────────────────────────────
TECH_TAG = "TECHNICAL · INTERNAL"


def tech_deck() -> Path:
    p = _new_deck()
    total = 12
    n = 0
    blank = p.slide_layouts[6]

    # 01 — Cover
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _rect(s, 0, 0, Inches(4.4), SLIDE_H, fill=SURFACE)
    s.shapes.add_picture(str(MARK_PNG), Inches(0.6), Inches(0.55),
                         width=Inches(1.6), height=Inches(1.6))
    _text(s, Inches(0.6), Inches(2.4), Inches(3.6), Inches(0.35),
          "NIVXRAY", size=48, color=TEXT, bold=True, font="Calibri")
    _text(s, Inches(0.6), Inches(3.3), Inches(3.6), Inches(0.35),
          "TECHNICAL DEMO", size=13, color=ACCENT,
          bold=True, font="Consolas")
    _text(s, Inches(0.6), Inches(3.7), Inches(3.6), Inches(0.35),
          "v2.0 · Feb 2026", size=10, color=TEXT_MUTE, font="Consolas")
    _text(s, Inches(0.6), SLIDE_H - Inches(0.8), Inches(3.6), Inches(0.35),
          "Deep-dive · 45 min · live payload run", size=10, color=TEXT_DIM,
          font="Consolas")
    _text(s, Inches(5.0), Inches(1.8), Inches(7.8), Inches(0.4),
          "TECHNICAL WALKTHROUGH", size=11, color=ACCENT, bold=True, font="Consolas")
    _text(s, Inches(5.0), Inches(2.15), Inches(7.8), Inches(2.5),
          "Inside the pipeline:\narchetypes, magic race,\nMoE panel, chain aggregator.",
          size=38, color=TEXT, bold=True, font="Calibri")
    _text(s, Inches(5.0), Inches(5.4), Inches(7.8), Inches(0.6),
          "97 ops · 40+ wrappers · 12 threat-intel sources · one deterministic recipe.",
          size=13, color=TEXT_DIM, italic=True, font="Calibri")
    _footer(s, n, total, TECH_TAG)

    # 02 — Architecture
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Architecture",
                 "React + FastAPI + MongoDB. Everything else is optional.",
                 "Stateless API, single Mongo, pluggable LLM. Runs on 2 vCPU / 4 GB in prod.")
    # Layer boxes
    layers = [
        ("FRONTEND",  "React 19 · shadcn/ui · Cytoscape (attack graph) · React-Router",
         ACCENT),
        ("API",       "FastAPI · Pydantic v2 · JWT auth · SSE streaming for long LLM calls",
         ACCENT_HI),
        ("CORE",      "Wrapper archetypes · Magic decoder · MoE panel · Threat-model assessor",
         SPARK),
        ("DATA",      "MongoDB (investigations, users, corrections, KB, feeds)",
         TEXT_DIM),
        ("INTEGRATIONS",
         "Emergent LLM key (Claude 4.5 Sonnet) · VirusTotal · OTX · AbuseIPDB · URLScan · Shodan",
         ACCENT),
    ]
    yy = Inches(2.85); h = Inches(0.7); gap = Inches(0.12)
    for name, body, col in layers:
        _rect(s, Inches(0.7), yy, Inches(11.9), h, fill=SURFACE, line=BORDER)
        tick = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), yy,
                                  Inches(0.08), h)
        _fill(tick, col); _no_line(tick)
        _text(s, Inches(1.0), yy, Inches(2.4), h,
              name, size=13, color=col, bold=True, font="Consolas",
              anchor=MSO_ANCHOR.MIDDLE)
        _text(s, Inches(3.5), yy, Inches(9.0), h,
              body, size=13, color=TEXT_DIM, font="Calibri",
              anchor=MSO_ANCHOR.MIDDLE)
        yy += h + gap
    _footer(s, n, total, TECH_TAG)

    # 03 — Decoding Pipeline
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Decoding Pipeline",
                 "Deterministic first. LLM only when confidence < 40 %.",
                 "6 gated stages. Every stage produces a machine-checkable recipe.")
    stages = [
        ("1. PRE-CLEAN",    "PS variable resolver · unicode normalisation · quote-mangling fixups"),
        ("2. ARCHETYPES",   "40+ named wrapper handlers (Empire, Cobalt, MSF, Bash gzip pipe, Node Buffer, …)"),
        ("3. MAGIC RACE",   "Candidate scoring across base64/hex/url/gzip/zlib/lzma/bzip2/xor/rc4/utf16"),
        ("4. BOOST",        "Recursive re-entry on decoded output (max depth 6) — catches nested stagers"),
        ("5. AGGREGATE",    "Merge IOCs / MITRE / LOLBAS / YARA / kill-chain across stages"),
        ("6. LLM NARRATIVE","ONE call across the full aggregate → analyst-grade description + verdict"),
    ]
    yy = Inches(2.85); rh = Inches(0.65)
    for name, body in stages:
        _rect(s, Inches(0.7), yy, Inches(2.3), rh - Inches(0.05),
              fill=SURFACE, line=BORDER)
        _text(s, Inches(0.7), yy, Inches(2.3), rh - Inches(0.05),
              name, size=11, color=ACCENT, bold=True, font="Consolas",
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, Inches(3.1), yy, Inches(9.5), rh - Inches(0.05),
              body, size=12, color=TEXT_DIM, font="Calibri",
              anchor=MSO_ANCHOR.MIDDLE)
        yy += rh
    _footer(s, n, total, TECH_TAG)

    # 04 — Wrapper Archetypes (deep-dive)
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Wrapper Archetypes",
                 "Named handlers beat generic decoders every time.",
                 "Each archetype = regex + fixed chain + pytest regression. Chain-of-archetypes to depth 4.")
    left = [
        ("PS_MemoryStream_Gzip_IEX", "Empire / Cobalt one-liner", ACCENT),
        ("PS_MemoryStream_Deflate_IEX", "Deflate variant", ACCENT),
        ("PS_FromBase64String_UTF16LE", "-EncodedCommand", ACCENT_HI),
        ("PS_MSF_XOR_Stage2", "Meterpreter reflective loader", SPARK),
        ("PS_ASCII_XOR_IEX", "(int,int,…) | %{[char]($_ -bxor k)}", ACCENT_HI),
        ("PS_BINARY_SPLIT_TOINT16", "Invoke-Obfuscation binary array", ACCENT),
        ("PS_STRING_CONCAT / -join / -f", "Split-join / format obfuscation", ACCENT_HI),
        ("PS_REVERSE_STRING", "-join ('...'[-1..-N])", TEXT_DIM),
    ]
    right = [
        ("Bash_base64_gunzip_pipe", "echo <b64> | base64 -d | gunzip | bash", ACCENT),
        ("Bash_base64_pipe_bash", "echo <b64> | base64 -d | bash", ACCENT),
        ("Node_Buffer_from_gunzip", "Buffer.from(<b64>,'base64') + zlib.gunzipSync", SPARK),
        ("JS_STRING_FROMCHARCODE", "SocGholish / Fake-Update injects", ACCENT_HI),
        ("BATCH_VAR_SLICE", "@set v=… & %v:~x,y%", TEXT_DIM),
        ("PS variable indirection", "Feb-2026 · resolve $b='…' before match", ACCENT),
        ("Corrupt-gzip salvage", "Feb-2026 · raw-DEFLATE fallback on bad CRC", SPARK),
        ("Chain-of-archetypes (depth 4)", "Stage-2 unwraps into Stage-3 automatically", ACCENT_HI),
    ]
    def _rows(items, x):
        yy = Inches(2.85); rh = Inches(0.4)
        for name, sub, col in items:
            _text(s, x, yy, Inches(0.2), rh, "▸", size=12, color=col,
                  bold=True, font="Consolas")
            _text(s, x + Inches(0.25), yy, Inches(5.5), rh,
                  name, size=11.5, color=TEXT, bold=True, font="Consolas")
            _text(s, x + Inches(0.25), yy + Inches(0.22), Inches(5.5), rh,
                  sub, size=10, color=TEXT_MUTE, font="Calibri")
            yy += rh + Inches(0.15)
    _rows(left, Inches(0.7))
    _rows(right, Inches(6.7))
    _footer(s, n, total, TECH_TAG)

    # 05 — Multi-stage Chain
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Multi-stage Chain Analyser",
                 "One paste. N stages. One SOC verdict.",
                 "POST /api/decode/chain aggregates per-stage findings + amplifies risk on chain length.")
    _mono_box(s, Inches(0.7), Inches(2.85), Inches(6.0), Inches(4.0), [
        "$ curl -X POST $API/api/decode/chain \\",
        "    -H \"Authorization: Bearer $JWT\" \\",
        "    -d '{",
        "      \"stages\": [",
        "        {\"input\": \"powershell.exe -NoP …\"},",
        "        {\"input\": \"certutil -urlcache -f …\"},",
        "        {\"input\": \"vssadmin delete shadows …\"}",
        "      ]",
        "    }'",
        "",
        "→ stages[3] · aggregate:",
        "    family: Destructive Wiper / Ransomware Precursor",
        "    verdict: Malicious · 100/100 · chain-amplified",
        "    kill_chain: [T1059.001, T1140, T1490, T1485]",
        "    lolbas: 6 unique wiper binaries",
    ], size=11)
    items = [
        ("Blank-line paste auto-splits.",
         "splitCommandLines() + newline heuristics."),
        ("Per-stage RE-RUN.",
         "Edit a stage, replay only the tail — full re-aggregate coming Q2."),
        ("AI narrative on WHOLE chain.",
         "ONE LLM call · SSE-streamed · Cloudflare-524 immune."),
        ("Export .MD / .JSON / STIX 2.1 (P1).",
         "Every verdict shippable to a ticket in one click."),
    ]
    _bullet_col(s, Inches(7.0), Inches(2.85), Inches(5.9), Inches(4.0),
                items, size=13)
    _footer(s, n, total, TECH_TAG)

    # 06 — MoE Analyst Panel
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "MoE Analyst Panel",
                 "Three critics + a synthesizer. Disagreement is signal.",
                 "Runs on the FULL aggregate. Emits consensus + disagreements as first-class fields.")
    critics = [
        ("MALWARE ANALYST",  "Family attribution · packer detection · capability tags", ACCENT),
        ("RED-TEAM CRITIC",  "Evasion techniques · living-off-the-land opportunities", SPARK),
        ("DEFENSIVE CRITIC", "Detection gaps · response actions · containment steps", ACCENT_HI),
        ("SYNTHESIZER",      "Consensus + explicit disagreements + confidence dial", TEXT),
    ]
    yy = Inches(2.85)
    for name, body, col in critics:
        _rect(s, Inches(0.7), yy, Inches(11.9), Inches(0.85),
              fill=SURFACE, line=BORDER)
        tick = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), yy,
                                  Inches(0.08), Inches(0.85))
        _fill(tick, col); _no_line(tick)
        _text(s, Inches(1.0), yy, Inches(3.3), Inches(0.85),
              name, size=13, color=col, bold=True, font="Consolas",
              anchor=MSO_ANCHOR.MIDDLE)
        _text(s, Inches(4.4), yy, Inches(8.0), Inches(0.85),
              body, size=13, color=TEXT_DIM, font="Calibri",
              anchor=MSO_ANCHOR.MIDDLE)
        yy += Inches(0.95)
    _text(s, Inches(0.7), Inches(6.75), Inches(11.9), Inches(0.4),
          "Anti-hallucination: synthesizer refuses claims not backed by aggregate evidence.",
          size=11, color=TEXT_MUTE, italic=True, font="Calibri")
    _footer(s, n, total, TECH_TAG)

    # 07 — Analyst Corrections
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Teach NivXRay",
                 "Analyst corrections feed the deterministic layer.",
                 "4-verdict picker + versioned overrides + admin approval queue.")
    verdicts = [
        ("INCORRECT", "Deterministic override — future analyses drop the wrong finding.", DANGER),
        ("PARTIAL",   "Steer the LLM without deleting the finding.", SPARK),
        ("CORRECT",   "Positive reinforcement — no override, tracked for accuracy.", ACCENT_HI),
        ("SUGGEST",   "Advisory improvement — LLM-inject only, never blocking.", ACCENT),
    ]
    cw = Inches(2.9); cx = Inches(0.7); cy = Inches(2.85)
    for i, (v, body, col) in enumerate(verdicts):
        x = cx + i * (cw + Inches(0.13))
        _rect(s, x, cy, cw, Inches(2.2), fill=SURFACE, line=BORDER)
        tick = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, cy, cw, Inches(0.06))
        _fill(tick, col); _no_line(tick)
        _text(s, x, cy + Inches(0.3), cw, Inches(0.5),
              v, size=15, color=col, bold=True, font="Consolas",
              align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.2), cy + Inches(0.9), cw - Inches(0.4),
              Inches(1.2), body, size=12, color=TEXT_DIM,
              align=PP_ALIGN.CENTER, font="Calibri")
    _text(s, Inches(0.7), Inches(5.35), Inches(11.9), Inches(1.6),
          "Admin dashboard at /admin/corrections shows: 10 metric buckets, "
          "by-surface heatmap, top reused, FP/FN signal, reviewer stats, "
          "avg approval time, 7-day trend.",
          size=13, color=TEXT_DIM, italic=True, font="Calibri")
    _footer(s, n, total, TECH_TAG)

    # 08 — Threat Intel Enrichment
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Threat-Intel Enrichment",
                 "Every extracted IOC is enriched — no manual pivoting.",
                 "Parallel calls to 8+ sources, deduped, cached, defanged for safe embed in reports.")
    sources = [
        ("VirusTotal",    "hash · URL · IP · domain reputation + AV verdicts"),
        ("AlienVault OTX","pulses · adversary campaigns · TTPs"),
        ("AbuseIPDB",     "IP abuse score · report count · category"),
        ("URLScan.io",    "screenshots · DOM · redirect chain"),
        ("Hybrid Analysis","sandbox verdict · related samples"),
        ("Shodan",        "port · service · certificate fingerprint"),
        ("GreyNoise",     "internet-noise classification (bg/bening/malicious)"),
        ("IPinfo / abuse.ch", "ASN / geo · Feodo / URLhaus / SSLBL feeds"),
    ]
    left = sources[:4]; right = sources[4:]
    yy = Inches(2.85)
    for i, (name, body) in enumerate(left):
        _text(s, Inches(0.7), yy + i * Inches(0.55), Inches(0.2),
              Inches(0.4), "▸", size=12, color=ACCENT, bold=True,
              font="Consolas")
        _text(s, Inches(0.95), yy + i * Inches(0.55), Inches(5.5),
              Inches(0.4), name, size=13, color=TEXT, bold=True, font="Consolas")
        _text(s, Inches(0.95), yy + i * Inches(0.55) + Inches(0.27),
              Inches(5.5), Inches(0.4), body, size=11, color=TEXT_MUTE,
              font="Calibri")
    for i, (name, body) in enumerate(right):
        _text(s, Inches(6.9), yy + i * Inches(0.55), Inches(0.2),
              Inches(0.4), "▸", size=12, color=SPARK, bold=True,
              font="Consolas")
        _text(s, Inches(7.15), yy + i * Inches(0.55), Inches(5.5),
              Inches(0.4), name, size=13, color=TEXT, bold=True, font="Consolas")
        _text(s, Inches(7.15), yy + i * Inches(0.55) + Inches(0.27),
              Inches(5.5), Inches(0.4), body, size=11, color=TEXT_MUTE,
              font="Calibri")
    _footer(s, n, total, TECH_TAG)

    # 09 — Key APIs
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Public API",
                 "One JWT. Nine core endpoints.",
                 "Full OpenAPI spec at /docs. Every response includes recipe + trace for replay.")
    _mono_box(s, Inches(0.7), Inches(2.85), Inches(11.9), Inches(4.0), [
        "POST  /api/auth/login                → JWT + role",
        "POST  /api/auth/change-password      → force-rotate seeded creds",
        "",
        "POST  /api/decode/smart              → hybrid archetype+magic race",
        "POST  /api/decode/chain              → N stages · aggregate verdict",
        "POST  /api/decode/chain/narrative    → ONE LLM narrative on full chain",
        "POST  /api/decode/chain/export       → .md / .json / .stix",
        "",
        "POST  /api/threat-model/analyze      → Mermaid → threat findings",
        "GET   /api/enrich/ioc                → VT / OTX / Abuse / URLScan / …",
        "",
        "GET   /api/corrections/analytics     → admin · 10 metric buckets",
        "POST  /api/corrections               → analyst refine (4-verdict)",
        "GET   /api/investigations/{id}       → owner-scoped replay",
    ], size=12)
    _footer(s, n, total, TECH_TAG)

    # 10 — Testing / Quality
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Quality Signal",
                 "Every archetype pinned by a pytest regression.",
                 "Deterministic layer is auditable — you can prove why a verdict fired.")
    stats = [
        ("115+", "pytest cases", ACCENT),
        ("40+",  "wrapper archetypes",       ACCENT_HI),
        ("97",   "decoder ops",     SPARK),
        ("0",    "known false-positive families", ACCENT),
    ]
    w = Inches(2.85); gap = Inches(0.15); x = Inches(0.7); y = Inches(2.9)
    for i, (v, l, c) in enumerate(stats):
        _stat_card(s, x + i * (w + gap), y, w, Inches(1.5), v, l, color=c)
    items = [
        ("Regression matrix.",
         "Every real-world captured payload becomes a pinned test — never breaks silently."),
        ("Corpus fuzzing.",
         "training/corpus/generator_v2.py — 900+ synthetic obfuscation shapes."),
        ("SEC audit passed.",
         "SEC-001/002/003 remediated: JWT secret rotated, admin seed idempotent, owner-scoped investigations."),
    ]
    _bullet_col(s, Inches(0.7), Inches(4.85), Inches(12.0), Inches(2.5),
                items, size=13)
    _footer(s, n, total, TECH_TAG)

    # 11 — Live Demo Script
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    _title_block(s, "Live Demo · 5 min",
                 "Watch the pipeline light up in real time.",
                 "Everything below is scripted so it always works — but pausable at every step.")
    steps = [
        ("1. Paste PS gzip loader",
         "Empire-style one-liner with variable indirection ($b='H4sI…')"),
        ("2. Hit NIVXRAY DECODE",
         "Archetype fires → stage 0 confidence 100 in < 1 s"),
        ("3. Reveal chain-mode auto-detect",
         "Multi-line paste auto-splits into 3 stages"),
        ("4. Aggregate SOC verdict",
         "Family: Destructive Wiper · Malicious 100/100 · MITRE mapped"),
        ("5. AI NARRATIVE (whole chain)",
         "ONE LLM call → analyst-grade description streamed via SSE"),
        ("6. Export .MD",
         "Ready-to-paste ticket update. Every step reproducible."),
    ]
    _bullet_col(s, Inches(0.7), Inches(2.85), Inches(12.0), Inches(4.4),
                steps, size=14)
    _footer(s, n, total, TECH_TAG)

    # 12 — Q&A / Contact
    n += 1; s = p.slides.add_slide(blank); _bg(s)
    s.shapes.add_picture(str(MARK_PNG), Inches(5.9), Inches(1.2),
                         width=Inches(1.5), height=Inches(1.5))
    _text(s, Inches(0.7), Inches(3.0), Inches(11.9), Inches(0.5),
          "QUESTIONS · CODE · ARCHITECTURE",
          size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
    _text(s, Inches(0.7), Inches(3.55), Inches(11.9), Inches(1.0),
          "Let's open the code.",
          size=36, color=TEXT, bold=True, align=PP_ALIGN.CENTER, font="Calibri")
    _text(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(0.5),
          "We'll drop into /app/backend/wrapper_archetypes.py or wherever your team wants to go.",
          size=14, color=TEXT_DIM, italic=True, align=PP_ALIGN.CENTER, font="Calibri")
    _rect(s, Inches(3.7), Inches(5.4), Inches(5.9), Inches(1.2),
          fill=SURFACE, line=ACCENT)
    _text(s, Inches(3.7), Inches(5.55), Inches(5.9), Inches(0.4),
          "engineering@nivxforge.com", size=18, color=ACCENT, bold=True,
          align=PP_ALIGN.CENTER, font="Consolas")
    _text(s, Inches(3.7), Inches(6.0), Inches(5.9), Inches(0.4),
          "github.com/nivxforge/nivxray  ·  discord.gg/nivxforge",
          size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER, font="Consolas")
    _footer(s, n, total, TECH_TAG)

    out = OUT_DIR / "NivXRay-Technical-Demo.pptx"
    p.save(out)
    return out


if __name__ == "__main__":
    a = sales_deck()
    print("Sales deck  :", a)
    b = tech_deck()
    print("Tech deck   :", b)
