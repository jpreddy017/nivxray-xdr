"""NivXRay · NAIDE Pitch Deck Builder
Produces deck/NivXRay_NAIDE_Deck.pptx from live screenshots + text content."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

OUT = "/app/deck_assets/NivXRay_NAIDE_Deck.pptx"
ASSETS = "/app/deck_assets"

# Palette (matches product)
BG_DARK      = RGBColor(0x08, 0x14, 0x0F)
BG_PANEL     = RGBColor(0x0D, 0x1E, 0x18)
ACCENT_GREEN = RGBColor(0x7E, 0xE6, 0xA8)
ACCENT_AMBER = RGBColor(0xFF, 0xB2, 0x6B)
ACCENT_RED   = RGBColor(0xFF, 0x9A, 0x9A)
TEXT_LIGHT   = RGBColor(0xE6, 0xFF, 0xE9)
TEXT_MUTE    = RGBColor(0x96, 0xC9, 0xAA)
TEXT_DIM     = RGBColor(0x4A, 0x8B, 0x63)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ---------- helpers ----------
def blank(): return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, color=BG_DARK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg

def text_box(slide, x, y, w, h, text, size=14, color=TEXT_LIGHT,
             bold=False, align=PP_ALIGN.LEFT, font="Consolas"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return tb

def eyebrow(slide, text, y=Inches(0.35)):
    text_box(slide, Inches(0.6), y, Inches(12), Inches(0.35),
             text, size=10, color=ACCENT_GREEN, bold=True)

def title(slide, text, y=Inches(0.7)):
    text_box(slide, Inches(0.6), y, Inches(12), Inches(0.9),
             text, size=32, color=TEXT_LIGHT, bold=True, font="Segoe UI")

def footer(slide, page_no, total=None):
    text_box(slide, Inches(0.6), Inches(7.15), Inches(6), Inches(0.3),
             "NivXRay · NAIDE · Deterministic Analyst Engine", size=9, color=TEXT_DIM)
    text_box(slide, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
             f"{page_no}", size=9, color=TEXT_DIM, align=PP_ALIGN.RIGHT)

def brand_stripe(slide):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.4),
                                prs.slide_width, Inches(0.1))
    r.fill.solid(); r.fill.fore_color.rgb = ACCENT_GREEN; r.line.fill.background()

# ============================ SLIDE 1 — COVER ============================
s = blank(); fill_bg(s)
text_box(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
         "NIVXRAY · NAIDE", size=14, color=ACCENT_GREEN, bold=True)
text_box(s, Inches(0.6), Inches(2.9), Inches(12), Inches(1.8),
         "Autonomous Investigation\n& Deterministic Engine",
         size=54, color=TEXT_LIGHT, bold=True, font="Segoe UI")
text_box(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.5),
         "From raw evidence to ticket-ready analyst brief · in seconds · zero LLM drift",
         size=18, color=TEXT_MUTE, font="Segoe UI")
text_box(s, Inches(0.6), Inches(5.9), Inches(12), Inches(0.4),
         "Pitch Deck · Product · Architecture · Use Cases · ROI",
         size=13, color=ACCENT_GREEN, bold=True)
brand_stripe(s)

# ============================ SLIDE 2 — THE PROBLEM ============================
s = blank(); fill_bg(s)
eyebrow(s, "THE PROBLEM"); title(s, "The SOC is drowning in evidence · but starved of clarity")
body = (
    "•  A single suspicious PowerShell blob takes an L1 analyst 20–40 minutes to peel, decode\n"
    "   and map to MITRE — often with 4-6 online tools open in parallel (CyberChef, VT,\n"
    "   AbuseIPDB, URLScan, custom regex, Excel).\n\n"
    "•  Threat-intel reports (Talos, Microsoft, Mandiant, Cisco) bury the actionable\n"
    "   evidence — IOCs, commands, MITRE — inside 15+ pages of prose.\n\n"
    "•  Modern LLM copilots are non-deterministic. Same input gives different verdicts on\n"
    "   different days. Auditors and IR teams cannot defend a non-reproducible finding.\n\n"
    "•  Existing decoders stop at 'decode'. They don't correlate IOCs, hashes, MITRE\n"
    "   techniques, LOLBAS binaries, YARA/Sigma ideas, or produce a ticket-ready brief.\n\n"
    "•  Analysts context-switch across 6+ tabs · miss 30–60% of embedded indicators."
)
text_box(s, Inches(0.6), Inches(1.8), Inches(12.3), Inches(5),
         body, size=15, color=TEXT_LIGHT, font="Segoe UI")
footer(s, 2)

# ============================ SLIDE 3 — INTRODUCING NAIDE ============================
s = blank(); fill_bg(s)
eyebrow(s, "INTRODUCING NAIDE")
title(s, "One paste. One click. A full deterministic investigation.")
text_box(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.5),
         "NivXRay Autonomous Investigation & Deterministic Engine",
         size=15, color=ACCENT_GREEN, bold=True)
bullets = (
    "•  Accepts ANYTHING an analyst can paste — a PowerShell blob, a raw URL, a threat-\n"
    "   report link, an EDR CSV row, a Sysmon XML, a defanged IOC list, or a multi-\n"
    "   stage command chain.\n\n"
    "•  Runs 60+ registered capabilities on a passive registry — decode, IOC extraction,\n"
    "   IDA (URL acquisition), DIE (semantic AST), ICE (incident correlation), MITRE\n"
    "   mapping, LOLBAS lookup, YARA-lite, IOC enrichment (VT / AbuseIPDB / OTX).\n\n"
    "•  Every finding is evidence-provenanced. Every verdict is reproducible bit-for-bit.\n"
    "   Optional LLM narrative sits ON TOP — never inside — the evidence graph.\n\n"
    "•  Ships an analyst brief that reads like a senior human analyst wrote it: Executive\n"
    "   Summary, Observed Behaviour, Attack Intent, MITRE Summary, IOC Intelligence,\n"
    "   Recommendations, and Evidence Confidence — all in ~2-8 seconds."
)
text_box(s, Inches(0.6), Inches(2.4), Inches(12.3), Inches(4.5),
         bullets, size=14, color=TEXT_LIGHT, font="Segoe UI")
footer(s, 3)

# ============================ SLIDE 4 — WHY DETERMINISTIC ============================
s = blank(); fill_bg(s)
eyebrow(s, "WHY DETERMINISTIC"); title(s, "Deterministic vs LLM-drift · the audit rule")
rows = [
    ("Reproducibility",   "Same input → same output, bit-for-bit",     "Non-deterministic · different every run"),
    ("Auditability",       "Every field carries a provenance record",   "Cannot show why a verdict was reached"),
    ("Latency",            "2-8 seconds for a full brief",              "10-60 seconds + token cost"),
    ("Cost per case",      "Zero LLM cost on the core path",            "$0.02-$0.20 per investigation"),
    ("Air-gapped ready",   "Runs fully offline — deterministic core",   "Needs cloud LLM endpoint"),
    ("SOC-analyst trust",  "Findings can be cross-checked in ~seconds", "'Why did it say that?' is unanswerable"),
]
y = Inches(1.85)
# header
text_box(s, Inches(0.6), y, Inches(3.4), Inches(0.35), "DIMENSION",       size=11, color=ACCENT_GREEN, bold=True)
text_box(s, Inches(4.2), y, Inches(4.7), Inches(0.35), "NIVXRAY · NAIDE", size=11, color=ACCENT_GREEN, bold=True)
text_box(s, Inches(9.2), y, Inches(3.9), Inches(0.35), "LLM-ONLY COPILOTS", size=11, color=ACCENT_AMBER, bold=True)
for i, (a, b, c) in enumerate(rows):
    y = Inches(2.3 + i * 0.72)
    text_box(s, Inches(0.6), y, Inches(3.4), Inches(0.6), a, size=13, color=TEXT_LIGHT, bold=True, font="Segoe UI")
    text_box(s, Inches(4.2), y, Inches(4.7), Inches(0.6), b, size=12, color=TEXT_MUTE, font="Segoe UI")
    text_box(s, Inches(9.2), y, Inches(3.9), Inches(0.6), c, size=12, color=TEXT_MUTE, font="Segoe UI")
footer(s, 4)

# ============================ SLIDE 5 — ARCHITECTURE ============================
s = blank(); fill_bg(s)
eyebrow(s, "ARCHITECTURE"); title(s, "Six-pillar deterministic core · IUE → IDA → DIE → ICE → IEDDE → L4")
arch = (
    "┌──────────────────────────────────────────────────────────────────────────────┐\n"
    "│  1 · IUE — Input Understanding Engine                                        │\n"
    "│      classifies every paste (url / powershell / cmd / prose / csv / xml)     │\n"
    "│      routes to the correct capability set — NEVER runs a decoder blindly     │\n"
    "├──────────────────────────────────────────────────────────────────────────────┤\n"
    "│  2 · IDA — Intelligent Document Acquisition   (Talos · Mandiant · MS · …)    │\n"
    "│      fetches the URL, extracts commands + IOCs + MITRE + timeline + hashes   │\n"
    "├──────────────────────────────────────────────────────────────────────────────┤\n"
    "│  3 · DIE — Deterministic Investigation Engine                                │\n"
    "│      PS-AST · CMD-AST · Bash-AST · Python-AST · JavaScript-AST · Semantic    │\n"
    "├──────────────────────────────────────────────────────────────────────────────┤\n"
    "│  4 · ICE — Incident Correlation Engine                                       │\n"
    "│      merges per-stage evidence into a single incident with tactics observed  │\n"
    "├──────────────────────────────────────────────────────────────────────────────┤\n"
    "│  5 · IEDDE — Intelligent Evidence-Driven Decoding Engine                     │\n"
    "│      surfaces WHY each decoding branch fired · with per-layer confidence     │\n"
    "├──────────────────────────────────────────────────────────────────────────────┤\n"
    "│  6 · L4 · Analyst Brief · MITRE Attack-Chain · Evidence Confidence           │\n"
    "└──────────────────────────────────────────────────────────────────────────────┘"
)
text_box(s, Inches(0.6), Inches(1.7), Inches(12.3), Inches(5.5),
         arch, size=11, color=ACCENT_GREEN, font="Consolas")
footer(s, 5)

# ============================ SLIDE 6 — COMPONENTS ============================
s = blank(); fill_bg(s)
eyebrow(s, "COMPONENTS"); title(s, "9 adapters · 10 analyzers · a passive capability registry")
comp = (
    "Adapters (input surface):\n"
    "•  URL Adapter · IDA acquires threat reports · pluggable vendor pack (Talos, MS, Mandiant, Cisco)\n"
    "•  Sysmon Event 1/3 · EVTX transport · CSV EDR · Base64/Hex/Gzip/Zlib/LZMA/AES/RC4 payloads\n"
    "•  PDF / DOCX / HTML · attachment upload · Sysmon Event 11 (File Create) · Event 22 (DNS)\n\n"
    "Analyzers (correlate + judge):\n"
    "•  IOC Enrichment (VirusTotal · AbuseIPDB · OTX · abuse.ch · URLScan · Shodan)\n"
    "•  MITRE ATT&CK Mapper · LOLBAS · YARA-lite · Sigma-hint generator\n"
    "•  Report Generator (Executive Summary · Ticket-ready Analyst Summary · Attack Intent)\n"
    "•  Artifact Intelligence (PE header · shellcode detect · defanged IOC refang)\n"
    "•  Behavioral Timeline · Query/Hunt sub-view · Predicted Process Tree · Verdict Card\n\n"
    "Passive registry:\n"
    "•  Every capability declares its accepted formats + version + implementation path\n"
    "•  Thin router dispatches; capabilities never call each other directly\n"
    "•  Determinism-lock via a canonical equivalence harness (145 zero-drift regression tests)"
)
text_box(s, Inches(0.6), Inches(1.8), Inches(12.3), Inches(5.2),
         comp, size=13, color=TEXT_LIGHT, font="Segoe UI")
footer(s, 6)

# ============================ SLIDE 7 — DEPLOYMENT SETUP 1 ============================
s = blank(); fill_bg(s)
eyebrow(s, "DEPLOYMENT MODEL 1"); title(s, "Upstream · NivXRay enriches logs BEFORE they reach SIEM/EDR")
diagram = (
    "  ┌────────────┐    ┌────────────┐    ┌────────────┐    ┌──────────────┐    ┌────────────┐\n"
    "  │  LOG       │    │  LOG       │    │            │    │              │    │            │\n"
    "  │  SOURCES   │──▶ │  COLLECTOR │──▶ │  NIVXRAY   │──▶ │  ENRICHED    │──▶ │  SIEM /    │\n"
    "  │            │    │  (agent /  │    │  · NAIDE   │    │  EVENT +     │    │  EDR       │\n"
    "  │  Sysmon    │    │  Fluentd / │    │            │    │  BRIEF       │    │  (Splunk, │\n"
    "  │  Windows   │    │  Vector /  │    │  · IUE     │    │  · MITRE     │    │  Sentinel,│\n"
    "  │  Linux     │    │  Filebeat) │    │  · DIE     │    │  · IOCs      │    │  Elastic, │\n"
    "  │  EDR CSV   │    │            │    │  · ICE     │    │  · Verdict   │    │  QRadar,  │\n"
    "  │  Netflow   │    │            │    │  · L4      │    │  · Recs      │    │  Chronicle)│\n"
    "  └────────────┘    └────────────┘    └────────────┘    └──────────────┘    └────────────┘"
)
text_box(s, Inches(0.4), Inches(1.9), Inches(12.6), Inches(3), diagram,
         size=10, color=ACCENT_GREEN, font="Consolas")
notes = (
    "Value delivered:\n"
    "•  Every event lands in SIEM already enriched with MITRE technique IDs + IOC verdict\n"
    "•  Detection rules trigger on richer signals — fewer false positives, higher fidelity\n"
    "•  Analysts see a ticket-ready brief in the alert body — not raw log bytes\n"
    "•  Reduces mean detection cost per case by ~60% (deterministic, no LLM tokens)"
)
text_box(s, Inches(0.6), Inches(5.0), Inches(12.3), Inches(2.2), notes,
         size=13, color=TEXT_LIGHT, font="Segoe UI")
footer(s, 7)

# ============================ SLIDE 8 — DEPLOYMENT SETUP 2 ============================
s = blank(); fill_bg(s)
eyebrow(s, "DEPLOYMENT MODEL 2"); title(s, "Downstream · analyst-triggered investigation from SIEM/EDR alert")
diagram = (
    "                                                 ┌──────────────┐\n"
    "                                                 │              │\n"
    "  ┌────────────┐    ┌────────────┐    ┌────────▶ │  NIVXRAY     │  ─┐\n"
    "  │            │    │            │    │  Alert   │  · NAIDE     │   │\n"
    "  │  SIEM /    │──▶ │  ALERT     │────┘  payload │              │   │\n"
    "  │  EDR       │    │  TRIGGERED │               │  autonomous  │   │\n"
    "  │            │    │            │    ◀───────── │  investigate │   │\n"
    "  │  Splunk    │    │  (attack   │    Enriched   │  (2-8 s)     │   │\n"
    "  │  Sentinel  │    │  ID + IOCs │    JSON       │              │   │\n"
    "  │  Chronicle)│    │  + hash)   │               │              │   │\n"
    "  └────────────┘    └────────────┘               └──────────────┘   │\n"
    "                                                          │        │\n"
    "                                                          ▼        ▼\n"
    "                                                 ┌───────────────────┐\n"
    "                                                 │  ANALYST UI       │\n"
    "                                                 │  · Brief          │\n"
    "                                                 │  · Attack Chain   │\n"
    "                                                 │  · IOC Intel      │\n"
    "                                                 │  · Recommendations│\n"
    "                                                 └───────────────────┘"
)
text_box(s, Inches(0.4), Inches(1.85), Inches(12.6), Inches(4.6), diagram,
         size=9, color=ACCENT_GREEN, font="Consolas")
notes = "Analyst opens the ticket · NivXRay has ALREADY completed the investigation · verify + escalate in under a minute."
text_box(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4), notes,
         size=13, color=ACCENT_GREEN, bold=True, font="Segoe UI")
footer(s, 8)

# ============================ SLIDE 9-15 — USE CASES ============================
use_cases = [
    ("USE CASE 1", "Threat-report URL triage · Talos, Mandiant, Microsoft",
     "Paste the URL of a public threat report. NAIDE acquires the article, extracts every\n"
     "command, IOC, MITRE technique, timeline event and hash — then produces a full\n"
     "deterministic analyst brief.\n\n"
     "Live example: Talos IR ransomware blog\n"
     "•  6 commands investigated · 62 IOCs · 11 MITRE techniques · 6 tactic swim-lanes\n"
     "•  Malware families: 7 identified · CVEs correlated · timeline reconstructed\n"
     "•  End-to-end time: ~2.4 seconds (vs 45-90 min analyst read + extract)",
     "02_url_pasted.png", "03_workspace_populated.png"),

    ("USE CASE 2", "Obfuscated PowerShell -encodedcommand payload",
     "Paste a PowerShell -e blob. NAIDE recursively peels every layer (base64, gzip,\n"
     "UTF-16LE, XOR, character-array). Each layer's decoded output is surfaced with a\n"
     "confidence score and MITRE mapping.\n\n"
     "•  Multi-layer chain-mode fires automatically when multi-encoding is detected\n"
     "•  LOLBAS binaries flagged with per-technique badge (T1059.001, T1218.005, etc.)\n"
     "•  YARA-lite + Sigma hunt ideas generated deterministically from the recovered payload\n"
     "•  Time: ~3-5 seconds for a 4-layer chain (vs 25 min manual CyberChef pipeline)",
     "01_input_workspace.png", "uc2_ps_result.png"),

    ("USE CASE 3", "Analyst-paste IOC sweep · defanged bundle",
     "Paste any mix of defanged URLs, hashes, IPs, CVEs. NAIDE refangs everything, dedupes,\n"
     "resolves each IOC against the MITRE technique catalog and (optionally) live OSINT.\n\n"
     "•  Hash policy — ONLY SHA-256 hashes accepted (MD5 / SHA-1 filtered at projection)\n"
     "•  Every IOC card shows: verdict · confidence · consensus · first-seen · related campaigns\n"
     "•  Providers marked 'pending' until VT / AbuseIPDB / abuse.ch keys are wired\n"
     "•  Time: ~1-2 seconds for a 50-IOC paste (vs manual VT/AbuseIPDB browsing)",
     "uc3_ioc_input.png", "uc3_ioc_result.png"),

    ("USE CASE 4", "EDR CSV / Sysmon Event 1 · endpoint telemetry",
     "Drop a CSV export from Sophos / SentinelOne / CrowdStrike / Symantec — or paste\n"
     "raw Sysmon Event XML. NAIDE parses the process trees, extracts commands, maps\n"
     "each to MITRE, and clusters the attack chain into tactic swim-lanes.\n\n"
     "•  Predicted Process Tree reconstructs parent-child from EDR rows\n"
     "•  Behavioral Timeline shows chronology of every timestamped event\n"
     "•  Query/Hunt scoped sub-view lets the analyst filter by tactic or actor\n"
     "•  Time: ~4-6 seconds for a 500-row CSV (vs multi-hour Excel triage)",
     "uc4_sysmon_input.png", "uc4_sysmon_result.png"),

    ("USE CASE 5", "Ransomware / multi-stage command chain",
     "Paste a multi-line ClickFix / EDR-bypass / Meterpreter runner. NAIDE splits into\n"
     "stages, decodes each stage's obfuscation, maps techniques per stage, and shows\n"
     "the full attack-chain trajectory across ATT&CK tactics.\n\n"
     "•  Automatic chain-mode routing when multi-command is detected\n"
     "•  Executive Summary reads like a senior analyst wrote it: risk, intent, impact\n"
     "•  Recommendations grouped as Immediate · Threat Hunting · Containment\n"
     "•  Time: ~3-8 seconds for a 4-stage chain (vs 60+ min manual pipeline)",
     "06_executive_summary.png", "05_attack_chain.png"),

    ("USE CASE 6", "Live SIEM alert enrichment (Setup-2 flow)",
     "SIEM detects suspicious PS execution → alert body pushed to NAIDE → 6 seconds\n"
     "later the analyst opens a ticket with the FULL brief already attached: MITRE\n"
     "mapping, IOC verdicts, LOLBAS badges, executive summary, containment steps.\n\n"
     "•  Analyst verifies in ~60 seconds instead of investigating from scratch\n"
     "•  Handoff export: SOC Brief .md · PDF report · STIX 2.1 bundle · JSON\n"
     "•  Every artifact is deterministic + reproducible — audit-friendly\n"
     "•  Time saved per alert: ~20-40 minutes of L1 analyst time",
     "03_workspace_populated.png", "06_executive_summary.png"),
]

for i, (eyeb, ttl, body, shot_a, shot_b) in enumerate(use_cases):
    s = blank(); fill_bg(s)
    eyebrow(s, eyeb)
    text_box(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             ttl, size=22, color=TEXT_LIGHT, bold=True, font="Segoe UI")
    text_box(s, Inches(0.6), Inches(1.6), Inches(6.2), Inches(4.5),
             body, size=12, color=TEXT_LIGHT, font="Segoe UI")
    for j, shot in enumerate([shot_a, shot_b]):
        img_path = os.path.join(ASSETS, shot)
        if os.path.exists(img_path):
            try:
                s.shapes.add_picture(img_path,
                    Inches(7.0), Inches(1.5 + j * 2.85),
                    Inches(6.0), Inches(2.7))
            except Exception:
                pass
    footer(s, 9 + i)

# ============================ SLIDE 16 — TIME SAVINGS ============================
# ============================ THREAT ANALYSIS TOUR (2 slides · 10 tabs) ============
tour_tabs = [
    ("GRAPH",   "Investigation Graph · every stage + IOC + MITRE / LOLBIN edge"),
    ("MITRE",   "MITRE ATT&CK · every technique with ID + name + tactic"),
    ("LOLBAS",  "Living-off-the-Land binaries · every abused binary + technique"),
    ("RULES",   "YARA-lite + Sigma seeds · deterministic rule ideas"),
    ("IOCS",    "IOC catalogue · SHA-256 · URL · IP · CVE · registry · path"),
    ("TI-HITS", "Threat-intel hits · VT / AbuseIPDB / OTX / abuse.ch (pending until keys wired)"),
    ("OSINT",   "OSINT enrichment · reverse-DNS · WHOIS · pDNS · relationships"),
    ("AI",      "AI-verdict pane · optional LLM overlay (never inside evidence graph)"),
    ("FLOW",    "Attack Flow · deterministic chain projection across kill-chain"),
    ("CHAIN",   "Decode chain · every stage's engine + confidence"),
]
for slide_ix in range(2):
    s = blank(); fill_bg(s)
    eyebrow(s, "THREAT ANALYSIS · SIDEBAR TOUR")
    title(s, ["10 tabs · every angle of the investigation (1/2)",
              "10 tabs · every angle of the investigation (2/2)"][slide_ix])
    subset = tour_tabs[slide_ix*5:(slide_ix+1)*5]
    for i, (tab, desc) in enumerate(subset):
        col = i % 3
        row = i // 3
        x = Inches(0.5 + col * 4.3)
        y = Inches(1.65 + row * 2.85)
        # tab label
        text_box(s, x, y, Inches(4.1), Inches(0.35),
                 tab, size=13, color=ACCENT_GREEN, bold=True, font="Consolas")
        # image
        img_path = os.path.join(ASSETS, f"tab_{tab.lower().replace('-','_')}.png")
        if os.path.exists(img_path):
            try:
                s.shapes.add_picture(img_path, x, y + Inches(0.4),
                                      Inches(4.1), Inches(1.9))
            except Exception:
                pass
        text_box(s, x, y + Inches(2.35), Inches(4.1), Inches(0.5),
                 desc, size=9, color=TEXT_MUTE, font="Segoe UI")
    footer(s, 16 + slide_ix)

# ============================ SLIDE 18 — TIME SAVINGS (moved) ============================
s = blank(); fill_bg(s)
eyebrow(s, "ROI · TIME SAVINGS"); title(s, "How much analyst time does NivXRay actually save?")
rows = [
    ("Threat-report URL triage",      "45-90 min",  "2-4 sec",     "~1,200×"),
    ("Obfuscated PS -e blob decode",  "20-40 min",  "3-5 sec",     "~400×"),
    ("50-IOC defanged bundle sweep",  "30-60 min",  "1-2 sec",     "~1,500×"),
    ("500-row EDR CSV triage",         "2-4 hours",  "4-6 sec",     "~1,200×"),
    ("Multi-stage command chain",     "60-90 min",  "3-8 sec",     "~900×"),
    ("SIEM-alert full enrichment",    "20-40 min",  "6-8 sec",     "~250×"),
]
y = Inches(1.9)
text_box(s, Inches(0.6), y, Inches(4.6), Inches(0.35),
         "TASK", size=11, color=ACCENT_GREEN, bold=True)
text_box(s, Inches(5.5), y, Inches(2.6), Inches(0.35),
         "MANUAL ANALYST", size=11, color=ACCENT_AMBER, bold=True)
text_box(s, Inches(8.4), y, Inches(2.4), Inches(0.35),
         "NIVXRAY · NAIDE", size=11, color=ACCENT_GREEN, bold=True)
text_box(s, Inches(11.1), y, Inches(2.0), Inches(0.35),
         "SPEED-UP", size=11, color=ACCENT_GREEN, bold=True)
for i, (a, b, c, d) in enumerate(rows):
    y = Inches(2.35 + i * 0.55)
    text_box(s, Inches(0.6), y, Inches(4.6), Inches(0.5), a, size=13, color=TEXT_LIGHT, font="Segoe UI")
    text_box(s, Inches(5.5), y, Inches(2.6), Inches(0.5), b, size=13, color=ACCENT_AMBER, font="Segoe UI")
    text_box(s, Inches(8.4), y, Inches(2.4), Inches(0.5), c, size=13, color=ACCENT_GREEN, bold=True, font="Segoe UI")
    text_box(s, Inches(11.1), y, Inches(2.0), Inches(0.5), d, size=13, color=TEXT_LIGHT, bold=True, font="Segoe UI")
text_box(s, Inches(0.6), Inches(5.9), Inches(12.3), Inches(1.2),
         "Payback: 1 senior L2 analyst FTE ≈ USD 120-180K/yr. NivXRay recovers 40-60% of\n"
         "L1/L2 triage hours — freeing analysts for hunt, tuning and IR-work that actually\n"
         "requires human judgement.",
         size=13, color=TEXT_MUTE, font="Segoe UI")
footer(s, 16)

# ============================ SLIDE 17 — INDUSTRY FIT ============================
s = blank(); fill_bg(s)
eyebrow(s, "WHERE IT FITS"); title(s, "Industry scope · every SOC that reads bytes")
body = (
    "Primary buyers:\n"
    "•  Managed Security Service Providers (MSSPs) — 24×7 SOC teams with 40-200 analysts\n"
    "•  Financial-services CSIRT · high-signal alerts · audit-heavy · needs reproducibility\n"
    "•  Healthcare + government · air-gapped deployments where cloud LLM is a non-starter\n"
    "•  Fortune-500 in-house SOCs · already have SIEM + EDR · missing the correlator/brief layer\n"
    "•  IR consultancies · Mandiant / Palo Alto / SecureWorks style engagement teams\n\n"
    "Integration surface:\n"
    "•  Splunk · Sentinel · Chronicle · Elastic · QRadar · Sumo Logic · Devo · Panther\n"
    "•  CrowdStrike · SentinelOne · Sophos · Symantec · Defender for Endpoint\n"
    "•  SOAR: Phantom · XSOAR · Swimlane · Torq · Tines — NivXRay is called as an action node\n"
    "•  Ticketing: ServiceNow · Jira SD · Zendesk — analyst brief attaches directly\n\n"
    "Compliance-friendly by design:\n"
    "•  SOC2 · ISO27001 · NIST-CSF · PCI-DSS · HIPAA — deterministic findings pass audit reviews\n"
    "•  STIX 2.1 output · ATT&CK-Navigator layer export · Sigma / YARA rule seeds"
)
text_box(s, Inches(0.6), Inches(1.75), Inches(12.3), Inches(5.5),
         body, size=13, color=TEXT_LIGHT, font="Segoe UI")
footer(s, 17)

# ============================ SLIDE 18 — DIFFERENTIATORS ============================
s = blank(); fill_bg(s)
eyebrow(s, "DIFFERENTIATORS"); title(s, "Why NivXRay wins against CyberChef, EDR-CoPilots and LLM tools")
items = [
    ("Zero LLM drift on the core path",      "Every finding survives a 145-test canonical determinism harness."),
    ("Evidence provenance on every field",   "Analysts show WHERE each MITRE tag / IOC / verdict came from."),
    ("2-8 second full brief",                "No spinning cursor. No 30-second LLM wait. Instant analyst context."),
    ("Air-gap deployable",                    "Deterministic core has no cloud dependency. Runs on a laptop."),
    ("Threat-report acquisition built-in",   "IDA fetches Talos/Mandiant/MS articles and extracts evidence."),
    ("SHA-256 IOC policy",                    "MD5 / SHA-1 stripped at wire boundary — audit-clean by default."),
    ("60+ registered capabilities",           "9 adapters · 10 analyzers · thin router · trivially extensible."),
    ("Deterministic + optional LLM layer",   "LLM narrative sits ON TOP — never inside — the evidence graph."),
]
for i, (h, d) in enumerate(items):
    y = Inches(1.85 + i * 0.62)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Emu(50000),
                              Inches(0.15), Inches(0.15))
    dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT_GREEN
    dot.line.fill.background()
    text_box(s, Inches(1.0), y, Inches(4.5), Inches(0.5),
             h, size=13, color=TEXT_LIGHT, bold=True, font="Segoe UI")
    text_box(s, Inches(5.6), y, Inches(7.5), Inches(0.5),
             d, size=12, color=TEXT_MUTE, font="Segoe UI")
footer(s, 18)

# ============================ SLIDE 19 — ROADMAP ============================
s = blank(); fill_bg(s)
eyebrow(s, "ROADMAP"); title(s, "What ships next")
r = (
    "Q1 2026 (in flight):\n"
    "•  M0f Production Cutover — thin-router replaces legacy pipeline (equivalence-proven)\n"
    "•  URL-acquisition beyond Talos — Sekoia, Trend Micro, Kaspersky, Palo Alto Unit 42\n"
    "•  OCR-Lite (SHADOW) — extract code screenshots embedded in threat reports\n\n"
    "Q2 2026:\n"
    "•  SOAR playbook nodes — Torq, Tines, XSOAR, Splunk SOAR (Phantom)\n"
    "•  Sysmon Event 22 (DNS) + Event 11 (File Create) adapter productionisation\n"
    "•  Attack Chain auto-scroll + focus-follows-technique interactions\n\n"
    "Q3-Q4 2026:\n"
    "•  Multi-tenant SaaS · per-tenant OSINT keys · role-based investigation isolation\n"
    "•  ATT&CK Navigator live layer publishing · YARA/Sigma rule review loop\n"
    "•  Custom analyst 'critic' agents on the LLM overlay — non-blocking, non-drift"
)
text_box(s, Inches(0.6), Inches(1.75), Inches(12.3), Inches(5.5),
         r, size=13, color=TEXT_LIGHT, font="Segoe UI")
footer(s, 19)

# ============================ SLIDE 20 — DEMO ACCESS ============================
s = blank(); fill_bg(s)
eyebrow(s, "TRY IT"); title(s, "Live demo · try any input right now")
demo = (
    "Preview:  https://greeting-app-5782.preview.emergentagent.com/\n"
    "Prod:     https://nivxray.nivxforge.com/\n\n"
    "Suggested test inputs:\n"
    "  1.  https://blog.talosintelligence.com/talos-ir-ransomware-engagements-...\n"
    "  2.  Any -EncodedCommand PowerShell blob\n"
    "  3.  A defanged IOC list from an incident report\n"
    "  4.  A Sysmon Event 1 XML block\n"
    "  5.  A Sophos / SentinelOne CSV export row\n\n"
    "For each: click AUTO INVESTIGATE. Scroll down. Read the Deterministic Analyst Brief."
)
text_box(s, Inches(0.6), Inches(1.85), Inches(12.3), Inches(5),
         demo, size=15, color=TEXT_LIGHT, font="Consolas")
footer(s, 20)

# ============================ SLIDE 21 — REGEN PROMPTS ============================
s = blank(); fill_bg(s)
eyebrow(s, "APPENDIX"); title(s, "Regenerate this deck · Claude & ChatGPT prompts")
p = (
    "Prompt (Claude / GPT — paste and iterate):\n\n"
    "\"You are a security-product presales architect. Produce a 20-slide deck for NivXRay\n"
    "(NAIDE — Autonomous Investigation & Deterministic Engine), a deterministic SOC-analyst\n"
    "engine that takes any pasted evidence (URL / PowerShell / CSV / Sysmon / IOC list) and\n"
    "returns a full ticket-ready analyst brief in 2-8 seconds.\n\n"
    "Cover: (1) SOC pain today, (2) NAIDE elevator pitch, (3) Deterministic vs LLM-drift,\n"
    "(4) 6-pillar architecture (IUE → IDA → DIE → ICE → IEDDE → L4), (5) components,\n"
    "(6) Deployment Model 1 (log-source enrichment upstream of SIEM/EDR),\n"
    "(7) Deployment Model 2 (SIEM/EDR alert → NAIDE → analyst),\n"
    "(8) 6 use cases with INPUT screenshot + OUTPUT/MITRE-graph/Executive-Summary shots +\n"
    "elapsed time, (9) ROI table showing 250-1500× speed-up vs manual analyst,\n"
    "(10) industry-fit + integrations (Splunk/Sentinel/Chronicle/SOAR),\n"
    "(11) differentiators, (12) roadmap, (13) demo URL + regen prompt appendix.\n\n"
    "Tone: technical, evidence-driven, no hyperbole. Palette: matrix-green on near-black.\n"
    "Font: Consolas for code, Segoe UI for body. Every claim must be verifiable from live UI.\""
)
text_box(s, Inches(0.6), Inches(1.75), Inches(12.3), Inches(5.5),
         p, size=12, color=TEXT_LIGHT, font="Consolas")
footer(s, 21)

prs.save(OUT)
print(f"SAVED: {OUT}")
print(f"SIZE:  {os.path.getsize(OUT)} bytes")
print(f"SLIDES: {len(prs.slides)}")
